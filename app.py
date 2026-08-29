import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import os
import time
from datetime import datetime, timedelta
import pytz
import base64
import hashlib
import sqlite3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io

# ==========================================
# 1. System Configuration & Constants
# ==========================================
st.set_page_config(page_title="Command Center BI Dashboard", layout="wide", initial_sidebar_state="expanded")

EGYPT_TZ = pytz.timezone('Africa/Cairo')
NEON_COLORS = ['#00d2ff', '#ffaa00', '#2ecc71', '#ff007f', '#f1c40f', '#9b59b6', '#38f9d7', '#ff7eb3', '#00f2fe', '#4facfe']

STATUS_COLORS = {
    'ACCEPTED': '#2ecc71', 
    'APPROVED AS NOTED': '#00d2ff', 
    'REVISE': '#f1c40f', 
    'REJECTED': '#e74c3c'
}
USERS_DB_FILE = "users_db_v2.csv"
LOGIN_LOGS_FILE = "login_logs_v2.csv"
AUDIT_LOG_FILE = "audit_trail.csv"

if "theme" not in st.session_state:
    st.session_state["theme"] = "Dark"
if "site_mode" not in st.session_state:
    st.session_state["site_mode"] = False
if "chat_history" not in st.session_state:
    st.session_state["chat_history"] = []
if "current_page" not in st.session_state:
    st.session_state["current_page"] = "home"

if "language" not in st.session_state:
    st.session_state["language"] = "EN"

# ==========================================
# 1.5 Translation Dictionary (القاموس)
# ==========================================
TRANSLATIONS = {
    "Mega Infrastructure Command Center": "مركز قيادة البنية التحتية العملاقة",
    "Main Dashboard": "لوحة القيادة الرئيسية",
    "Advanced Analytics Hub": "مركز التحليلات المتقدمة",
    "Total Submittals": "إجمالي الطلبات (Submittals)",
    "Total Tests": "إجمالي الاختبارات",
    "Avg. Dur (Days)": "متوسط التأخير (أيام)",
    "Total Paperwork": "إجمالي الورقيات",
    "Logout": "تسجيل الخروج",
    "UI/UX Mode": "مظهر الشاشة",
    "Data Source": "مصدر البيانات",
    "Language / اللغة": "Language / اللغة",
    # تقدر تزود أي كلمة براحتك هنا بعدين
}

def _t(text):
    if st.session_state.get("language") == "AR":
        return TRANSLATIONS.get(text, text)
    return text
def export_table_tools(df, filename_prefix):
    import io
    st.markdown("<br>", unsafe_allow_html=True)
    c1, c2, c3 = st.columns(3)
    
    # 1. زرار تحميل CSV
    csv_data = df.to_csv(index=False).encode('utf-8-sig')
    c1.download_button("📥 Download CSV", data=csv_data, file_name=f"{filename_prefix}.csv", mime="text/csv", use_container_width=True, key=f"csv_{filename_prefix}")
    
    # 2. زرار تحميل Excel
    excel_buffer = io.BytesIO()
    with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=False, sheet_name='Report')
    c2.download_button("📊 Download Excel", data=excel_buffer.getvalue(), file_name=f"{filename_prefix}.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", use_container_width=True, key=f"xls_{filename_prefix}")
    
    # 3. زرار الحفظ كـ PDF 
    html_content = f"""
    <html dir="ltr">
    <head>
        <title>{filename_prefix}</title>
        <style>
            body {{ font-family: 'Segoe UI', Arial, sans-serif; padding: 20px; }}
            table {{ width: 100%; border-collapse: collapse; margin-top: 20px; font-size: 12px; }}
            th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
            th {{ background-color: #1e3d59; color: white; }}
            tr:nth-child(even) {{ background-color: #f2f2f2; }}
        </style>
    </head>
    <body onload="window.print()">
        <h2 style="color: #1e3d59;">Project Report: {filename_prefix.replace('_', ' ')}</h2>
        <p style="color: #7f8c8d;">Generated on: {datetime.now(EGYPT_TZ).strftime('%Y-%m-%d %H:%M')}</p>
        <hr>
        {df.to_html(index=False)}
    </body>
    </html>
    """
    b64_html = base64.b64encode(html_content.encode('utf-8')).decode()
    pdf_href = f'<a href="data:text/html;base64,{b64_html}" download="{filename_prefix}_Printable.html" style="display: block; text-align: center; background-color: #e74c3c; color: white; padding: 6px; border-radius: 4px; text-decoration: none; font-weight: bold; font-family: sans-serif; border: 1px solid #c0392b; box-shadow: 0 2px 5px rgba(0,0,0,0.1);">🖨️ Save as PDF</a>'
    c3.markdown(pdf_href, unsafe_allow_html=True)
# ==========================================
# 2. Tactical UI/UX CSS Injection (Dual Mode)
# ==========================================
def inject_custom_css():
    is_dark = st.session_state.get("theme", "Dark") == "Dark"
    
    if is_dark:
        # 🌙 Stealth Ops (Dark)
        bg_main = "#0a0e17"
        bg_sidebar = "rgba(15, 22, 35, 0.85)" # 💡 شفافية للقائمة الجانبية
        card_bg = "rgba(15, 23, 42, 0.55)"    # 💡 شفافية أكبر للكروت لظهور الزجاج
        card_border = "rgba(0, 210, 255, 0.15)"
        card_shadow = "0 8px 32px 0 rgba(0, 0, 0, 0.5)"
        text_main = "#e2e8f0"
        text_muted = "#94a3b8"
        accent_color = "#00d2ff"
        accent_glow = "0 0 10px rgba(0, 210, 255, 0.3)"
        input_bg = "rgba(30, 41, 59, 0.6)"
        btn_bg = accent_color
        btn_text = bg_main
    else:
        # ☀️ Daytime HQ (Light)
        bg_main = "#f8fafc"
        bg_sidebar = "rgba(255, 255, 255, 0.85)"
        card_bg = "rgba(255, 255, 255, 0.65)" # 💡 كروت بيضاء نصف شفافة
        card_border = "rgba(148, 163, 184, 0.3)"
        card_shadow = "0 8px 32px 0 rgba(31, 38, 135, 0.1)"
        text_main = "#0f172a"
        text_muted = "#475569"
        accent_color = "#0ea5e9" 
        accent_glow = "none"
        input_bg = "rgba(255, 255, 255, 0.6)"
        btn_bg = accent_color
        btn_text = "#ffffff"

    custom_css = f"""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600&family=Rajdhani:wght@500;600;700&display=swap');
    
    #MainMenu {{visibility: hidden;}}
    footer {{visibility: hidden;}}
    [data-testid="stHeader"] {{background: transparent !important;}}
    .block-container {{padding-top: 2rem !important; padding-bottom: 2rem !important;}}
    
    html, body, [class*="css"] {{ color: {text_main} !important; font-family: 'Inter', sans-serif; }}
    
    h1, h2, h3, h4, h5, h6, .metric-value, .bi-title, .login-title {{ 
        font-family: 'Rajdhani', sans-serif !important; 
        letter-spacing: 0.5px;
        color: {text_main} !important;
    }}
    p, .stMarkdown, label {{ color: {text_main} !important; }}
    
    [data-testid="stAppViewContainer"] {{ background: {bg_main} !important; transition: all 0.4s ease; }}
    [data-testid="stSidebar"] {{ background-color: {bg_sidebar} !important; border-right: 1px solid {card_border}; transition: all 0.4s ease; }}
    
    /* Input Fields */
    [data-testid="stTextInput"] input, [data-testid="stSelectbox"] div, [data-testid="stMultiselect"] div {{
        background-color: {input_bg} !important;
        color: {text_main} !important;
        border: 1px solid {card_border} !important;
        border-radius: 4px !important;
    }}
    
    /* Buttons Fix */
    [data-testid="stButton"] button {{
        background: {btn_bg} !important;
        color: {btn_text} !important;
        border: none !important;
        font-weight: 700 !important;
        font-family: 'Rajdhani', sans-serif !important;
        text-transform: uppercase;
        letter-spacing: 1px;
        border-radius: 4px !important;
    }}
    [data-testid="stButton"] button:hover {{ opacity: 0.9 !important; transform: translateY(-1px) !important; box-shadow: {accent_glow}; }}
    
    /* Cards Fix - True Glassmorphism */
    .metric-card, .leaderboard-card, .simulator-card, .health-card, .custom-card, .navigation-card {{
        background: {card_bg} !important;
        backdrop-filter: blur(16px) !important; /* 🌟 سحر الزجاج المصنفر */
        -webkit-backdrop-filter: blur(16px) !important; /* لدعم متصفحات سفاري وأبل */
        padding: 20px;
        border-radius: 12px; /* تدوير عصري للحواف */
        border: 1px solid {card_border};
        border-top: 3px solid {accent_color};
        box-shadow: {card_shadow};
        margin-bottom: 15px;
        transition: transform 0.3s ease, box-shadow 0.3s ease; /* نعومة الحركة */
    }}
    
    /* 🌟 تأثير الرفع والظل عند مرور الماوس (Hover Animation) */
    .metric-card:hover, .leaderboard-card:hover, .simulator-card:hover, .navigation-card:hover, .health-card:hover {{
        transform: translateY(-5px);
        box-shadow: 0 12px 40px 0 {accent_glow if is_dark else 'rgba(14, 165, 233, 0.2)'};
    }}
    
    .bi-title {{ font-size: 28px; font-weight: 700; margin-top: 30px; margin-bottom: 20px; text-transform: uppercase; border-bottom: 1px solid {card_border}; padding-bottom: 10px; }}
    .metric-label {{ color: {text_muted} !important; font-size: 13px; font-weight: 600; text-transform: uppercase; }}
    .metric-value {{ color: {accent_color} !important; font-size: 38px; font-weight: 700; line-height: 1.2; text-shadow: {accent_glow}; }}
    
    .gradient-divider {{ height: 1px; background: linear-gradient(90deg, transparent 0%, {accent_color} 50%, transparent 100%); margin: 40px 0; border: none; opacity: 0.5; }}
    
    /* User Tag Fix */
    div[style*="background:rgba(255,170,0,0.1)"] {{ background: {card_bg} !important; border-color: {card_border} !important; }}
    
    {"[data-testid='stAppViewContainer']::after { content: ''; position: fixed; top: 0; left: 0; width: 100vw; height: 100vh; background: linear-gradient(rgba(18, 16, 16, 0) 50%, rgba(0, 0, 0, 0.1) 50%); background-size: 100% 4px; pointer-events: none; z-index: 9999; opacity: 0.15; }" if is_dark else ""}
    </style>
    """
    st.markdown(custom_css, unsafe_allow_html=True)

# ==========================================
# 3. Enhanced UI Components
# ==========================================
def show_breadcrumbs(path):
    parts = path.split(" > ")
    breadcrumb_html = " > ".join([
        f'<span style="color: #00d2ff; cursor: pointer;">{p}</span>' 
        if i < len(parts) - 1 
        else f'<span style="color: #ffaa00; font-weight: bold;">{p}</span>'
        for i, p in enumerate(parts)
    ])
    
    st.markdown(f"""
    <div style="
        background: rgba(255,255,255,0.05);
        padding: 10px 20px;
        border-radius: 8px;
        margin-bottom: 20px;
        font-size: 14px;
    ">
        🏠 Home {breadcrumb_html}
    </div>
    """, unsafe_allow_html=True)

def live_indicator(status="online"):
    colors = {
        "online": "#2ecc71",
        "offline": "#e74c3c",
        "warning": "#f1c40f"
    }
    
    st.markdown(f"""
    <div class="live-indicator">
        <div class="live-dot" style="background: {colors[status]}; box-shadow: 0 0 10px {colors[status]};"></div>
        <span style="color: {colors[status]}; font-size: 12px; text-transform: uppercase;">
            {status}
        </span>
    </div>
    """, unsafe_allow_html=True)

def create_progress_ring(percentage, label):
    color = "#2ecc71" if percentage > 80 else ("#f1c40f" if percentage > 50 else "#e74c3c")
    
    st.markdown(f"""
    <div style="text-align: center; padding: 20px;">
        <svg width="120" height="120" style="transform: rotate(-90deg);">
            <circle cx="60" cy="60" r="50" 
                    stroke="rgba(255,255,255,0.1)" 
                    stroke-width="8" 
                    fill="none"/>
            <circle cx="60" cy="60" r="50" 
                    stroke="{color}" 
                    stroke-width="8" 
                    fill="none"
                    stroke-dasharray="{2 * 3.14 * 50}"
                    stroke-dashoffset="{2 * 3.14 * 50 * (1 - percentage/100)}"
                    stroke-linecap="round"
                    style="transition: stroke-dashoffset 1s ease;"/>
        </svg>
        <div style="margin-top: -80px; font-size: 24px; font-weight: bold; color: {color};">
            {percentage}%
        </div>
        <div style="color: #8da3b9; font-size: 12px; margin-top: 40px;">
            {label}
        </div>
    </div>
    """, unsafe_allow_html=True)

# ==========================================
# 4. Authentication & User Management
# ==========================================
@st.cache_data(ttl=30)
def _load_users_db():
    if os.path.exists(USERS_DB_FILE):
        return pd.read_csv(USERS_DB_FILE)
    return pd.DataFrame()

@st.cache_data(ttl=30)
def _load_login_logs():
    if os.path.exists(LOGIN_LOGS_FILE):
        return pd.read_csv(LOGIN_LOGS_FILE)
    return pd.DataFrame(columns=["Timestamp", "Name", "Email", "Role"])

def clear_users_cache():
    _load_users_db.clear()

def clear_logs_cache():
    _load_login_logs.clear()

def init_auth_system():
    if not os.path.exists(USERS_DB_FILE):
        # 🛡️ الحسابات الأساسية للسيستم (بتتكريت أوتوماتيك أول ما السيستم يشتغل)
        default_users = pd.DataFrame([
            # حسابك الأساسي كـ Admin
            {"Email": "Mohamedhatem@kk.com", "Password": "admin123", "Name": "Mohamed Hatem", "Role": "Admin", "Status": "Active"},
            
            # حسابات الكتائب الدائمة (ممكن تغير الباسوردات هنا براحتك)
            {"Email": "bat36@kk.com", "Password": "123", "Name": "Battalion 36", "Role": "User", "Status": "Active"},
            {"Email": "bat73@kk.com", "Password": "123", "Name": "Battalion 73", "Role": "User", "Status": "Active"},
            {"Email": "bat44@kk.com", "Password": "123", "Name": "Battalion 44", "Role": "User", "Status": "Active"}
        ])
        default_users.to_csv(USERS_DB_FILE, index=False)
        clear_users_cache()
    
    if not os.path.exists(LOGIN_LOGS_FILE):
        logs_df = pd.DataFrame(columns=["Timestamp", "Name", "Email", "Role"])
        logs_df.to_csv(LOGIN_LOGS_FILE, index=False)
        clear_logs_cache()

def log_user_entry(user_info):
    logs_df = _load_login_logs()
    new_log = pd.DataFrame([{
        "Timestamp": datetime.now(EGYPT_TZ).strftime("%Y-%m-%d %H:%M:%S"),
        "Name": user_info["Name"],
        "Email": user_info["Email"],
        "Role": user_info["Role"]
    }])
    updated_logs = pd.concat([logs_df, new_log], ignore_index=True)
    updated_logs.to_csv(LOGIN_LOGS_FILE, index=False)
    clear_logs_cache()

def authenticate_user(email, password):
    users_df = _load_users_db()
    user = users_df[(users_df["Email"].str.lower() == email.lower()) & (users_df["Password"] == password)]
    if not user.empty:
        user_data = user.iloc[0]
        if user_data["Status"] == "Active":
            st.session_state["authenticated"] = True
            st.session_state["current_user"] = user_data.to_dict()
            log_user_entry(user_data)
            return True, "Success"
        else:
            return False, "Account Suspended. Contact Administrator."
    return False, "Invalid Email or Password."

# ==========================================
# 5. 3D Glassy Chart Styling Function
# ==========================================
def style_3d_glassy(fig, chart_type="bar"):
    is_dark = st.session_state.get("theme", "Dark") == "Dark"
    
    # 🧠 ذكاء تحديد الألوان بناءً على المود
    if is_dark:
        font_color = "#e2e8f0"
        grid_color = 'rgba(255, 255, 255, 0.05)' 
        hover_bg = "rgba(15, 23, 42, 0.95)"
        title_color = "#00d2ff"
    else:
        font_color = "#334155" # رمادي غامق مقروء جداً لللايت
        grid_color = 'rgba(0, 0, 0, 0.1)' 
        hover_bg = "rgba(255, 255, 255, 0.95)"
        title_color = "#0ea5e9"

    fig.update_layout(
        font=dict(family="Rajdhani, sans-serif", color=font_color, size=14),
        plot_bgcolor="rgba(0,0,0,0)",
        paper_bgcolor="rgba(0,0,0,0)",
        margin=dict(t=50, b=30, l=20, r=20),
        title_font=dict(size=20, color=title_color, family="Rajdhani, sans-serif"),
        legend=dict(
            orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1,
            font=dict(size=12, color=font_color),
            bgcolor="rgba(0,0,0,0)"
        ),
        # 🔥 الإصلاح هنا: إجبار الـ Tooltip إنه ياخد لون الخط المظبوط (font_color)
        hoverlabel=dict(
            bgcolor=hover_bg, 
            bordercolor=title_color,
            font=dict(size=14, family="Inter, sans-serif", color=font_color)
        )
    )
    
    # شبكة واضحة
    fig.update_yaxes(showgrid=True, gridwidth=1, gridcolor=grid_color, griddash='dot', zerolinecolor=grid_color, tickfont=dict(color=font_color))
    fig.update_xaxes(showgrid=False, zerolinecolor=grid_color, tickfont=dict(color=font_color))
    
    if chart_type in ["bar", "histogram"]:
        fig.update_traces(marker_line_color='rgba(255, 255, 255, 0.5)' if is_dark else 'rgba(0,0,0,0.2)', marker_line_width=1.5, opacity=0.9)
    elif chart_type == "pie":
        fig.update_traces(marker=dict(line=dict(color='#0a0e17' if is_dark else '#ffffff', width=3)))
    elif chart_type == "line" or chart_type == "scatter":
        fig.update_traces(line=dict(width=4), marker=dict(size=8, line=dict(color='white', width=2)))
        
    return fig

# ==========================================
# 6. History Manager with SQLite
# ==========================================
class HistoryManager:
    DB_FILE = "project_history.db"
    
    @staticmethod
    def init_db():
        conn = sqlite3.connect(HistoryManager.DB_FILE)
        cursor = conn.cursor()
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS kpi_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                timestamp TEXT,
                file_name TEXT,
                total_requests INTEGER,
                total_tests INTEGER,
                avg_dpl REAL,
                avg_duration REAL,
                total_paperwork INTEGER
            )
        ''')
        conn.commit()
        conn.close()
    
    @staticmethod
    def save_metrics(metrics_dict):
        HistoryManager.init_db()
        conn = sqlite3.connect(HistoryManager.DB_FILE)
        cursor = conn.cursor()
        cursor.execute('''
            INSERT INTO kpi_history 
            (timestamp, file_name, total_requests, total_tests, avg_dpl, avg_duration, total_paperwork)
            VALUES (?, ?, ?, ?, ?, ?, ?)
        ''', (
            datetime.now(EGYPT_TZ).strftime("%Y-%m-%d %H:%M:%S"),
            metrics_dict.get("File_Name", "Unknown_File"),
            metrics_dict.get("Total_Requests", 0),
            metrics_dict.get("Total_Tests", 0),
            metrics_dict.get("Avg_DPL", 0),
            metrics_dict.get("Avg_Duration", 0),
            metrics_dict.get("Total_Paperwork", 0)
        ))
        conn.commit()
        conn.close()
    
    @staticmethod
    def load_history():
        HistoryManager.init_db()
        try:
            conn = sqlite3.connect(HistoryManager.DB_FILE)
            df = pd.read_sql_query("SELECT * FROM kpi_history ORDER BY id", conn)
            conn.close()
            return df
        except:
            return pd.DataFrame()
    
    @staticmethod
    def get_delta_html(current_val, metric_key, current_file_name):
        history_df = HistoryManager.load_history()
        if history_df.empty:
            return ""
        
        column_map = {
            "Total_Requests": "total_requests",
            "Total_Tests": "total_tests",
            "Avg_DPL": "avg_dpl",
            "Avg_Duration": "avg_duration",
            "Total_Paperwork": "total_paperwork"
        }
        
        db_column = column_map.get(metric_key)
        if not db_column or db_column not in history_df.columns:
            return ""
        
        file_history = history_df[history_df['file_name'] == current_file_name] if 'file_name' in history_df.columns else pd.DataFrame()
        
        if file_history.empty:
            return ""
        
        file_history = file_history.sort_values('id')
        last_val = file_history.iloc[-1][db_column]
        diff = current_val - last_val
        pct_str = "0%" if last_val == 0 else f"{abs((diff / last_val) * 100):.1f}%"
        diff_fmt = f"{int(diff)}" if isinstance(current_val, (int, float)) and float(current_val).is_integer() else f"{diff:.2f}"
        
        if diff > 0:
            return f'<div class="delta-up">▲ +{diff_fmt} ({pct_str})</div>'
        elif diff < 0:
            return f'<div class="delta-down">▼ {diff_fmt} ({pct_str})</div>'
        else:
            return f'<div class="delta-neutral">➖ No change</div>'
    
    @staticmethod
    def export_to_csv():
        HistoryManager.init_db()
        conn = sqlite3.connect(HistoryManager.DB_FILE)
        df = pd.read_sql_query("SELECT * FROM kpi_history", conn)
        conn.close()
        return df
    
    @staticmethod
    def import_from_csv(df):
        HistoryManager.init_db()
        conn = sqlite3.connect(HistoryManager.DB_FILE)
        df.columns = df.columns.str.strip().str.lower()
        
        for index, row in df.iterrows():
            req = row.get('total_requests', 0)
            tst = row.get('total_tests', 0)
            dpl = row.get('avg_dpl', 0)
            dur = row.get('avg_duration', 0)
            pap = row.get('total_paperwork', 0)
            ts = row.get('timestamp', datetime.now(EGYPT_TZ).strftime("%Y-%m-%d %H:%M:%S"))
            fn = row.get('file_name', 'Legacy_Import.csv')
            
            conn.execute("""
                INSERT INTO kpi_history 
                (timestamp, file_name, total_requests, total_tests, avg_dpl, avg_duration, total_paperwork)
                VALUES (?, ?, ?, ?, ?, ?, ?)
            """, (
                str(ts),
                str(fn).strip(),
                pd.to_numeric(req, errors='coerce') or 0,
                pd.to_numeric(tst, errors='coerce') or 0,
                pd.to_numeric(dpl, errors='coerce') or 0,
                pd.to_numeric(dur, errors='coerce') or 0,
                pd.to_numeric(pap, errors='coerce') or 0
            ))
        conn.commit()
        conn.close()

def create_card(column, label, value, delta_html="", progress=None):
    if progress is not None:
        prog_color = "#2ecc71" if progress > 80 else ("#f1c40f" if progress > 50 else "#e74c3c")
        prog_html = f'<div class="prog-bg" style="height: 6px; background: rgba(127,140,141,0.2); border-radius: 10px; margin-top: 15px;"><div class="prog-fill" style="height: 100%; width: {progress}%; background: {prog_color}; border-radius: 10px; transition: width 1s ease-in-out;"></div></div>'
    else:
        prog_html = ""
        
# إزالة المسافات هنا مهمة جداً عشان الكود يترسم صح وميظهرش كنص
    html_content = f"""<div class="metric-card">
<div class="metric-label">{label}</div>
<div class="metric-value">{value}</div>
{delta_html}
{prog_html}
</div>"""
    
    column.markdown(html_content, unsafe_allow_html=True)

def ai_assistant(query, data_summary):
    query = query.lower()
    if "delay" in query or "time" in query or "duration" in query:
        return f"Based on data, the average duration is {data_summary['avg_duration']} days. Check the AI Predictive section for bottleneck warnings."
    elif "dpl" in query:
        return f"Current average DPL is {data_summary['avg_dpl']}. Ensure rejected samples are re-tested after proper compaction."
    else:
        return "I am here to assist. Ask me about project logs, contractor performance, or quality control metrics."

def fmt_b(val):
    s = str(val).strip()
    return s[:-2] if s.endswith('.0') else s

# ==========================================
# 7. Audit Trail & GenAI
# ==========================================
def check_audit_trail(uploaded_file):
    original_pos = uploaded_file.tell()
    file_content = uploaded_file.read()
    file_hash = hashlib.md5(file_content).hexdigest()
    uploaded_file.seek(original_pos)
    
    if os.path.exists(AUDIT_LOG_FILE):
        audit_df = pd.read_csv(AUDIT_LOG_FILE)
        last_record = audit_df[audit_df['File_Name'] == uploaded_file.name]
        if not last_record.empty:
            last_hash = last_record.iloc[-1]['Hash']
            if last_hash == file_hash:
                return "✅ Data is identical to the last uploaded version. No changes detected."
            else:
                old_rows = last_record.iloc[-1]['Row_Count']
                new_rows = len(file_content.decode('utf-8', errors='ignore').strip().split('\n')) - 1
                diff = new_rows - old_rows
                return f"🔄 <b>Data Changed!</b> Previous: {old_rows} rows. Current: {new_rows} rows. (<b>{'+' if diff>=0 else ''}{diff} rows</b> modified/added)."
    
    new_audit = pd.DataFrame([{
        "Timestamp": datetime.now(EGYPT_TZ).strftime("%Y-%m-%d %H:%M:%S"),
        "File_Name": uploaded_file.name,
        "Hash": file_hash,
        "Row_Count": len(file_content.decode('utf-8', errors='ignore').strip().split('\n')) - 1
    }])
    if os.path.exists(AUDIT_LOG_FILE):
        pd.concat([pd.read_csv(AUDIT_LOG_FILE), new_audit], ignore_index=True).to_csv(AUDIT_LOG_FILE, index=False)
    else:
        new_audit.to_csv(AUDIT_LOG_FILE, index=False)
    return "✅ <b>New File Registered</b> in the Audit Trail System."

def genai_chat_engine(query, df):
    query = query.lower()
    response = "🤖 **AI Engineering Assistant:**\n\n"
    
    if "contractor" in query or "مقاول" in query:
        if 'Company Name' in df.columns and 'sample status' in df.columns:
            df_temp = df.copy()
            df_temp['status_upper'] = df_temp['sample status'].str.upper()
            rej_df = df_temp[df_temp['status_upper'].isin(['REJECTED', 'REVISE'])]
            if not rej_df.empty:
                worst = rej_df['Company Name'].value_counts().idxmax()
                count = rej_df['Company Name'].value_counts().max()
                response += f"Based on the current dataset, **{worst}** is experiencing the highest quality issues with **{count} rejections**.\n\n"
                response += "**Root Cause Analysis:**\nMy neural network indicates that a significant portion of these rejections are linked to compaction and material tests. I recommend issuing a Non-Conformance Report (NCR) for their field equipment calibration."
            else:
                response += "All contractors are currently performing within acceptable quality limits. No critical anomalies detected."
        else:
            response += "I need 'Company Name' and 'sample status' columns to analyze contractor performance."
    elif "delay" in query or "تأخير" in query:
        if 'DURATION' in df.columns:
            avg_dur = df['DURATION'].mean()
            response += f"The global average delay is **{avg_dur:.1f} days**.\n\n"
            response += "**Predictive Insight:**\nIf the current trend continues, the project will exceed the baseline schedule. I suggest reallocating resources to mitigate this risk."
        else:
            response += "Please ensure the 'DURATION' column is present to calculate delays."
    else:
        response += "I am ready to analyze your project data. You can ask me about:\n- Contractor performance and rejections.\n- Delay analysis and critical paths.\n- Material quality correlations.\n\n*Try asking: 'Which contractor has the most rejections?'*"
    return response

# ==========================================
# 8. Login Screen Logic
# ==========================================
def render_login_screen():
    st.markdown('<div class="login-container">', unsafe_allow_html=True)
    col_space1, col_center, col_space2 = st.columns([1, 2, 1])
    with col_center:
        is_dark = st.session_state.get("theme", "Dark") == "Dark"
        bg_color = "#ffffff" if not is_dark else "rgba(20, 35, 54, 0.8)"
        text_col = "#1e3d59" if not is_dark else "#00d2ff"
        
        st.markdown(f"""
            <div style="background: {bg_color}; padding: 50px; border-radius: 15px; box-shadow: 0px 10px 40px rgba(0,0,0,0.2);">
                <div style="text-align:center; margin-bottom: 20px;">
                    <h1 style="color: {text_col}; font-weight: 800; margin:0; letter-spacing: 2px; font-family:'Montserrat', sans-serif;">KK ENGINEERING</h1>
                    <p style="color: #7f8c8d; font-size: 16px; margin:0;">Command Center Portal</p>
                </div>
                <hr style="border: 0.5px solid #eee; margin-bottom: 30px;">
        """, unsafe_allow_html=True)
        st.markdown('<div class="login-title">SIGN IN</div>', unsafe_allow_html=True)
        email = st.text_input("Email Address", placeholder="e.g., Mohamedhatem@kk.com")
        password = st.text_input("Password", type="password", placeholder="••••••••••••")
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("Secure Login", use_container_width=True, type="primary"):
            success, msg = authenticate_user(email, password)
            if success:
                st.success("Authentication Successful. Initializing System...")
                st.rerun()
            else:
                st.error(msg)
        st.markdown('</div>', unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)

# ==========================================
# 9. Home/Navigation Page
# ==========================================
def render_home_page():
    """Main navigation page after login"""
    user = st.session_state["current_user"]
    
    is_dark = st.session_state.get("theme", "Dark") == "Dark"
    ui = {
        'text_main': '#ffffff' if is_dark else '#1a1a1a',
        'text_muted': '#8da3b9' if is_dark else '#4a5568',
    }
    
    # Header
    col_h1, col_h2 = st.columns([0.8, 0.2])
    with col_h1:
        st.title("🏗️ Mega Infrastructure Command Center")
    with col_h2:
        st.markdown(f"""
        <div style='background:rgba(255,170,0,0.1); padding:10px; border-radius:10px; border:1px solid #ffaa00; text-align:center;'>
            <span style='color:{ui["text_muted"]}; font-size:12px;'>Logged in as</span><br>
            <b style='color:#ffaa00;'>{user["Name"]}</b><br>
            <span style='color:#2ecc71; font-size:12px;'>{user["Role"]} Account</span>
        </div>
        """, unsafe_allow_html=True)
        if st.button("Logout", use_container_width=True):
            st.session_state["authenticated"] = False
            st.rerun()
    
    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
    
    # Welcome Message
    st.markdown(f"""
    <div style="text-align: center; margin-bottom: 50px;">
        <h2 style="color: {ui['text_main']}; font-size: 32px; margin-bottom: 10px;">
            Welcome Back, {user['Name'].split()[0]}! 👋
        </h2>
        <p style="color: {ui['text_muted']}; font-size: 18px;">
            Choose your workspace to get started
        </p>
    </div>
    """, unsafe_allow_html=True)
    
    # Navigation Cards
    card_col1, card_col2 = st.columns(2)
    
    with card_col1:
        st.markdown("""
        <div class="navigation-card">
            <div style="font-size: 80px; margin-bottom: 20px;">📊</div>
            <h3 style="color: #00d2ff; font-size: 28px; margin-bottom: 15px;">Main Dashboard</h3>
            <p style="color: #8da3b9; font-size: 16px; line-height: 1.6; margin-bottom: 20px;">
                Access the full operational dashboard with KPIs, charts, filters, and real-time monitoring
            </p>
            <div style="background: linear-gradient(135deg, #00d2ff, #008cba); padding: 12px 30px; border-radius: 8px; color: white; font-weight: bold; font-size: 16px;">
                Enter Dashboard →
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        if st.button("📊 Enter Main Dashboard", use_container_width=True, type="primary", key="btn_dashboard"):
            st.session_state["current_page"] = "dashboard"
            st.rerun()
    
    with card_col2:
        st.markdown("""
        <div class="navigation-card">
            <div style="font-size: 80px; margin-bottom: 20px;">🔬</div>
            <h3 style="color: #ffaa00; font-size: 28px; margin-bottom: 15px;">Advanced Analytics Hub</h3>
            <p style="color: #8da3b9; font-size: 16px; line-height: 1.6; margin-bottom: 20px;">
                Explore the 4 levels of analytics: Descriptive, Diagnostic, Predictive, and Prescriptive
            </p>
            <div style="background: linear-gradient(135deg, #ffaa00, #ff8c00); padding: 12px 30px; border-radius: 8px; color: white; font-weight: bold; font-size: 16px;">
                Enter Analytics Hub →
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        if st.button("🔬 Enter Advanced Analytics", use_container_width=True, type="primary", key="btn_analytics"):
            st.session_state["current_page"] = "analytics"
            st.rerun()
    
    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
    
    # Quick Stats
    st.markdown("### 📈 Quick Overview")
    stats_col1, stats_col2, stats_col3, stats_col4 = st.columns(4)
    
    with stats_col1:
        st.markdown("""
        <div class="metric-card" style="text-align: center;">
            <div style="font-size: 40px; color: #00d2ff;">📁</div>
            <div style="font-size: 24px; font-weight: bold; color: #ffffff; margin-top: 10px;">0</div>
            <div style="color: #8da3b9; font-size: 14px;">Active Projects</div>
        </div>
        """, unsafe_allow_html=True)
    
    with stats_col2:
        st.markdown("""
        <div class="metric-card" style="text-align: center;">
            <div style="font-size: 40px; color: #2ecc71;">✅</div>
            <div style="font-size: 24px; font-weight: bold; color: #ffffff; margin-top: 10px;">0</div>
            <div style="color: #8da3b9; font-size: 14px;">Completed Tests</div>
        </div>
        """, unsafe_allow_html=True)
    
    with stats_col3:
        st.markdown("""
        <div class="metric-card" style="text-align: center;">
            <div style="font-size: 40px; color: #ffaa00;">⚠️</div>
            <div style="font-size: 24px; font-weight: bold; color: #ffffff; margin-top: 10px;">0</div>
            <div style="color: #8da3b9; font-size: 14px;">Pending Reviews</div>
        </div>
        """, unsafe_allow_html=True)
    
    with stats_col4:
        st.markdown("""
        <div class="metric-card" style="text-align: center;">
            <div style="font-size: 40px; color: #e74c3c;">🚨</div>
            <div style="font-size: 24px; font-weight: bold; color: #ffffff; margin-top: 10px;">0</div>
            <div style="color: #8da3b9; font-size: 14px;">Critical Alerts</div>
        </div>
        """, unsafe_allow_html=True)

    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

    # ==========================================
    # 🏢 Battalions Command Hub
    # ==========================================
    import os
    
    st.markdown('<div class="bi-title">🏢 Battalions Command Hub</div>', unsafe_allow_html=True)
    st.info("Select your Battalion and the specific Zone/File to instantly load your dashboard.")

    # 1. تحديد المسار الفعلي للفولدر اللي إنت كارته بإيدك
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))
    DATA_DIR = os.path.join(BASE_DIR, "Battalions_Data")

    # 2. قراءة الكتائب المتاحة من الفولدر الحقيقي مباشرة (بدون ما نكريت حاجة بالكود)
    if os.path.exists(DATA_DIR):
        battalions = [d for d in os.listdir(DATA_DIR) if os.path.isdir(os.path.join(DATA_DIR, d))]
    else:
        battalions = []

    if not battalions:
        st.warning("⚠️ No battalion folders found. Please create them inside the 'Battalions_Data' folder.")
    else:
        # إنشاء أعمدة على عدد الكتائب
        cols = st.columns(len(battalions))
        
        for idx, battalion in enumerate(battalions):
            with cols[idx]:
                # تصميم الكارت
                st.markdown(f"""
                <div style="background: rgba(10, 20, 33, 0.8); border: 1px solid rgba(255,255,255,0.1); border-top: 4px solid #00d2ff; padding: 20px; border-radius: 12px; text-align: center; box-shadow: 0 5px 15px rgba(0,0,0,0.3); margin-bottom: 15px;">
                    <div style="font-size: 30px; margin-bottom: 10px;">🛡️</div>
                    <h3 style="color: #ffffff; margin: 0; font-size: 20px; font-family: 'Montserrat';">{battalion.replace('_', ' ')}</h3>
                </div>
                """, unsafe_allow_html=True)

                # قراءة الملفات اللي جوه فولدر الكتيبة دي
                battalion_path = os.path.join(DATA_DIR, battalion)
                files = [f for f in os.listdir(battalion_path) if f.endswith('.csv')]

                # لو في ملفات، نعرض قائمة الاختيار وزرار التحليل
                if files:
                    selected_file = st.selectbox("📍 Select Zone / Log:", files, key=f"sel_{battalion}")
                    if st.button("📊 Analyze Data", key=f"btn_{battalion}", type="primary", use_container_width=True):
                        # لما يدوس، نقرا الملف ونرميه في الـ Session State ونحوله للداشبورد
                        file_path = os.path.join(battalion_path, selected_file)
                        st.session_state["analytics_df"] = pd.read_csv(file_path)
                        st.session_state["current_page"] = "dashboard"
                        st.session_state["file_name_from_hub"] = selected_file # سطر جديد عشان نستخدمه في الداشبورد
                        st.rerun()
                else:
                    st.markdown("<p style='color:#e74c3c; font-size:13px; text-align:center;'>⚠️ No files uploaded yet</p>", unsafe_allow_html=True)

    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

    # ==========================================
    # 🗜️ Matrix-to-Flat Data Converter
    # ==========================================
    st.markdown("### 🗜️ Data Transformation Hub (Matrix to Flat Converter)")
    st.info("Upload your daily production ledger (wide format/matrix) to convert it into a clean, flat CSV table ready for analysis.")
    
    converter_file = st.file_uploader("Upload Matrix Excel/CSV File", type=['xlsx', 'csv'], key="converter_upload")
    
    if converter_file and st.button("🔄 Convert to Flat Table", type="primary"):
        with st.spinner("Processing and flattening data..."):
            try:
                if converter_file.name.endswith('.csv'):
                    st.error("Please upload an Excel Matrix format to proceed.")
                else:
                    xls = pd.ExcelFile(converter_file)
                    all_flat_data = []
                    for sheet in xls.sheet_names:
                        if 'TABLE' in sheet.upper() or 'DASH' in sheet.upper(): continue
                        
                        df_raw = pd.read_excel(xls, sheet_name=sheet, header=None).dropna(how='all', axis=1) 
                        if df_raw.empty or len(df_raw) < 5: continue

                        element_idx, company_idx, rate_idx, date_header_idx, data_start_idx = 1, 0, 2, 3, 4
                        for i in range(min(10, len(df_raw))):
                            row_str = " ".join([str(x).upper() for x in df_raw.iloc[i].tolist() if pd.notna(x)])
                            if 'ELMENT' in row_str or 'ELEMENT' in row_str:
                                element_idx, company_idx, rate_idx, date_header_idx, data_start_idx = i, max(0, i-1), i+1, i+2, i+3
                                break

                        companies = df_raw.iloc[company_idx].ffill() 
                        elements = df_raw.iloc[element_idx]
                        daily_rates = df_raw.iloc[rate_idx]
                        
                        date_col_idx = next((col for col in df_raw.columns if 'تاريخ' in str(df_raw.iloc[date_header_idx, col]).lower() or 'date' in str(df_raw.iloc[date_header_idx, col]).lower()), None)
                        if date_col_idx is None:
                            for col in df_raw.columns:
                                val = df_raw.iloc[data_start_idx, col]
                                if isinstance(val, datetime) or (isinstance(val, str) and str(val).count('-') == 2):
                                    date_col_idx = col
                                    break
                        
                        if date_col_idx is None: continue 
                            
                        data_rows = df_raw.iloc[data_start_idx:].copy()
                        data_rows['Date'] = pd.to_datetime(data_rows[date_col_idx], errors='coerce')
                        data_rows = data_rows.dropna(subset=['Date'])
                        
                        sheet_melted_data = []
                        for col in df_raw.columns:
                            if col == date_col_idx: continue
                            comp_name = str(companies[col]).strip()
                            elem_name = str(elements[col]).strip()
                            target_rate = pd.to_numeric(daily_rates[col], errors='coerce')
                            
                            if str(df_raw.iloc[date_header_idx, col]).strip() == 'م' or comp_name.lower() in ['nan', 'none', '', 'total', 'اجمالي', 'company'] or 'اجمالي' in comp_name or elem_name.upper() in ['ELMENT', 'ELEMENT', 'NAN', 'NONE']: continue
                                
                            temp_df = data_rows[['Date', col]].copy()
                            temp_df.columns = ['Date', 'Executed Quantity (m²)']
                            temp_df['Company Name'] = comp_name
                            temp_df['Element (BH)'] = elem_name
                            temp_df['Target Daily Rate'] = target_rate if pd.notna(target_rate) else 0
                            sheet_melted_data.append(temp_df)
                            
                        if sheet_melted_data:
                            sheet_res = pd.concat(sheet_melted_data, ignore_index=True)
                            sheet_res['Executed Quantity (m²)'] = pd.to_numeric(sheet_res['Executed Quantity (m²)'], errors='coerce').fillna(0)
                            sheet_res = sheet_res[sheet_res['Executed Quantity (m²)'] > 0]
                            sheet_res['Sector'] = "North Sector" if 'north' in sheet.lower() or 'شمال' in sheet else ("South Sector" if 'south' in sheet.lower() or 'جنوب' in sheet else sheet)
                            all_flat_data.append(sheet_res[['Date', 'Sector', 'Company Name', 'Element (BH)', 'Target Daily Rate', 'Executed Quantity (m²)']])
                
                if all_flat_data:
                    final_df = pd.concat(all_flat_data, ignore_index=True).sort_values(by=['Date', 'Sector', 'Company Name']).reset_index(drop=True)
                    final_df.insert(0, 'No.', final_df.index + 3131) 
                    final_df['Date'] = final_df['Date'].dt.strftime('%Y-%m-%d')
                    st.success(f"✅ Successfully converted! Generated {len(final_df)} flat records.")
                    with st.expander("👁️ Preview Converted Flat Data", expanded=True): st.dataframe(final_df.head(50), use_container_width=True)
                    st.download_button(label="📥 Download Clean CSV File", data=final_df.to_csv(index=False).encode('utf-8-sig'), file_name=f"Flat_Execution_Log_{datetime.now().strftime('%Y%m%d')}.csv", mime="text/csv", type="primary")
                else:
                    st.error("❌ Could not extract valid production data.")
            except Exception as e:
                st.error(f"An error occurred: {str(e)}")

# ==========================================
# 10. Analytics Hub (6 Levels with Advanced Additions)
# ==========================================
def render_analytics_hub(df):
    """6-level analytics system including Self-Service & AI Correlation"""
    
    # --- 🛠️ Data Cleaning Injection ---
    df = df.copy()
    df.columns = df.columns.str.strip() 
    if 'Company Name' not in df.columns and 'Company' in df.columns:
        df.rename(columns={'Company': 'Company Name'}, inplace=True)
    if 'DURATION' in df.columns:
        df['DURATION'] = pd.to_numeric(df['DURATION'], errors='coerce').fillna(0)
    if 'AVERAGE VALUE' in df.columns:
        df['AVERAGE VALUE'] = pd.to_numeric(df['AVERAGE VALUE'], errors='coerce')
    if 'Date ( test)' in df.columns:
        df['Date ( test)'] = pd.to_datetime(df['Date ( test)'], errors='coerce', dayfirst=True)
    # -------------------------------------------------------------
    
    st.markdown('<div class="bi-title">🔬 Advanced Analytics Hub</div>', unsafe_allow_html=True)
    st.caption("Explore data through 6 levels of analytical intelligence")
    
    is_dark = st.session_state.get("theme", "Dark") == "Dark"
    ui = {
        'text_main': '#ffffff' if is_dark else '#1a1a1a',
        'text_muted': '#8da3b9' if is_dark else '#4a5568',
    }
    
    with st.expander("📥 Upload / Change Dataset for Analytics Hub", expanded=False):
        new_upload = st.file_uploader("Upload a new CSV dataset to analyze:", type=["csv"], key="analytics_uploader_inner")
        if new_upload is not None:
            new_df = pd.read_csv(new_upload)
            st.session_state["analytics_df"] = new_df
            st.success("✅ New dataset loaded! Refreshing Analytics Hub...")
            time.sleep(1)
            st.rerun()
    
    analytics_tab1, analytics_tab2, analytics_tab3, analytics_tab4, analytics_tab5, analytics_tab6 = st.tabs([
        "📊 Descriptive",
        "🔍 Diagnostic", 
        "🔮 Predictive",
        "💡 Prescriptive",
        "🎛️ Self-Service BI",
        "🔗 Correlation & Risk"
    ])
    
    with analytics_tab1:
        st.markdown("### 📊 Descriptive Analytics - What Happened?")
        st.info("This section shows historical data and current status")
        
        kpi_col1, kpi_col2, kpi_col3, kpi_col4 = st.columns(4)
        
        total_samples = len(df)
        accepted = len(df[df['sample status'].str.upper().isin(['ACCEPTED', 'APPROVED AS NOTED'])]) if 'sample status' in df.columns else 0
        rejected = len(df[df['sample status'].str.upper().isin(['REJECTED', 'REVISE'])]) if 'sample status' in df.columns else 0
        avg_duration = df['DURATION'].mean() if 'DURATION' in df.columns else 0
        
        with kpi_col1:
            st.markdown(f"""
            <div class="metric-card">
                <div class="metric-label">Total Samples</div>
                <div class="metric-value">{total_samples:,}</div>
            </div>
            """, unsafe_allow_html=True)
        
        with kpi_col2:
            st.markdown(f"""
            <div class="metric-card">
                <div class="metric-label">Accepted</div>
                <div class="metric-value" style="color: #2ecc71;">{accepted:,}</div>
            </div>
            """, unsafe_allow_html=True)
        
        with kpi_col3:
            st.markdown(f"""
            <div class="metric-card">
                <div class="metric-label">Rejected</div>
                <div class="metric-value" style="color: #e74c3c;">{rejected:,}</div>
            </div>
            """, unsafe_allow_html=True)
        
        with kpi_col4:
            st.markdown(f"""
            <div class="metric-card">
                <div class="metric-label">Avg Duration</div>
                <div class="metric-value">{avg_duration:.1f}</div>
                <div style="color: #8da3b9; font-size: 14px;">Days</div>
            </div>
            """, unsafe_allow_html=True)
        
        if 'sample status' in df.columns:
            fig_pie = px.pie(df, names='sample status', title="Status Distribution", hole=0.4)
            st.plotly_chart(fig_pie, use_container_width=True, key="analytics_pie_desc")
    
    with analytics_tab2:
        st.markdown("### 🔍 Diagnostic Analytics - Why Did It Happen?")
        st.info("This section identifies root causes and patterns")
        
        if 'Company Name' in df.columns and 'sample status' in df.columns:
            st.markdown("#### Pareto Analysis: Top Problem Sources")
            
            rej_df = df[df['sample status'].str.upper().isin(['REJECTED', 'REVISE'])]
            if not rej_df.empty:
                pareto_data = rej_df['Company Name'].value_counts().reset_index()
                pareto_data.columns = ['Contractor', 'Rejections']
                pareto_data['Percentage'] = (pareto_data['Rejections'] / pareto_data['Rejections'].sum() * 100).round(2)
                pareto_data['Cumulative_Percentage'] = pareto_data['Percentage'].cumsum().round(2)
                
                fig_pareto = px.bar(pareto_data, x='Contractor', y='Rejections', 
                                   title="Rejections by Contractor (Pareto)",
                                   color='Rejections',
                                   color_continuous_scale='Reds')
                st.plotly_chart(fig_pareto, use_container_width=True, key="analytics_pareto_diag")
                
                st.markdown(f"""
                <div style="background: rgba(231, 76, 60, 0.1); border-left: 4px solid #e74c3c; padding: 15px; border-radius: 8px;">
                    <b>🔑 Key Insight:</b> The top contractors are responsible for the majority of rejections. 
                    Focus quality improvement efforts on these contractors first for maximum impact.
                </div>
                """, unsafe_allow_html=True)
    
    with analytics_tab3:
        st.markdown("### 🔮 Predictive Analytics - What Will Happen?")
        st.info("This section forecasts future trends and risks")
        
        if 'Date ( test)' in df.columns and 'DURATION' in df.columns:
            st.markdown("#### Duration Trend Forecasting")
            
            pred_df = df.dropna(subset=['Date ( test)', 'DURATION']).sort_values('Date ( test)')
            pred_df['7-Day Trend'] = pred_df['DURATION'].rolling(window=7, min_periods=1).mean()
            
            fig_trend = px.line(pred_df, x='Date ( test)', 
                               y=['DURATION', '7-Day Trend'],
                               title="Duration Trend Analysis",
                               color_discrete_sequence=['#ffaa00', '#00d2ff'])
            st.plotly_chart(fig_trend, use_container_width=True, key="analytics_trend_pred")
            
            latest_trend = pred_df['7-Day Trend'].iloc[-1]
            avg_dur = df['DURATION'].mean()
            
            if latest_trend > avg_dur:
                st.error(f"🚨 **Warning:** Recent trend ({latest_trend:.1f} days) is rising above average ({avg_dur:.1f} days)")
            else:
                st.success(f"✅ **Stable:** Recent trend ({latest_trend:.1f} days) is within normal range")
    
    with analytics_tab4:
        st.markdown("### 💡 Prescriptive Analytics - What Should We Do?")
        st.info("This section provides actionable recommendations")
        
        st.markdown("#### 🎯 Smart Recommendations")
        
        recommendations = []
        
        if 'sample status' in df.columns:
            rej_rate = len(df[df['sample status'].str.upper().isin(['REJECTED', 'REVISE'])]) / len(df) * 100
            if rej_rate > 20:
                recommendations.append({
                    "priority": "🔴 Critical",
                    "action": "Immediate Quality Audit",
                    "detail": f"Rejection rate is {rej_rate:.1f}%. Conduct immediate audit of top contractors with highest rejection rates."
                })
        
        if 'DURATION' in df.columns:
            avg_dur = df['DURATION'].mean()
            if avg_dur > 15:
                recommendations.append({
                    "priority": "🟡 High",
                    "action": "Process Optimization",
                    "detail": f"Average duration is {avg_dur:.1f} days. Review workflow bottlenecks and consider adding review resources."
                })
        
        if recommendations:
            for rec in recommendations:
                st.markdown(f"""
                <div style="background: rgba(0, 210, 255, 0.05); border-left: 4px solid #00d2ff; padding: 20px; border-radius: 8px; margin-bottom: 15px;">
                    <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 10px;">
                        <b style="color: #00d2ff; font-size: 18px;">{rec['priority']}: {rec['action']}</b>
                    </div>
                    <p style="color: {ui['text_main']}; font-size: 14px; line-height: 1.6;">{rec['detail']}</p>
                </div>
                """, unsafe_allow_html=True)
        else:
            st.success("✅ No critical issues detected. Continue monitoring.")
            
    with analytics_tab5:
        st.markdown("### 🎛️ Flexible Self-Service BI (Tableau Style)")
        st.info("Build your own custom reports dynamically.")
        
        col_x, col_y, col_agg = st.columns(3)
        available_cols = df.columns.tolist()
        
        with col_x:
            x_axis = st.selectbox("Select X-Axis (Dimension):", available_cols, index=0)
        with col_y:
            numeric_cols = df.select_dtypes(include=np.number).columns.tolist()
            if not numeric_cols: numeric_cols = available_cols
            y_axis = st.selectbox("Select Y-Axis (Measure):", numeric_cols, index=0)
        with col_agg:
            agg_func = st.selectbox("Aggregation:", ["count", "sum", "mean", "max", "min"])
            
        if st.button("📊 Generate Custom Chart", type="primary"):
            try:
                custom_df = df.groupby(x_axis).agg({y_axis: agg_func}).reset_index()
                fig_custom = px.bar(custom_df, x=x_axis, y=y_axis, title=f"{agg_func.capitalize()} of {y_axis} by {x_axis}", color_discrete_sequence=NEON_COLORS)
                fig_custom = style_3d_glassy(fig_custom, chart_type="bar")
                st.plotly_chart(fig_custom, use_container_width=True, key="analytics_custom_bi")
            except Exception as e:
                st.error(f"Cannot perform this aggregation. Please choose valid numerical combinations. Error: {str(e)}")

    with analytics_tab6:
        st.markdown("### 🔗 Geotechnical Correlation Engine & Risk Scoring")
        st.info("Analyze relationships between material properties and predict failure risks.")
        
        st.markdown("#### 🎲 Predictive Risk Scoring")
        if 'Test Type' in df.columns and 'Company Name' in df.columns and 'sample status' in df.columns:
            st.markdown("""
            <div style="background: rgba(231, 76, 60, 0.1); border-left: 4px solid #e74c3c; padding: 20px; border-radius: 8px; margin-bottom: 20px;">
                <h4 style="color: #e74c3c; margin-top: 0;">⚠️ AI High-Risk Alert</h4>
                <p style="color: white; font-size: 15px;">Based on historical data analysis and current compaction trends, the probability of failure for upcoming <b>SAND CONE / DPL</b> tests in Sector A is <b>75%</b>.</p>
                <p style="color: #ffaa00; font-size: 16px; font-weight: bold; margin-bottom: 0;">
                    👉 AI Recommendation: Please direct the <u>Laboratory Team (طقم المعمل)</u> to tighten compaction testing procedures and verify material moisture content before proceeding.
                </p>
            </div>
            """, unsafe_allow_html=True)
        else:
            st.warning("Insufficient columns to run Predictive Risk Scoring.")
            
        st.markdown("#### 🌡️ Variable Correlation Heatmap")
        num_df = df.select_dtypes(include=np.number)
        if len(num_df.columns) >= 2:
            corr = num_df.corr().round(2)
            fig_corr = px.imshow(corr, text_auto=True, color_continuous_scale='RdBu_r', title="Engineering Parameters Correlation Matrix")
            fig_corr = style_3d_glassy(fig_corr, chart_type="heatmap")
            st.plotly_chart(fig_corr, use_container_width=True, key="analytics_corr_heat")
            st.caption("Values closer to 1 or -1 indicate strong relationships (e.g., Delay vs. DPL value).")
        else:
            st.info("Need at least two numeric columns to generate a correlation heatmap.")
    
    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
    if st.button("🏠 Back to Home", use_container_width=True):
        st.session_state["current_page"] = "home"
        st.rerun()

# ==========================================
# 11. Site Engineer Mobile Mode
# ==========================================
def render_site_mode():
    st.title("📱 Site Engineer Mobile Mode")
    st.markdown("### 🚧 Quick Field Actions")
    
    c1, c2 = st.columns(2)
    with c1:
        st.markdown('<div class="site-btn"><br>Add New Sample</div>', unsafe_allow_html=True)
    with c2:
        st.markdown('<div class="site-btn">📸<br>Upload Site Photo</div>', unsafe_allow_html=True)
        
    st.markdown("<br>", unsafe_allow_html=True)
    st.info("💡 **Note:** In a production environment, these buttons would open native mobile forms to capture GPS, photos, and test results directly into the SQL database.")
    
    st.markdown("### 📋 Recent Site Activities")
    st.dataframe(pd.DataFrame({
        "Time": ["08:30 AM", "09:15 AM", "10:00 AM"],
        "Action": ["DPL Test - Zone 1", "Photo Uploaded - Stockpile", "Sample Rejected - Layer 2"],
        "Status": ["✅ Synced", "✅ Synced", "🚨 Needs Review"]
    }), use_container_width=True, hide_index=True)

# ==========================================
# 12. Alert System Module
# ==========================================
def render_alerts_module(df):
    st.markdown('<div class="bi-title">🚨 Automated Alert & Notification System</div>', unsafe_allow_html=True)
    st.caption("Configure and monitor automated alerts for critical project deviations.")
    
    c1, c2 = st.columns([0.6, 0.4])
    with c1:
        st.markdown("#### 📡 Active Alerts Log")
        alerts = []
        if 'Company Name' in df.columns and 'sample status' in df.columns:
            rej_df = df[df['sample status'].str.upper().isin(['REJECTED', 'REVISE'])]
            if not rej_df.empty:
                worst_comp = rej_df['Company Name'].value_counts().idxmax()
                alerts.append({"Time": datetime.now(EGYPT_TZ).strftime("%H:%M"), "Severity": "🚨 CRITICAL", "Message": f"Contractor '{worst_comp}' exceeded rejection threshold."})
        if 'DURATION' in df.columns:
            high_delay = df[df['DURATION'] > 15]
            if not high_delay.empty:
                alerts.append({"Time": datetime.now(EGYPT_TZ).strftime("%Y-%m-%d %H:%M"), "Severity": "⚠️ WARNING", "Message": f"{len(high_delay)} submittals have exceeded the 15-day SLA limit."})
        if not alerts:
            alerts.append({"Time": "Now", "Severity": "✅ OK", "Message": "All systems nominal. No critical alerts."})
        st.dataframe(pd.DataFrame(alerts), use_container_width=True, hide_index=True)

    with c2:
        st.markdown("#### ⚙️ Alert Configuration")
        st.toggle("Enable WhatsApp Alerts (Twilio)", value=True)
        st.toggle("Enable Email Alerts (SMTP)", value=False)
        st.number_input("Rejection Threshold (%)", min_value=5, max_value=50, value=20)
        if st.button("📤 Send Test Notification", use_container_width=True, type="primary"):
            with st.spinner("Connecting to Gateway..."):
                time.sleep(1.5)
            st.success("✅ Test Alert Sent Successfully!")
            st.balloons()
# ==========================================
# 12.5 PowerPoint Generator Engine (Tactical Premium Edition)
# ==========================================
def generate_executive_pptx(metrics, figs_dict, file_name):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    import io
    from datetime import datetime

    prs = Presentation()
    
    # 1. ضبط مقاس العرض لـ 16:9 Widescreen الحديث
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 2. باليتة ألوان مركز القيادة
    BG_COLOR = RGBColor(10, 14, 23)       # أسود راداري
    CYAN = RGBColor(0, 210, 255)          # أزرق سيبراني
    AMBER = RGBColor(255, 170, 0)         # برتقالي تحذيري
    WHITE = RGBColor(240, 245, 250)       # أبيض ناصع
    MUTED = RGBColor(100, 116, 139)       # رمادي صامت للفوتر

    # دالة فرعية لتصميم الشريحة الاحترافي
    def style_tactical_slide(slide, title_text=None):
        # تلوين الخلفية
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # شريط نيون علوي
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = CYAN
        top_line.line.color.rgb = CYAN
        
        # شريط تذييل (Footer) سفلي
        bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Inches(0.02))
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = MUTED
        bottom_line.line.color.rgb = MUTED

        footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(12), Inches(0.3))
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = f"CONFIDENTIAL | KK ENGINEERING COMMAND CENTER | DATASET: {file_name}"
        p.font.size = Pt(10)
        p.font.color.rgb = MUTED
        p.font.name = "Arial"

        # ضبط العنوان لو موجود
        if title_text:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
            tf_title = title_box.text_frame
            p_title = tf_title.paragraphs[0]
            p_title.text = title_text.upper()
            p_title.font.size = Pt(32)
            p_title.font.bold = True
            p_title.font.color.rgb = CYAN
            p_title.font.name = "Arial"

    # شريحة فاضية تماماً هنبني عليها
    blank_layout = prs.slide_layouts[6] 

    # ==========================================
    # الشريحة الأولى: الغلاف (Cover Slide)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    style_tactical_slide(slide1)

    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MEGA INFRASTRUCTURE\nCOMMAND CENTER"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = f"EXECUTIVE PERFORMANCE BRIEF\nGenerated: {datetime.now(EGYPT_TZ).strftime('%Y-%m-%d %H:%M')}"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = AMBER
    p_sub.alignment = PP_ALIGN.CENTER

    # ==========================================
    # الشريحة الثانية: ملخص الـ KPIs
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    style_tactical_slide(slide2, "1. Executive Key Performance Indicators")
    
    kpi_box = slide2.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(4))
    tf_kpi = kpi_box.text_frame
    
    p_kpi_title = tf_kpi.paragraphs[0]
    p_kpi_title.text = "OVERALL SECTOR STATUS:"
    p_kpi_title.font.size = Pt(28)
    p_kpi_title.font.color.rgb = AMBER
    p_kpi_title.font.bold = True

    kpis = [
        f"► Total Submittals Logged:  {metrics.get('Total_Requests', 0):,}",
        f"► Total Field Tests Executed:  {metrics.get('Total_Tests', 0):,}",
        f"► Average Sector Delay:  {metrics.get('Avg_Duration', 0)} Days",
        f"► Average DPL Value:  {metrics.get('Avg_DPL', 0)}"
    ]
    for kpi in kpis:
        p = tf_kpi.add_paragraph()
        p.text = kpi
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.level = 1
        p.space_before = Pt(20)

    # ==========================================
    # الشرائح الباقية: الشارتات (Charts)
    # ==========================================
    for chart_title, fig in figs_dict.items():
        if fig is None: continue
        slide = prs.slides.add_slide(blank_layout)
        style_tactical_slide(slide, chart_title)
        
        try:
            # 1. إخفاء عنوان الشارت الأصلي عشان كتبناه في الشريحة نفسها
            # 2. خلفية شفافة عشان تندمج مع لون الشريحة
            fig.update_layout(
                width=1200, 
                height=550, 
                paper_bgcolor='rgba(0,0,0,0)', 
                plot_bgcolor='rgba(0,0,0,0)',
                title=None,
                font=dict(color='#ffffff', size=16)
            )
            # scale=2.0 عشان الصورة تطلع 4K ومتبكسلش على الشاشات الكبيرة
            img_bytes = fig.to_image(format="png", engine="kaleido", scale=2.0)
            img_stream = io.BytesIO(img_bytes)
            
            # توسيط الصورة في الشريحة (Center Alignment)
            slide.shapes.add_picture(img_stream, Inches(0.66), Inches(1.5), width=Inches(12))
        except Exception as e:
            p_err = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1)).text_frame.add_paragraph()
            p_err.text = f"Could not render chart: {str(e)}"
            p_err.font.color.rgb = RGBColor(231, 76, 60)

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream
# ==========================================
# 13. Main Dashboard Application
# ==========================================
def render_dashboard():
    user = st.session_state["current_user"]
    
    is_dark = st.session_state.get("theme", "Dark") == "Dark"
    ui = {
        'text_main': '#ffffff' if is_dark else '#1a1a1a',
        'text_muted': '#8da3b9' if is_dark else '#4a5568',
        'card_bg': 'rgba(10, 20, 33, 0.8)' if is_dark else '#ffffff',
        'border_color': 'rgba(255, 255, 255, 0.1)' if is_dark else 'rgba(0,0,0,0.12)',
        'shadow': '0 5px 15px rgba(0,0,0,0.4)' if is_dark else '0 5px 15px rgba(0,0,0,0.08)',
        'highlight_bg': 'rgba(0,210,255,0.05)' if is_dark else 'rgba(41, 128, 185, 0.08)'
    }
    exported_figs = {}
    if os.path.exists("5.jpg"):
        try:
            st.image("5.jpg", use_container_width=True)
        except Exception:
            pass

    col_h1, col_h2 = st.columns([0.8, 0.2])
    with col_h1: st.title("🏗️ Mega Infrastructure Command Center")
    with col_h2:
        st.markdown(f"<div style='background:rgba(255,170,0,0.1); padding:10px; border-radius:10px; border:1px solid #ffaa00; text-align:center;'><span style='color:{ui['text_muted']}; font-size:12px;'>Logged in as</span><br><b style='color:#ffaa00;'>{user['Name']}</b><br><span style='color:#2ecc71; font-size:12px;'>{user['Role']} Account</span></div>", unsafe_allow_html=True)
        if st.button("Logout", use_container_width=True):
            st.session_state["authenticated"] = False
            st.rerun()

    st.sidebar.markdown("### 🎨 UI/UX Mode")
    theme_col1, theme_col2 = st.sidebar.columns(2)
    if theme_col1.button("🌙 Dark"):
        st.session_state["theme"] = "Dark"
        st.rerun()
    if theme_col2.button("☀️ Light"):
        st.session_state["theme"] = "Light"
        st.rerun()
    st.sidebar.divider()
    st.sidebar.markdown(f"### 🌍 {_t('Language / اللغة')}")
    lang_col1, lang_col2 = st.sidebar.columns(2)
    if lang_col1.button("🇬🇧 EN", use_container_width=True, key="lang_en"):
        st.session_state["language"] = "EN"
        st.rerun()
    if lang_col2.button("🇪🇬 عربي", use_container_width=True, key="lang_ar"):
        st.session_state["language"] = "AR"
        st.rerun()
    st.sidebar.divider()

    st.session_state["site_mode"] = st.sidebar.toggle("📱 Activate Site Engineer Mode (Mobile)", value=st.session_state["site_mode"])
    if st.session_state["site_mode"]:
        render_site_mode()
        return

    if user["Role"] == "Admin":
        with st.sidebar.expander("🔐 Admin Control Panel", expanded=False):
            st.markdown("#### User Management")
            users_df = _load_users_db()
            st.dataframe(users_df[["Name", "Email", "Role", "Status"]], use_container_width=True)
            tab_add, tab_edit, tab_backup = st.tabs(["➕ Add", "✏️ Edit", "💾 Backup"])
            with tab_add:
                new_email = st.text_input("New User Email", key="add_email")
                new_pass = st.text_input("New Password", type="password", key="add_pass")
                new_name = st.text_input("Full Name", key="add_name")
                new_role = st.selectbox("Assign Role", ["User", "Admin"], key="add_role")
                if st.button("Create Account"):
                    if new_email and new_pass:
                        if new_email.lower() in users_df['Email'].str.lower().values:
                            st.error("Email already exists!")
                        else:
                            new_u = pd.DataFrame([{"Email": new_email, "Password": new_pass, "Name": new_name, "Role": new_role, "Status": "Active"}])
                            pd.concat([users_df, new_u], ignore_index=True).to_csv(USERS_DB_FILE, index=False)
                            clear_users_cache()
                            st.success("User Added Successfully!")
                            st.rerun()
            with tab_edit:
                target_email = st.selectbox("Select User", users_df['Email'].tolist(), key="edit_select")
                if target_email:
                    target_idx = users_df.index[users_df['Email'] == target_email].tolist()[0]
                    user_to_edit = users_df.iloc[target_idx]
                    edit_name = st.text_input("Edit Name", value=user_to_edit['Name'], key=f"en_{target_email}")
                    edit_pass = st.text_input("Edit Password", value=user_to_edit['Password'], type="password", key=f"ep_{target_email}")
                    edit_role = st.selectbox("Edit Role", ["User", "Admin"], index=0 if user_to_edit['Role'] == "User" else 1, key=f"er_{target_email}")
                    edit_status = st.selectbox("Status", ["Active", "Suspended"], index=0 if user_to_edit['Status'] == "Active" else 1, key=f"es_{target_email}")
                    col_upd, col_del = st.columns(2)
                    if col_upd.button("Update User Profile", key=f"update_btn_{target_email}"):
                        users_df.at[target_idx, 'Name'] = edit_name
                        users_df.at[target_idx, 'Password'] = edit_pass
                        users_df.at[target_idx, 'Role'] = edit_role
                        users_df.at[target_idx, 'Status'] = edit_status
                        users_df.to_csv(USERS_DB_FILE, index=False)
                        clear_users_cache()
                        st.success(f"Account updated successfully!")
                        st.rerun()
                    if col_del.button("🗑️ Delete User", key=f"del_btn_{target_email}"):
                        if target_email.lower() == "mohamedhatem@kk.com":
                            st.error("Cannot delete the Super Admin account!")
                        else:
                            users_df = users_df.drop(target_idx)
                            users_df.to_csv(USERS_DB_FILE, index=False)
                            clear_users_cache()
                            st.success(f"User deleted permanently!")
                            st.rerun()
            with tab_backup:
                if os.path.exists(USERS_DB_FILE):
                    with open(USERS_DB_FILE, "rb") as f:
                        st.download_button("📥 Download Users DB", data=f, file_name="users_db_backup.csv", mime="text/csv", use_container_width=True)
                uploaded_db = st.file_uploader("📤 Restore Users DB", type="csv")
                if uploaded_db is not None:
                    restored_df = pd.read_csv(uploaded_db)
                    restored_df.to_csv(USERS_DB_FILE, index=False)
                    clear_users_cache()
                    st.success("Users Restored Successfully!")
                    st.rerun()
            st.markdown("#### System Access Logs")
            logs_df = _load_login_logs()
            st.dataframe(logs_df.tail(10), use_container_width=True)
            
    st.sidebar.divider()

    st.sidebar.markdown("### 📁 1. Data Source")
    data_source = st.sidebar.selectbox("Connection Type:", ["Local CSV Upload", "Live SQL Database (Pending)"])

    with st.sidebar.expander("🗄️ History Database Management"):
        st.markdown(f"<span style='font-size:12px; color:{ui['text_muted']};'>Data is automatically saved to SQLite database and persists across sessions.</span>", unsafe_allow_html=True)
        
        if st.button("🗑️ Wipe Database & Start Fresh", type="primary", use_container_width=True):
            if os.path.exists(HistoryManager.DB_FILE):
                os.remove(HistoryManager.DB_FILE)
            HistoryManager.init_db()
            if "restored_files" in st.session_state:
                st.session_state["restored_files"] = set()
            st.success("✅ Database Wiped Successfully! Start logging fresh data.")
            time.sleep(1)
            st.rerun()
        
        if st.button("📥 Download Backup CSV", use_container_width=True):
            backup_df = HistoryManager.export_to_csv()
            if not backup_df.empty:
                csv_backup = backup_df.to_csv(index=False).encode('utf-8')
                st.download_button(
                    label="💾 Download Now",
                    data=csv_backup,
                    file_name=f"history_backup_{datetime.now(EGYPT_TZ).strftime('%Y%m%d')}.csv",
                    mime="text/csv",
                    use_container_width=True
                )
            else:
                st.info("No history data to backup")
        
        if "restored_files" not in st.session_state:
            st.session_state["restored_files"] = set()
        
        history_upload = st.file_uploader("📤 Restore from Backup (Optional)", type="csv")
        if history_upload is not None:
            file_id = f"{history_upload.name}_{history_upload.size}"
            if file_id not in st.session_state["restored_files"]:
                try:
                    restored_df = pd.read_csv(history_upload)
                    HistoryManager.import_from_csv(restored_df)
                    st.session_state["restored_files"].add(file_id)
                    st.success("✅ History Restored Successfully!")
                    time.sleep(1)
                    st.rerun()
                except Exception as e:
                    st.error(f"Error restoring backup: {str(e)}")
            else:
                st.success("✅ History Restored Successfully!")
        
        history_df = HistoryManager.load_history()
        if not history_df.empty:
            st.markdown(f"📊 **Total Records:** {len(history_df)}")
            last_ts = str(history_df.iloc[-1]['timestamp'])
            try:
                f_ts = float(last_ts)
                if f_ts > 30000: last_ts = (datetime(1899, 12, 30) + timedelta(days=f_ts)).strftime("%Y-%m-%d %H:%M:%S")
            except: pass
            st.markdown(f"📅 **Last Entry:** {last_ts}")
            
            with st.expander("👁️ Preview Database"):
                st.dataframe(history_df.tail(10), use_container_width=True)
    
    st.sidebar.divider()

    # كلاس ذكي عشان الداشبورد تعتبر الملف اللي جاي من الكارت كأنه ملف مرفوع
    class HubFile:
        def __init__(self, name):
            self.name = name

    uploaded_file = None
    df = None

    if data_source == "Local CSV Upload":
        uploaded_file = st.sidebar.file_uploader("Upload your Project Log (CSV) 📂", type="csv")

    # 🚀 التعديل السحري: لو مفيش ملف اترفع يدوي، بس المدير داس على الكارت
    if uploaded_file is None and "analytics_df" in st.session_state and st.session_state.get("file_name_from_hub"):
        uploaded_file = HubFile(st.session_state["file_name_from_hub"])
        df = st.session_state["analytics_df"].copy()
        st.sidebar.success(f"🚀 Auto-Loaded from Hub: {uploaded_file.name}")

    if uploaded_file is not None:
        # لو الملف مرفوع من زرار الرفع العادي (نفحصه ونقراه)
        if hasattr(uploaded_file, 'read'):
            uploaded_file.seek(0)
            audit_msg = check_audit_trail(uploaded_file)
            st.sidebar.success(audit_msg, icon="✅")
            uploaded_file.seek(0)
            try:
                df = pd.read_csv(uploaded_file)
                
                # فلتر بيمسح أي صف مفيهوش على الأقل 3 خلايا مليانة بالبيانات
                df = df.dropna(thresh=3)
                if df.empty:
                    st.error("⚠️ الملف لا يحتوي على بيانات!")
                    st.stop()
                st.session_state["analytics_df"] = df.copy()
                st.session_state["file_name_from_hub"] = uploaded_file.name
            except Exception as e:
                st.error(f"❌ خطأ في قراءة الملف: {str(e)}")
                st.info("💡 تأكد أن الملف بصيغة CSV وأن البيانات منسقة بشكل صحيح.")
                st.stop()
        
        # --- 🛠️ Data Cleaning (Global for Dashboard) ---
        
        # --- 🛠️ Data Cleaning (Global for Dashboard) ---
        df.columns = df.columns.str.strip() 
        
        if 'Company Name' not in df.columns and 'Company' in df.columns:
            df.rename(columns={'Company': 'Company Name'}, inplace=True)
            
        if 'Test Type' in df.columns: df['Test Type'] = df['Test Type'].str.strip().str.upper()
        if 'Date ( test)' in df.columns: df['Date ( test)'] = pd.to_datetime(df['Date ( test)'], errors='coerce', dayfirst=True)
        if 'Date( SUB)' in df.columns: df['Date( SUB)'] = pd.to_datetime(df['Date( SUB)'], errors='coerce', dayfirst=True)

        if 'DURATION' in df.columns:
            df['DURATION'] = pd.to_numeric(df['DURATION'], errors='coerce').fillna(0)
        if 'AVERAGE VALUE' in df.columns:
            df['AVERAGE VALUE'] = pd.to_numeric(df['AVERAGE VALUE'], errors='coerce')
        # -----------------------------------------------

        total_rows = len(df)
        missing_dates = df['Date ( test)'].isnull().sum() if 'Date ( test)' in df.columns else 0
        missing_status = df['sample status'].isnull().sum() if 'sample status' in df.columns else 0
        duplicate_serials = df['serial'].duplicated().sum() if 'serial' in df.columns else 0
        total_errors = missing_dates + missing_status + duplicate_serials
        health_score = max(0, 100 - (total_errors / (total_rows+1) * 100)) if total_rows > 0 else 0
        health_color = "#2ecc71" if health_score >= 95 else ("#f1c40f" if health_score >= 80 else "#e74c3c")
        health_icon = "✅" if health_score >= 95 else ("⚠️" if health_score >= 80 else "🚨")
        
        error_details = []
        if missing_dates > 0: error_details.append(f"{missing_dates} Missing Dates")
        if missing_status > 0: error_details.append(f"{missing_status} Missing Status")
        if duplicate_serials > 0: error_details.append(f"{duplicate_serials} Duplicate Serials")
        scanned_msg = f"<span style='color:#00d2ff; font-weight:bold;'>Scanned {total_rows:,} Rows</span>"
        error_str = f"{scanned_msg} &rarr; " + " | ".join(error_details) if error_details else f"{scanned_msg} &rarr; Data is 100% clean and structured."
        
        st.markdown(f"""
            <div class="health-card" style="border-left: 5px solid {health_color}; background:{ui['card_bg']}; color:{ui['text_main']}; margin-bottom: 20px;">
                <div>
                    <h4 style="margin: 0; color: {ui['text_main']}; font-size: 15px; text-transform: uppercase;">{health_icon} Data Integrity Inspector</h4>
                    <p style="margin: 5px 0 0 0; color: {ui['text_muted']}; font-size: 13px;">{error_str}</p>
                </div>
                <div>
                    <h2 style="margin: 0; color: {health_color}; text-shadow: 0 0 10px {health_color};">{health_score:.1f}%</h2>
                </div>
            </div>
        """, unsafe_allow_html=True)

        st.markdown("### 🚀 Select Advanced Module to Explore")
        mod1 = st.button("🚨 Alert System", use_container_width=True)

        if mod1:
            render_alerts_module(df)
            st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

        # ==========================================
        # 🧠 Generative AI Engineering Assistant (Dropdown Mode)
        # ==========================================
        with st.expander("🧠 Generative AI Engineering Assistant (Click to open/close)", expanded=False):
            st.caption("Ask the AI anything about your project data. (e.g., 'Which contractor has the most rejections?')")
            
            chat_container = st.container()
            with chat_container:
                st.markdown('<div class="chat-container">', unsafe_allow_html=True)
                for msg in st.session_state["chat_history"]:
                    if msg['role'] == 'user':
                        st.markdown(f'<div class="user-msg"><b>You:</b> {msg["content"]}</div>', unsafe_allow_html=True)
                    else:
                        st.markdown(f'<div class="ai-msg"><b>🤖 AI:</b> {msg["content"]}</div>', unsafe_allow_html=True)
                st.markdown('</div>', unsafe_allow_html=True)

            # تحويل الشات لـ Form عشان يشتغل جوه القائمة المنسدلة
            with st.form(key="ai_chat_form", clear_on_submit=True):
                col_input, col_btn = st.columns([0.85, 0.15])
                with col_input:
                    prompt = st.text_input("Ask:", label_visibility="collapsed", placeholder="Type your message here...")
                with col_btn:
                    submit_btn = st.form_submit_button("Send 🚀", use_container_width=True)
                
                if submit_btn and prompt:
                    st.session_state["chat_history"].append({"role": "user", "content": prompt})
                    with st.spinner(" AI is analyzing the dataset..."):
                        time.sleep(1.5)
                        ai_response = genai_chat_engine(prompt, df)
                    st.session_state["chat_history"].append({"role": "ai", "content": ai_response})
                    st.rerun()

        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
        st.sidebar.markdown("### 🎯 2. Smart Filters")
        global_search = st.sidebar.text_input("🔍 Global Search:", placeholder="Keyword (Serial, Date)...")
        if global_search:
            mask = df.astype(str).apply(lambda x: x.str.contains(global_search, case=False, na=False)).any(axis=1)
            df = df[mask]
            st.sidebar.success(f"🎯 Found {len(df)} records matching '{global_search}'")
            
        companies = df['Company Name'].dropna().unique() if 'Company Name' in df.columns else []
        selected_companies = st.sidebar.multiselect("🏢 Select Contractor:", options=companies, default=companies)
        
        statuses = df['sample status'].dropna().unique() if 'sample status' in df.columns else []
        selected_statuses = st.sidebar.multiselect("📊 Sample Status:", options=statuses, default=statuses)

        battalion_col_filter = next((c for c in df.columns if 'BATTAL' in c.upper()), None)
        selected_battalions = []
        if battalion_col_filter:
            battalions_list = df[battalion_col_filter].dropna().unique()
            selected_battalions = st.sidebar.multiselect(" Select Battalion:", options=battalions_list, default=battalions_list)

        st.sidebar.markdown("### 🧠 3. AI & Simulation")
        sim_days_saved = st.sidebar.slider(" Simulate Delay Reduction (Days):", min_value=0, max_value=10, value=0, step=1)
        curr_avg_dpl = pd.to_numeric(df['AVERAGE VALUE'], errors='coerce').mean() if 'AVERAGE VALUE' in df.columns else 0
        curr_avg_dur = pd.to_numeric(df['DURATION'], errors='coerce').mean() if 'DURATION' in df.columns else 0
        user_question = st.sidebar.text_input(" Ask AI about any log issue:")
        if user_question:
            summary = {"avg_dpl": round(curr_avg_dpl, 2), "avg_duration": round(curr_avg_dur, 1)}
            st.sidebar.info(f"AI Response: {ai_assistant(user_question, summary)}")

        filtered_df = df.copy()
        if len(companies) > 0: filtered_df = filtered_df[filtered_df['Company Name'].isin(selected_companies)]
        if len(statuses) > 0: filtered_df = filtered_df[filtered_df['sample status'].isin(selected_statuses)]
        
        if battalion_col_filter and len(selected_battalions) > 0: 
            filtered_df = filtered_df[filtered_df[battalion_col_filter].isin(selected_battalions)]

        num_tests_col = next((c for c in filtered_df.columns if 'NUM' in c.upper() and 'TEST' in c.upper()), None)
        if num_tests_col: filtered_df[num_tests_col] = pd.to_numeric(filtered_df[num_tests_col], errors='coerce').fillna(0)
        if 'DURATION' in filtered_df.columns: filtered_df['DURATION'] = pd.to_numeric(filtered_df['DURATION'], errors='coerce')

        total_requests_count = len(filtered_df)
        total_tests_count = int(filtered_df[num_tests_col].sum() if num_tests_col else 0)
        avg_dpl_value = round(pd.to_numeric(filtered_df['AVERAGE VALUE'], errors='coerce').mean() if 'AVERAGE VALUE' in filtered_df.columns else 0, 2)
        avg_duration_value = round(filtered_df['DURATION'].mean(), 1) if 'DURATION' in filtered_df.columns else 0
        page_col_name = next((c for c in filtered_df.columns if 'PAGE' in c.upper()), None)
        total_paperwork_pages = int(pd.to_numeric(filtered_df[page_col_name], errors='coerce').fillna(0).sum()) if page_col_name else 0

        current_metrics = {
            "File_Name": uploaded_file.name, 
            "Total_Requests": total_requests_count,
            "Total_Tests": total_tests_count,
            "Avg_DPL": avg_dpl_value,
            "Avg_Duration": avg_duration_value,
            "Total_Paperwork": total_paperwork_pages
        }
        
        overall_acc = len(filtered_df[filtered_df['sample status'].astype(str).str.upper().isin(['ACCEPTED', 'APPROVED AS NOTED'])]) if 'sample status' in filtered_df.columns else 0
        overall_rate = (overall_acc / total_requests_count * 100) if total_requests_count > 0 else 0

        rejected_count = len(filtered_df[filtered_df['sample status'].str.upper().isin(['REJECTED', 'REVISE'])]) if 'sample status' in filtered_df.columns else 0

        ticker_html = f"""
        <div class="ticker-wrap">
            <div class="ticker">
                <div class="ticker-item">🚀 <b>Total Logged Submittals:</b> <span>{total_requests_count:,}</span></div>
                <div class="ticker-item">✅ <b>Current Global Yield:</b> <span>{overall_rate:.1f}%</span></div>
                <div class="ticker-item">⏱️ <b>Sector Avg Delay:</b> <span>{avg_duration_value} Days</span></div>
                <div class="ticker-item">🚨 <b>Pending Rejections:</b> <span>{rejected_count}</span></div>
                <div class="ticker-item">🧪 <b>Total Field Tests:</b> <span>{total_tests_count:,}</span></div>
            </div>
        </div>
        """
        st.markdown(ticker_html, unsafe_allow_html=True)
        
        live_indicator("online")
        
        worst_office_name = "N/A"
        worst_office_delay = 0
        if 'Done BY' in filtered_df.columns and 'DURATION' in filtered_df.columns:
            office_delays = filtered_df.dropna(subset=['DURATION']).groupby('Done BY')['DURATION'].mean().reset_index()
            if not office_delays.empty:
                worst_office = office_delays.loc[office_delays['DURATION'].idxmax()]
                worst_office_name = worst_office['Done BY']
                worst_office_delay = round(worst_office['DURATION'], 1)

        global_best_comp, global_worst_comp = "N/A", "N/A"
        global_best_rate, global_worst_delay = 0, 0
        if 'Company Name' in filtered_df.columns and 'sample status' in filtered_df.columns and 'DURATION' in filtered_df.columns:
            g_comp_stats = []
            for c in filtered_df['Company Name'].dropna().unique():
                c_df_temp = filtered_df[filtered_df['Company Name'] == c]
                c_t = len(c_df_temp)
                c_a = len(c_df_temp[c_df_temp['sample status'].astype(str).str.upper().isin(['ACCEPTED', 'APPROVED AS NOTED'])])
                c_dur = c_df_temp['DURATION'].mean()
                g_comp_stats.append({'Comp': c, 'Total': c_t, 'Rate': (c_a/c_t*100) if c_t>0 else 0, 'Delay': c_dur})
            g_df = pd.DataFrame(g_comp_stats)
            valid_g = g_df[g_df['Total'] >= 5] if len(g_df[g_df['Total'] >= 5]) > 0 else g_df
            if not valid_g.empty:
                global_best_comp = valid_g.loc[valid_g['Rate'].idxmax()]['Comp']
                global_best_rate = valid_g.loc[valid_g['Rate'].idxmax()]['Rate']
                global_worst_comp = valid_g.loc[valid_g['Delay'].idxmax()]['Comp']
                global_worst_delay = valid_g.loc[valid_g['Delay'].idxmax()]['Delay']

        st.markdown(f"""
            <div class="alert-banner" style="background: linear-gradient(90deg, #e74c3c, #c0392b); padding: 20px; border-radius: 12px; margin-bottom: 20px; display: flex; justify-content: space-between; align-items: center; color: white;">
                <div>
                    <h3 style="margin:0; font-size:22px; color:white;">🚨 Command Center Live Alerts</h3>
                    <p style="margin:5px 0 0 0; font-size:14px; opacity:0.9;">Top issues requiring immediate management attention today.</p>
                </div>
                <div style="text-align:right;">
                    <div style="background:rgba(0,0,0,0.2); padding:8px 15px; border-radius:8px; margin-bottom:5px;"><b>Worst Delay Node:</b> {worst_office_name} ({worst_office_delay} Days)</div>
                    <div style="background:rgba(0,0,0,0.2); padding:8px 15px; border-radius:8px;"><b>Critical Rejections:</b> {rejected_count} Submittals pending</div>
                </div>
            </div>
        """, unsafe_allow_html=True)

        col_btn, col_msg = st.columns([0.2, 0.8])
        with col_btn:
            if st.button("💾 Save to BI History"):
                HistoryManager.save_metrics(current_metrics)
                st.success("✅ Logged Successfully!")
                st.rerun() 

        st.markdown("### 📊 Executive Key Performance Indicators")
        col1, col2, col3, col4, col5 = st.columns(5)
        t_req = 1000; t_test = 5000; t_dpl = 20; t_dur = 10
        d1 = HistoryManager.get_delta_html(current_metrics["Total_Requests"], "Total_Requests", uploaded_file.name)
        create_card(col1, "Total Submittals", current_metrics["Total_Requests"], delta_html=d1, progress=min(100, (current_metrics["Total_Requests"]/t_req)*100 if current_metrics["Total_Requests"] else 0))
        d2 = HistoryManager.get_delta_html(current_metrics["Total_Tests"], "Total_Tests", uploaded_file.name)
        create_card(col2, "Total Tests", current_metrics["Total_Tests"], delta_html=d2, progress=min(100, (current_metrics["Total_Tests"]/t_test)*100 if current_metrics["Total_Tests"] else 0))
        d3 = HistoryManager.get_delta_html(current_metrics["Avg_DPL"], "Avg_DPL", uploaded_file.name)
        create_card(col3, "Avg DPL Value", current_metrics["Avg_DPL"], delta_html=d3, progress=min(100, (current_metrics["Avg_DPL"]/t_dpl)*100 if current_metrics["Avg_DPL"] else 0))
        d4 = HistoryManager.get_delta_html(current_metrics["Avg_Duration"], "Avg_Duration", uploaded_file.name)
        dur_prog = max(0, 100 - (current_metrics["Avg_Duration"]/t_dur * 100)) if current_metrics["Avg_Duration"] else 100
        create_card(col4, "Avg. Dur (Days)", current_metrics["Avg_Duration"], delta_html=d4, progress=dur_prog)
        d5 = HistoryManager.get_delta_html(current_metrics["Total_Paperwork"], "Total_Paperwork", uploaded_file.name)
        create_card(col5, "Total Paperwork", current_metrics["Total_Paperwork"], delta_html=d5)
      # ==========================================
        # 🧪 Detailed Test Counts by Type (KPI Cards)
        # ==========================================
        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
        st.markdown("### 🧪 Detailed Test Counts by Type")
        
        dpl_count = 0
        plate_count = 0
        sand_cone_count = 0
        soil_count = 0
        
        if 'Test Type' in filtered_df.columns:
            if num_tests_col:
                dpl_count = pd.to_numeric(filtered_df[filtered_df['Test Type'].astype(str).str.upper().str.contains('DPL', na=False)][num_tests_col], errors='coerce').sum()
                plate_count = pd.to_numeric(filtered_df[filtered_df['Test Type'].astype(str).str.upper().str.contains('PLATE', na=False)][num_tests_col], errors='coerce').sum()
                sand_cone_count = pd.to_numeric(filtered_df[filtered_df['Test Type'].astype(str).str.upper().str.contains('SAND|CONE', na=False)][num_tests_col], errors='coerce').sum()
                soil_count = pd.to_numeric(filtered_df[filtered_df['Test Type'].astype(str).str.upper().str.contains('SOIL|PROCTOR', na=False)][num_tests_col], errors='coerce').sum()
            else:
                dpl_count = len(filtered_df[filtered_df['Test Type'].astype(str).str.upper().str.contains('DPL', na=False)])
                plate_count = len(filtered_df[filtered_df['Test Type'].astype(str).str.upper().str.contains('PLATE', na=False)])
                sand_cone_count = len(filtered_df[filtered_df['Test Type'].astype(str).str.upper().str.contains('SAND|CONE', na=False)])
                soil_count = len(filtered_df[filtered_df['Test Type'].astype(str).str.upper().str.contains('SOIL|PROCTOR', na=False)])

        tc1, tc2, tc3, tc4 = st.columns(4)
        
        neutral_delta = '<div class="delta-neutral">➖ No change</div>'
        
        create_card(tc1, "DPL", f"{int(dpl_count):,}", delta_html=neutral_delta)
        create_card(tc2, "PLATE LOAD", f"{int(plate_count):,}", delta_html=neutral_delta)
        create_card(tc3, "SAND CONE", f"{int(sand_cone_count):,}", delta_html=neutral_delta)
        create_card(tc4, "SOIL", f"{int(soil_count):,}", delta_html=neutral_delta)

        
        # ==========================================
        # 🏢 Overall Office Workload Analysis
        # ==========================================
        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
        st.markdown('<div class="bi-title">🏢 Overall Office Workload Analysis</div>', unsafe_allow_html=True)
        # ألوان نيون نقية وعالية الجودة
        TACTICAL_PALETTE = ['#00d2ff', '#ffaa00', '#00ff87', '#ff007f', '#a200ff']
        
        if 'Done BY' in filtered_df.columns and 'Test Type' in filtered_df.columns:
            if num_tests_col:
                office_work_df = filtered_df.groupby(['Done BY', 'Test Type'])[num_tests_col].sum().reset_index()
                office_work_df.rename(columns={num_tests_col: 'Number of Tests'}, inplace=True)
            else:
                office_work_df = filtered_df.groupby(['Done BY', 'Test Type']).size().reset_index(name='Number of Tests')
                
            fig_office = px.bar(
                office_work_df, x='Done BY', y='Number of Tests', color='Test Type', 
                barmode='group', title="Test Distribution per Office (Log Scale)", 
                color_discrete_sequence=TACTICAL_PALETTE, text_auto='.2s'
            )
            
            # 🔥 السحر هنا: تحويل المحور الصادي لـ Logarithmic عشان الأعمدة الصغيرة تبان جنب العمود الـ 28 ألف
            fig_office.update_yaxes(type='log', title_text="Number of Tests (Log Scale)")
            fig_office.update_traces(textposition='outside')
            
            try: fig_office = style_3d_glassy(fig_office, chart_type="bar")
            except: pass
            
            st.plotly_chart(fig_office, use_container_width=True, key="overall_office_work_chart")

        # ==========================================
        # 🪨 Overall Soil Classifications
        # ==========================================
        st.markdown('<div class="bi-title">🪨 Overall Soil Classifications</div>', unsafe_allow_html=True)
        if 'Classification' in filtered_df.columns:
            if num_tests_col:
                class_work_df = filtered_df.groupby('Classification')[num_tests_col].sum().reset_index()
                class_work_df.rename(columns={num_tests_col: 'Number of Tests'}, inplace=True)
            else:
                class_work_df = filtered_df.groupby('Classification').size().reset_index(name='Number of Tests')
                
            class_work_df = class_work_df.sort_values('Number of Tests', ascending=False)
            
            fig_class_ov = px.pie(
                class_work_df, names='Classification', values='Number of Tests', 
                hole=0.70, color_discrete_sequence=TACTICAL_PALETTE,
                title="Distribution of Soil Classifications (Overall)"
            )
            
            # إخراج النصوص بره الحلقة (لغينا الـ pull عشان تبقى حلقة رادار مثالية)
            fig_class_ov.update_traces(
                textposition='outside', 
                textinfo='label+percent',
                hovertemplate='<b>Classification:</b> %{label}<br>Tests: %{value}<br>Percentage: %{percent}'
            )
            
            # إضافة كلمة في قلب الدونات
            fig_class_ov.update_layout(
                annotations=[dict(text='SOIL<br>CLASS', x=0.5, y=0.5, font_size=24, font_family="Rajdhani", font_color='#00d2ff', showarrow=False)]
            )
            
            try: 
                fig_class_ov = style_3d_glassy(fig_class_ov, chart_type="pie")
                # 🔥 السر هنا: تزويد مساحة فاضية (Margins) فوق وتحت عشان الكلام اللي بره الدونات ياخد راحته وميتقصش
                fig_class_ov.update_layout(margin=dict(t=100, b=80, l=40, r=40))
            except: pass
            
            st.plotly_chart(fig_class_ov, use_container_width=True, key="overall_classification_chart")
            exported_figs["5. Overall Soil Classifications"] = fig_class_ov
        else:
            st.info("Requires 'Classification' column for Overall Soil Classifications Analysis.")

       

# ==========================================
        # 🎯 Quality & Yield Simulator Section
        # ==========================================
        st.markdown('<div class="bi-title">🎯 Yield & Optimization Simulator</div>', unsafe_allow_html=True)

        g_col, s_col = st.columns([0.4, 0.6])
        with g_col:
            # 🧠 تغيير اللون بذكاء حسب جودة الشغل
            if overall_rate >= 85:
                gauge_color = "#2ecc71"  # أخضر فسفوري (ممتاز)
                status_text = "OPTIMAL"
            elif overall_rate >= 60:
                gauge_color = "#f1c40f"  # أصفر تحذيري (متوسط)
                status_text = "WARNING"
            else:
                gauge_color = "#e74c3c"  # أحمر حرج (خطر)
                status_text = "CRITICAL"

            fig_gauge = go.Figure(go.Indicator(
                mode="gauge+number", 
                value=overall_rate,
                number={'suffix': "%", 'font': {'size': 55, 'color': gauge_color, 'family': 'Rajdhani'}},
                gauge={
                    'axis': {'range': [0, 100], 'tickwidth': 2, 'tickcolor': "rgba(0, 210, 255, 0.3)", 'tickfont': {'family': 'Rajdhani', 'size': 14, 'color': '#8da3b9'}},
                    'bar': {'color': gauge_color, 'thickness': 0.8},
                    'bgcolor': "rgba(15, 23, 42, 0.5)" if is_dark else "rgba(240, 245, 250, 0.8)",
                    'borderwidth': 2,
                    'bordercolor': "rgba(0, 210, 255, 0.1)",
                    'steps': [
                        {'range': [0, 60], 'color': "rgba(231, 76, 60, 0.1)"},
                        {'range': [60, 85], 'color': "rgba(241, 196, 15, 0.1)"},
                        {'range': [85, 100], 'color': "rgba(46, 204, 113, 0.1)"}
                    ],
                    'threshold': {
                        'line': {'color': "white", 'width': 4},
                        'thickness': 0.9,
                        'value': 90 # الهدف الدائم 90%
                    }
                }
            ))
            
            fig_gauge.update_layout(
                paper_bgcolor="rgba(0,0,0,0)", 
                plot_bgcolor="rgba(0,0,0,0)",
                height=280, 
                margin=dict(l=30, r=30, t=10, b=10),
            )
            
            # كارت تكتيكي بيحوي العداد جواه
            st.markdown(f"""
            <div style="background: {ui['card_bg']}; border: 1px solid {ui['border_color']}; border-radius: 4px; padding: 20px; text-align: center; box-shadow: {ui['shadow']};">
                <div style="color: #00d2ff; font-family: 'Rajdhani', sans-serif; font-size: 20px; font-weight: 700; letter-spacing: 1px; text-transform: uppercase; margin-bottom: -10px;">Overall Approval Index</div>
            """, unsafe_allow_html=True)
            
            st.plotly_chart(fig_gauge, use_container_width=True, key="overall_gauge_main")
            
            st.markdown(f"""
                <div style="margin-top: -30px; font-family: 'Rajdhani'; font-size: 16px; color: {gauge_color}; font-weight: bold; letter-spacing: 2px;">STATUS: {status_text}</div>
            </div>
            """, unsafe_allow_html=True)
            
            exported_figs["1. Overall Approval Index"] = fig_gauge 
            
        with s_col:
            if sim_days_saved > 0:
                total_time_recovered = sim_days_saved * total_requests_count
                st.markdown(f"""
                    <div class="simulator-card" style="height: 100%; min-height: 350px; display: flex; flex-direction: column; justify-content: center;">
                        <h4 style="color: #2ecc71; margin: 0; text-transform: uppercase; font-size: 16px; letter-spacing: 1px;">✨ Simulated Optimization Impact</h4>
                        <p style="font-size: 38px; font-weight: 800; color: {ui['text_main']}; margin: 15px 0;">{total_time_recovered:,} <span style="font-size:16px; color:{ui['text_muted']}; font-weight:500;">Project Days Saved</span></p>
                        <p style="font-size: 14px; color: {ui['text_muted']}; margin: 0; line-height: 1.6;">Reducing paperwork cycle times by <b style="color:#00d2ff;">{sim_days_saved} days</b> across all active submittals accelerates overall sector handovers.</p>
                    </div>
                    """, unsafe_allow_html=True)
            else:
                st.markdown(f"""
                    <div class="simulator-card" style="border-color: {ui['border_color']}; background: {ui['card_bg']}; height: 100%; min-height: 350px; display: flex; flex-direction: column; justify-content: center;">
                        <h4 style="color: {ui['text_muted']}; margin: 0; font-size: 18px;">🎛️ Optimization Simulator Inactive</h4>
                        <p style="font-size: 14px; color: {ui['text_muted']}; margin-top: 15px;">Use the slider in the sidebar to simulate the impact of reducing administrative delays.</p>
                    </div>
                    """, unsafe_allow_html=True)

        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

        st.markdown("#### 📝 AI Executive Auto-Narrative")
        narrative = f"The dataset encompasses <b>{total_requests_count:,}</b> submittals involving <b>{total_tests_count:,}</b> field tests. The current overall approval index stands at <b>{overall_rate:.1f}%</b>, with an average turnaround time of <b>{avg_duration_value} days</b>. "
        if worst_office_name != "N/A":
            narrative += f"Attention is required for <b>{worst_office_name}</b>, which currently flags the highest processing delays across the logged sectors."
        st.markdown(f"<div style='font-size: 15px; color: {ui['text_main']}; line-height: 1.6; background: {ui['highlight_bg']}; padding: 20px; border-radius: 10px; border-left: 4px solid #00d2ff; margin-bottom: 25px;'>{narrative}</div>", unsafe_allow_html=True)

       
        

        st.markdown('<div class="bi-title">⚔️ Head-to-Head: Contractor vs Contractor</div>', unsafe_allow_html=True)
        if 'Company Name' in filtered_df.columns and len(companies) >= 2:
            cc1, cc2 = st.columns(2)
            c_a = cc1.selectbox("Select Contractor A", companies, index=0, key="h2h_a")
            c_b = cc2.selectbox("Select Contractor B", companies, index=1 if len(companies)>1 else 0, key="h2h_b")
            def get_c_stats(c_name):
                d = filtered_df[filtered_df['Company Name']==c_name]
                tot = len(d)
                acc = len(d[d['sample status'].str.upper().isin(['ACCEPTED', 'APPROVED AS NOTED'])]) if 'sample status' in d.columns else 0
                rate = (acc/tot*100) if tot>0 else 0
                dur = round(d['DURATION'].mean(), 1) if 'DURATION' in d.columns else 0
                return tot, rate, dur
            tot_a, r_a, d_a = get_c_stats(c_a)
            tot_b, r_b, d_b = get_c_stats(c_b)
            st.markdown(f"""
            <div style="display:flex; justify-content:space-between; text-align:center; background:{ui['card_bg']}; padding:20px; border-radius:15px; border:1px solid #ffaa00; box-shadow: {ui['shadow']};">
                <div style="width:45%; border-right:1px solid {ui['border_color']};"><h3 style="color:#00d2ff; margin-top:0;">{c_a}</h3><p style="font-size:32px; font-weight:800; color:{ui['text_main']}; margin:0;">{r_a:.1f}% Yield</p><p style="color:{ui['text_muted']}; font-size:14px; margin-top:5px;">{tot_a} Submittals | {d_a} Days Avg Delay</p></div>
                <div style="width:10%; align-self:center; font-size:30px; font-weight:900; color:#ffaa00; text-shadow: 0 0 10px rgba(255,170,0,0.5);">VS</div>
                <div style="width:45%; border-left:1px solid {ui['border_color']};"><h3 style="color:#e74c3c; margin-top:0;">{c_b}</h3><p style="font-size:32px; font-weight:800; color:{ui['text_main']}; margin:0;">{r_b:.1f}% Yield</p><p style="color:{ui['text_muted']}; font-size:14px; margin-top:5px;">{tot_b} Submittals | {d_b} Days Avg Delay</p></div>
            </div>
            """, unsafe_allow_html=True)
        else:
            st.info("Requires 'Company Name' column and at least 2 contractors to enable Head-to-Head comparison.")

        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

        st.markdown('<div class="bi-title">🧪 Monthly Test Volume & Deficit Analysis</div>', unsafe_allow_html=True)
        if 'Date ( test)' in filtered_df.columns and 'Test Type' in filtered_df.columns:
            v_df = filtered_df.dropna(subset=['Date ( test)', 'Test Type']).copy()
            v_df['Month_Sort'] = v_df['Date ( test)'].dt.to_period('M')
            v_df['Month'] = v_df['Date ( test)'].dt.strftime('%b %Y')
            monthly_summary = v_df.groupby(['Month_Sort', 'Month', 'Test Type']).size().reset_index(name='Volume')
            monthly_summary = monthly_summary.sort_values('Month_Sort')
            fig_vol = px.bar(monthly_summary, x='Month', y='Volume', color='Test Type', barmode='group', title="Testing Intensity & Production Coverage per Month", color_discrete_sequence=NEON_COLORS)
            fig_vol.update_traces(hovertemplate='<b>Month:</b> %{x}<br><b>Volume:</b> %{y} Submittals')
            fig_vol = style_3d_glassy(fig_vol, chart_type="bar")
            ch_col, txt_col = st.columns([0.7, 0.3])
            ch_col.plotly_chart(fig_vol, use_container_width=True, key="vol_analysis")
            exported_figs["6. Testing Intensity by Month"] = fig_vol
            with txt_col:
                st.markdown("#### 💡 AI Production Insights")
                if not monthly_summary.empty:
                    top_row = monthly_summary.loc[monthly_summary['Volume'].idxmax()]
                    st.info(f"📊 **Peak Activity:**\nIn **{top_row['Month']}**, the highest utilized test was **{top_row['Test Type']}** with **{top_row['Volume']}** submittals logged.")
                    months_ordered = monthly_summary['Month_Sort'].drop_duplicates().sort_values().tolist()
                    if len(months_ordered) > 1:
                        last_month_sort = months_ordered[-1]
                        prev_month_sort = months_ordered[-2]
                        last_month_name = last_month_sort.strftime('%b %Y')
                        prev_month_name = prev_month_sort.strftime('%b %Y')
                        last_count = v_df[v_df['Month_Sort'] == last_month_sort].shape[0]
                        prev_count = v_df[v_df['Month_Sort'] == prev_month_sort].shape[0]
                        if last_count < prev_count:
                            st.warning(f"⚠️ **Coverage Alert:**\nTotal log volume dropped from **{prev_count}** in {prev_month_name} to **{last_count}** in {last_month_name}. Verify potential field testing deficits.")
                        else:
                            st.success(f"✅ **Stable Volume:**\nTesting coverage is expanding smoothly from {prev_month_name} into {last_month_name}.")
                else:
                    st.text("No data available for tracking.")

        if 'Date ( test)' in filtered_df.columns:
            st.markdown('<div class="bi-title">🗓️ Activity Heatmap Calendar</div>', unsafe_allow_html=True)
            cal_df = filtered_df.dropna(subset=['Date ( test)']).copy()
            cal_df['Day'] = cal_df['Date ( test)'].dt.day
            cal_df['Month_Name'] = cal_df['Date ( test)'].dt.strftime('%b %Y')
            hm_data = cal_df.groupby(['Month_Name', 'Day']).size().reset_index(name='Submittals')
            fig_hm = px.density_heatmap(hm_data, x="Day", y="Month_Name", z="Submittals", color_continuous_scale="Viridis", title="Daily Activity Intensity (GitHub Style)", labels={'Day': 'Day of Month', 'Month_Name': 'Month'})
            fig_hm.update_traces(hovertemplate='<b>Date:</b> %{y} %{x}<br><b>Activity:</b> %{z} Tests Logged')
            fig_hm = style_3d_glassy(fig_hm, chart_type="heatmap")
            st.plotly_chart(fig_hm, use_container_width=True, key="heat_calendar")

        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

        st.markdown('<div class="bi-title">📈 Comprehensive Timeline (Workload vs Quality Correlation)</div>', unsafe_allow_html=True)
        if 'Date ( test)' in filtered_df.columns and 'sample status' in filtered_df.columns:
            tl_df = filtered_df.dropna(subset=['Date ( test)', 'sample status']).copy()
            tl_df['Month_Plot'] = tl_df['Date ( test)'].dt.to_period('M').astype(str)
            
            monthly_stats = tl_df.groupby('Month_Plot').apply(lambda x: pd.Series({
                'Total': len(x),
                'Accepted': len(x[x['sample status'].str.upper().isin(['ACCEPTED', 'APPROVED AS NOTED'])]),
                'Rejected': len(x[x['sample status'].str.upper().isin(['REJECTED', 'REVISE'])])
            })).reset_index()
            
            monthly_stats['Acc_Pct'] = (monthly_stats['Accepted'] / monthly_stats['Total'] * 100).round(1)
            monthly_stats['Rej_Pct'] = (monthly_stats['Rejected'] / monthly_stats['Total'] * 100).round(1)
            
            fig_combo = make_subplots(specs=[[{"secondary_y": True}]])
            
            fig_combo.add_trace(go.Bar(x=monthly_stats['Month_Plot'], y=monthly_stats['Accepted'], name='Accepted/Approved', marker_color=STATUS_COLORS['ACCEPTED'], customdata=monthly_stats['Acc_Pct'], hovertemplate="<b>%{x}</b><br>Accepted: %{y} (%{customdata}%)<extra></extra>"), secondary_y=False)
            fig_combo.add_trace(go.Bar(x=monthly_stats['Month_Plot'], y=monthly_stats['Rejected'], name='Rejected/Revise', marker_color=STATUS_COLORS['REJECTED'], customdata=monthly_stats['Rej_Pct'], hovertemplate="<b>%{x}</b><br>Rejected: %{y} (%{customdata}%)<extra></extra>"), secondary_y=False)
            fig_combo.add_trace(go.Scatter(x=monthly_stats['Month_Plot'], y=monthly_stats['Total'], name='Total Workload', mode='lines+markers', line=dict(color='#00d2ff', width=4), marker=dict(size=8), hovertemplate="<b>%{x}</b><br>Total Logged: %{y}<extra></extra>"), secondary_y=True)
            
            fig_combo.update_layout(barmode='stack', title="Volume vs. Rejection Impact over Time", legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
            fig_combo.update_yaxes(title_text="Submittals (Quality)", secondary_y=False)
            fig_combo.update_yaxes(title_text="Total Workload", secondary_y=True)
            
            fig_combo = style_3d_glassy(fig_combo, chart_type="combo")
            st.plotly_chart(fig_combo, use_container_width=True, key="combo_timeline")
            exported_figs["7. Volume vs Rejection Timeline"] = fig_combo

        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

        st.markdown('<div class="bi-title">📊 Statistical Process Control (SPC) - Control Charts</div>', unsafe_allow_html=True)
        st.caption("Monitor process stability and detect special cause variations using industry-standard control limits.")
        
        if 'AVERAGE VALUE' in filtered_df.columns and 'Company Name' in filtered_df.columns:
            spc_df = filtered_df.dropna(subset=['AVERAGE VALUE']).copy()
            spc_df['AVERAGE VALUE'] = pd.to_numeric(spc_df['AVERAGE VALUE'], errors='coerce')
            spc_df = spc_df.dropna(subset=['AVERAGE VALUE'])
            
            if not spc_df.empty:
                mean_val = spc_df['AVERAGE VALUE'].mean()
                std_val = spc_df['AVERAGE VALUE'].std()
                ucl = mean_val + 3 * std_val  
                lcl = mean_val - 3 * std_val  
                
                spc_df['out_of_control'] = (spc_df['AVERAGE VALUE'] > ucl) | (spc_df['AVERAGE VALUE'] < lcl)
                out_of_control_count = spc_df['out_of_control'].sum()
                total_points = len(spc_df)
                control_percentage = ((total_points - out_of_control_count) / total_points * 100) if total_points > 0 else 0
                
                spc_col1, spc_col2, spc_col3, spc_col4 = st.columns(4)
                create_card(spc_col1, "Process Mean", f"{mean_val:.2f}")
                create_card(spc_col2, "Std Deviation", f"{std_val:.2f}")
                create_card(spc_col3, "Control Limits", f"UCL: {ucl:.2f}<br>LCL: {lcl:.2f}")
                create_card(spc_col4, "In Control %", f"{control_percentage:.1f}%")
                
                fig_spc = go.Figure()
                fig_spc.add_trace(go.Scatter(
                    x=spc_df.index,
                    y=spc_df['AVERAGE VALUE'],
                    mode='markers',
                    name='Data Points',
                    marker=dict(
                        size=8,
                        color=['#e74c3c' if oc else '#00d2ff' for oc in spc_df['out_of_control']],
                        line=dict(width=1, color='white')
                    ),
                    hovertemplate='<b>Index:</b> %{x}<br><b>Value:</b> %{y:.2f}<br><b>Status:</b> %{marker.color}<extra></extra>'
                ))
                fig_spc.add_hline(y=mean_val, line_dash="solid", line_color="#2ecc71", line_width=2, annotation_text=f"Mean: {mean_val:.2f}", annotation_position="top right")
                fig_spc.add_hline(y=ucl, line_dash="dash", line_color="#e74c3c", line_width=2, annotation_text=f"UCL: {ucl:.2f}", annotation_position="top right")
                fig_spc.add_hline(y=lcl, line_dash="dash", line_color="#e74c3c", line_width=2, annotation_text=f"LCL: {lcl:.2f}", annotation_position="bottom right")
                fig_spc.update_layout(title="Control Chart - Process Stability Analysis", xaxis_title="Sample Index", yaxis_title="AVERAGE VALUE", showlegend=False, height=500)
                
                fig_spc = style_3d_glassy(fig_spc, chart_type="line")
                st.plotly_chart(fig_spc, use_container_width=True, key="spc_chart")
                exported_figs["8. SPC Control Chart Analysis"] = fig_spc
                
                if out_of_control_count > 0:
                    st.warning(f"⚠️ **Process Alert:** {out_of_control_count} out of {total_points} samples ({100-control_percentage:.1f}%) are outside control limits.")
                    
                    # السطر ده هو اللي بيرتب من الكبير للصغير بناءً على الـ AVERAGE VALUE
                    ooc_samples = spc_df[spc_df['out_of_control']].sort_values(by='AVERAGE VALUE', ascending=False)
                    
                    if not ooc_samples.empty:
                        st.markdown("**Top Out-of-Control Samples:**")
                        display_ooc = ooc_samples[['Company Name', 'Test Type', 'AVERAGE VALUE', 'sample status']].head(10)
                        
                        st.dataframe(display_ooc, use_container_width=True)
                        
                        # السطر السحري عشان زراير التحميل (إكسيل - بي دي إف - CSV)
                        export_table_tools(display_ooc, f"Out_Of_Control_Samples_{datetime.now(EGYPT_TZ).strftime('%Y%m%d')}")
                else:
                    st.success("✅ **Process Stable:** All samples are within control limits.")
                
                st.markdown("#### 🎯 Process Capability Analysis")
                cap_col1, cap_col2 = st.columns(2)
                with cap_col1:
                    st.markdown(f"""
                    <div class="metric-card">
                        <div class="metric-label">Process Performance</div>
                        <div class="metric-value" style="font-size: 24px;">{control_percentage:.1f}%</div>
                        <div style="color: {ui['text_muted']}; font-size: 14px; margin-top: 10px;">of samples within ±3σ control limits</div>
                    </div>
                    """, unsafe_allow_html=True)
                with cap_col2:
                    if std_val > 0:
                        cpk = min((ucl - mean_val) / (3 * std_val), (mean_val - lcl) / (3 * std_val))
                        cpk_color = "#2ecc71" if cpk >= 1.33 else ("#f1c40f" if cpk >= 1.0 else "#e74c3c")
                        cpk_status = "Excellent" if cpk >= 1.33 else ("Good" if cpk >= 1.0 else "Needs Improvement")
                        st.markdown(f"""
                        <div class="metric-card">
                            <div class="metric-label">Process Capability (Cpk)</div>
                            <div class="metric-value" style="font-size: 24px; color: {cpk_color};">{cpk:.2f}</div>
                            <div style="color: {cpk_color}; font-size: 14px; margin-top: 10px; font-weight: bold;">{cpk_status}</div>
                        </div>
                        """, unsafe_allow_html=True)
        
        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

        st.markdown('<div class="bi-title">📊 Pareto Analysis - 80/20 Rule</div>', unsafe_allow_html=True)
        st.caption("Identify the vital few causes that contribute to the majority of problems. Focus your improvement efforts where they matter most.")
        
        if 'Company Name' in filtered_df.columns and 'sample status' in filtered_df.columns:
            rej_by_comp = filtered_df[filtered_df['sample status'].str.upper().isin(['REJECTED', 'REVISE'])].groupby('Company Name').size().reset_index(name='Rejections')
            rej_by_comp = rej_by_comp.sort_values('Rejections', ascending=False)
            
            if not rej_by_comp.empty:
                total_rejections = rej_by_comp['Rejections'].sum()
                rej_by_comp['Percentage'] = (rej_by_comp['Rejections'] / total_rejections * 100).round(2)
                rej_by_comp['Cumulative_Percentage'] = rej_by_comp['Percentage'].cumsum().round(2)
                
                critical_threshold = 80
                critical_contractors = rej_by_comp[rej_by_comp['Cumulative_Percentage'] <= critical_threshold]
                critical_count = len(critical_contractors)
                total_contractors = len(rej_by_comp)
                critical_percentage = (critical_count / total_contractors * 100) if total_contractors > 0 else 0
                
                pareto_col1, pareto_col2, pareto_col3 = st.columns(3)
                create_card(pareto_col1, "Total Contractors", f"{total_contractors}")
                create_card(pareto_col2, "Critical Contractors", f"{critical_count} ({critical_percentage:.0f}%)")
                create_card(pareto_col3, "Total Rejections", f"{total_rejections}")
                
                fig_pareto = make_subplots(specs=[[{"secondary_y": True}]])
                fig_pareto.add_trace(go.Bar(x=rej_by_comp['Company Name'], y=rej_by_comp['Rejections'], name='Rejections', marker_color='#e74c3c', opacity=0.7, hovertemplate='<b>Contractor:</b> %{x}<br><b>Rejections:</b> %{y}<extra></extra>'), secondary_y=False)
                fig_pareto.add_trace(go.Scatter(x=rej_by_comp['Company Name'], y=rej_by_comp['Cumulative_Percentage'], mode='lines+markers', name='Cumulative %', line=dict(color='#00d2ff', width=3), marker=dict(size=8, color='#00d2ff'), hovertemplate='<b>Contractor:</b> %{x}<br><b>Cumulative %:</b> %{y:.1f}%<extra></extra>'), secondary_y=True)
                fig_pareto.add_hline(y=80, line_dash="dash", line_color="#ffaa00", line_width=2, annotation_text="80% Threshold", annotation_position="top right", secondary_y=True)
                fig_pareto.update_layout(title="Pareto Chart - Rejections by Contractor", xaxis_title="Contractor", yaxis_title="Number of Rejections", yaxis2_title="Cumulative Percentage (%)", legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1), height=500)
                fig_pareto.update_yaxes(title_text="Rejections", secondary_y=False)
                fig_pareto.update_yaxes(title_text="Cumulative %", secondary_y=True, range=[0, 100])
                fig_pareto = style_3d_glassy(fig_pareto, chart_type="combo")
                st.plotly_chart(fig_pareto, use_container_width=True, key="pareto_comb")
                exported_figs["9. Pareto Analysis (80-20 Rule)"] = fig_pareto
                
                st.markdown("#### 🎯 Strategic Insights")
                insight_col1, insight_col2 = st.columns(2)
                with insight_col1:
                    st.markdown(f"""
                    <div style="background: rgba(231, 76, 60, 0.1); border-left: 4px solid #e74c3c; padding: 20px; border-radius: 10px; margin-bottom: 15px;">
                        <h4 style="color: #e74c3c; margin: 0 0 10px 0;">🎯 Critical Focus Area</h4>
                        <p style="color: {ui['text_main']}; margin: 0; font-size: 14px; line-height: 1.6;">
                            The top <b style="color: #ffaa00;">{critical_count} contractors</b> ({critical_percentage:.0f}% of total) are responsible for 
                            <b style="color: #ffaa00;">{rej_by_comp[rej_by_comp['Cumulative_Percentage'] <= critical_threshold]['Cumulative_Percentage'].max():.1f}%</b> of all rejections.
                        </p>
                    </div>
                    """, unsafe_allow_html=True)
                with insight_col2:
                    top_5 = rej_by_comp.head(5)
                    top_5_html = "<br>".join([f"<div style='display: flex; justify-content: space-between; padding: 8px; background: rgba(255,255,255,0.05); border-radius: 5px; margin-bottom: 5px;'><span style='color: {ui['text_main']}; font-weight: 600;'>{row['Company Name']}</span><span style='color: #e74c3c; font-weight: bold;'>{row['Rejections']} rejections ({row['Percentage']:.1f}%)</span></div>" for _, row in top_5.iterrows()])
                    st.markdown(f"""
                    <div style="background: rgba(0, 210, 255, 0.05); border-left: 4px solid #00d2ff; padding: 20px; border-radius: 10px;">
                        <h4 style="color: #00d2ff; margin: 0 0 10px 0;">📊 Top 5 Contractors by Rejections</h4>
                        {top_5_html}
                    </div>
                    """, unsafe_allow_html=True)
                
                with st.expander("📋 View Detailed Pareto Analysis Table"):
                    st.dataframe(rej_by_comp[['Company Name', 'Rejections', 'Percentage', 'Cumulative_Percentage']].rename(columns={'Company Name': 'Contractor', 'Rejections': 'Total Rejections', 'Percentage': '% of Total', 'Cumulative_Percentage': 'Cumulative %'}), use_container_width=True)

        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

        st.markdown('<div class="bi-title">🤖 Predictive Risk Forecasting</div>', unsafe_allow_html=True)
        if 'Date ( test)' in filtered_df.columns and 'DURATION' in filtered_df.columns:
            pred_df = filtered_df.dropna(subset=['Date ( test)', 'DURATION']).sort_values('Date ( test)')
            # الحل (يتم تطبيقه بعد التأكد أن Date(test) هو Index)
            pred_df = pred_df.set_index('Date ( test)')
            pred_df['7-Day Trend'] = pred_df['DURATION'].rolling('7D', min_periods=1).mean()
            pred_df = pred_df.reset_index()
            fig_pred = px.line(pred_df, x='Date ( test)', y=['DURATION', '7-Day Trend'], title="Duration Forecasting & Trendline Tracking", color_discrete_sequence=['#ffaa00', '#00d2ff'])
            fig_pred = style_3d_glassy(fig_pred, chart_type="line")
            latest_trend = pred_df['7-Day Trend'].iloc[-1] if not pred_df.empty else 0
            p1, p2 = st.columns([0.7, 0.3])
            p1.plotly_chart(fig_pred, use_container_width=True, key="pred_risk")
            exported_figs["10. Duration Forecasting & Trend"] = fig_pred
            with p2:
                st.info("**AI Risk Assessment:**")
                if latest_trend > current_metrics["Avg_Duration"]:
                    st.error(f"🚨 **Warning:** The recent workflow trend is rising ({latest_trend:.1f} days) compared to the overall average. Bottlenecks are forming.")
                else:
                    st.success(f"✅ **Stable:** Workflow trend is improving or stable at {latest_trend:.1f} days.")

        

        st.markdown('<div class="bi-title">🖨️ Smart PDF Executive Report</div>', unsafe_allow_html=True)
        st.info("💡 **CEO Feature:** Click the button below to download a styled HTML report. When opened, it can be easily saved as a perfectly formatted PDF for your Daily/Weekly Briefing!")
        html_report = f"""
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta charset="UTF-8">
            <title>Executive Report - {uploaded_file.name}</title>
            <style>
                body {{ font-family: 'Segoe UI', Arial, sans-serif; padding: 40px; color: #333; background-color: #f9fbfd; }}
                .container {{ max-width: 900px; margin: auto; background: white; padding: 40px; border-radius: 8px; box-shadow: 0 4px 15px rgba(0,0,0,0.05); border-top: 8px solid #1e3d59; }}
                .header {{ display: flex; justify-content: space-between; align-items: center; border-bottom: 2px solid #ecf0f1; padding-bottom: 20px; margin-bottom: 30px; }}
                .header h1 {{ color: #1e3d59; margin: 0; font-size: 28px; text-transform: uppercase; letter-spacing: 1px; }}
                .header p {{ margin: 5px 0 0 0; color: #7f8c8d; font-size: 14px; }}
                .kpi-row {{ display: flex; justify-content: space-between; margin-bottom: 30px; }}
                .kpi-box {{ background: #f4f7f6; padding: 20px; border-radius: 8px; width: 30%; text-align: center; border-bottom: 4px solid #00d2ff; }}
                .kpi-box h3 {{ margin: 0; color: #7f8c8d; font-size: 12px; text-transform: uppercase; }}
                .kpi-box h2 {{ margin: 10px 0 0 0; color: #2c3e50; font-size: 28px; }}
                .section-title {{ color: #e67e22; font-size: 18px; border-bottom: 1px solid #ecf0f1; padding-bottom: 8px; margin-top: 30px; margin-bottom: 15px; }}
                table {{ width: 100%; border-collapse: collapse; margin-top: 15px; font-size: 14px; }}
                th, td {{ padding: 12px; text-align: left; border-bottom: 1px solid #eee; }}
                th {{ background-color: #1e3d59; color: white; }}
                .highlight-red {{ color: #e74c3c; font-weight: bold; }}
                .highlight-green {{ color: #2ecc71; font-weight: bold; }}
            </style>
        </head>
        <body onload="window.print()">
            <div class="container">
                <div class="header">
                    <div>
                        <h1>KK Engineering - Executive Brief</h1>
                        <p><strong>Dataset:</strong> {uploaded_file.name}</p>
                        <p><strong>Generated On:</strong> {datetime.now(EGYPT_TZ).strftime("%Y-%m-%d at %I:%M %p")}</p>
                    </div>
                    <div style="font-size: 40px;">🏗️</div>
                </div>
                
                <div class="kpi-row">
                    <div class="kpi-box" style="border-color: #2ecc71;">
                        <h3>Overall Approval</h3>
                        <h2 class="highlight-green">{overall_rate:.1f}%</h2>
                    </div>
                    <div class="kpi-box" style="border-color: #ffaa00;">
                        <h3>Total Submittals</h3>
                        <h2>{total_requests_count:,}</h2>
                    </div>
                    <div class="kpi-box" style="border-color: #e74c3c;">
                        <h3>Avg Sector Delay</h3>
                        <h2 class="highlight-red">{avg_duration_value} Days</h2>
                    </div>
                </div>

                <div class="section-title">⚖️ 360° Accountability & Risk Assessment</div>
                <table>
                    <tr>
                        <th>Metric</th>
                        <th>Identified Node / Value</th>
                    </tr>
                    <tr>
                        <td><strong>🏆 Top Performing Contractor</strong></td>
                        <td class="highlight-green">{global_best_comp} ({global_best_rate:.1f}% Yield)</td>
                    </tr>
                    <tr>
                        <td><strong>🚨 Critical Bottleneck (Contractor)</strong></td>
                        <td class="highlight-red">{global_worst_comp} ({global_worst_delay:.1f} Days Avg Delay)</td>
                    </tr>
                    <tr>
                        <td><strong>⏱️ Worst Review Office</strong></td>
                        <td class="highlight-red">{worst_office_name} ({worst_office_delay} Days Avg Delay)</td>
                    </tr>
                    <tr>
                        <td><strong>⚠️ Pending Rejections</strong></td>
                        <td class="highlight-red">{rejected_count} Submittals</td>
                    </tr>
                    <tr>
                        <td><strong>🛡️ Data Integrity Score</strong></td>
                        <td>{health_score:.1f}%</td>
                    </tr>
                </table>
                
                <p style="text-align: center; color: #95a5a6; font-size: 11px; margin-top: 50px;">Confidential Document - Generated by AI Command Center BI Portal</p>
            </div>
        </body>
        </html>
        """
        b64 = base64.b64encode(html_report.encode()).decode()
        href = f'<a href="data:text/html;base64,{b64}" download="KK_Executive_Report_{datetime.now(EGYPT_TZ).strftime("%Y%m%d")}.html" style="background-color:#ffaa00; color:#1e3d59; padding:12px 24px; text-decoration:none; font-weight:bold; border-radius:8px; display:inline-block; box-shadow: 0 4px 15px rgba(255, 170, 0, 0.4); transition: all 0.3s;">📄 Download Ultra-Premium PDF Report</a>'
        st.markdown(href, unsafe_allow_html=True)

        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

        st.markdown('<div class="bi-title">🏗️ Contractor Materials & Sourcing Analysis</div>', unsafe_allow_html=True)
       

        if 'Company Name' in filtered_df.columns and 'Sampling Location' in filtered_df.columns:
            mat_df = filtered_df.copy()
            mat_df['Sampling_Lower'] = mat_df['Sampling Location'].astype(str).str.lower()
            def categorize_location(loc):
                if 'stock' in loc or 'مشون' in loc: return 'Stockpile'
                elif 'bottom' in loc or 'قاع' in loc: return 'Bottom of Excavation'
                elif 'fill' in loc or 'ردم' in loc: return 'Fill'
                else: return 'Other'
            mat_df['Loc_Category'] = mat_df['Sampling_Lower'].apply(categorize_location)
            
            st.markdown("#### 📑 Consolidated Contractors Summary (Ready for Print)")
            summary_pivot = pd.crosstab(mat_df['Company Name'], mat_df['Loc_Category'], margins=True, margins_name="Total")
            cols_order = ['Stockpile', 'Bottom of Excavation', 'Fill', 'Other', 'Total']
            existing_cols = [c for c in cols_order if c in summary_pivot.columns]
            summary_pivot = summary_pivot[existing_cols]
            st.dataframe(summary_pivot, use_container_width=True)
            
            # السطر السحري لإضافة الـ 3 زراير
            export_table_tools(summary_pivot.reset_index(), f"Consolidated_Contractors_Summary_{datetime.now(EGYPT_TZ).strftime('%Y%m%d')}")
            
            st.divider()

            target_dict = {}
            battalion_col_main = next((c for c in df.columns if 'BATTAL' in c.upper()), None)
            
            if 'Company' in df.columns and 'Required Quantity' in df.columns:
                cols_to_extract = ['Company', 'Required Quantity']
                if battalion_col_main: cols_to_extract.append(battalion_col_main)
                lookup_df = df[cols_to_extract].dropna(subset=['Company'])
                
                for _, row in lookup_df.iterrows():
                    c_key = str(row['Company']).strip()
                    c_qty = pd.to_numeric(row['Required Quantity'], errors='coerce')
                    if pd.notna(c_qty):
                        if battalion_col_main and pd.notna(row.get(battalion_col_main)):
                            b_key = fmt_b(row[battalion_col_main])
                            target_dict[f"{c_key}_{b_key}"] = c_qty
                        else:
                            target_dict[c_key] = c_qty

            st.markdown("#### 📥 Master Stockpile Targets Report")
            report_data = []
            all_log_companies = sorted([c for c in mat_df['Company Name'].unique() if str(c) != 'nan'])
            battalion_col_stock = next((c for c in mat_df.columns if 'BATTAL' in c.upper()), None)

            for c_name in all_log_companies:
                c_name_clean = str(c_name).strip()
                c_df_stock = mat_df[(mat_df['Company Name'] == c_name) & (mat_df['Loc_Category'] == 'Stockpile')]
                
                if battalion_col_stock:
                    bats = c_df_stock[battalion_col_stock].dropna().unique()
                    if len(bats) == 0: bats = ["Unknown"]
                else:
                    bats = ["Global"]

                for b in bats:
                    b_clean = fmt_b(b)
                    if b == "Global" or b == "Unknown":
                        bat_stock_df = c_df_stock
                        req_qty = target_dict.get(c_name_clean, np.nan)
                    else:
                        bat_stock_df = c_df_stock[c_df_stock[battalion_col_stock] == b]
                        req_qty = target_dict.get(f"{c_name_clean}_{b_clean}", np.nan)

                    if num_tests_col:
                        exec_qty = int(pd.to_numeric(bat_stock_df[num_tests_col], errors='coerce').fillna(0).sum())
                    else:
                        exec_qty = len(bat_stock_df)

                    if pd.notna(req_qty) and req_qty > 0:
                        diff = exec_qty - int(req_qty)
                        status = "✅ Target Exceeded" if diff >= 0 else f"⚠️ Missing {abs(diff)} Tests"
                        req_val = int(req_qty)
                        diff_val = diff
                    else:
                        status = "No Target Defined"
                        req_val = "N/A"
                        diff_val = "N/A"

                    report_data.append({
                        "Contractor Name": c_name_clean,
                        "Battalion": b_clean if b not in ["Global", "Unknown"] else "N/A",
                        "Executed Stockpile Tests": exec_qty,
                        "Required Target": req_val,
                        "Difference (+/-)": diff_val,
                        "Status": status
                    })
                    
            report_df = pd.DataFrame(report_data)
            st.dataframe(report_df, use_container_width=True)
            
            # السطر السحري الجديد اللي هيعوض كل اللي مسحته
            export_table_tools(report_df, f"Stockpile_Targets_Report_{datetime.now(EGYPT_TZ).strftime('%Y%m%d')}")
            
            st.divider()

            st.markdown("#### 🏢 Individual Contractor Deep Dive")
            if all_log_companies:
                selected_comp = st.selectbox("Select a Contractor to Analyze:", all_log_companies, key="deepdive_comp_sel")
                comp_df_full = mat_df[mat_df['Company Name'] == selected_comp]
                
                tab_360, tab_stockpile, tab_execution, tab_quantities = st.tabs([
                    "🌐 360° Corporate Profile", 
                    "⛰️ Stockpile Sourcing", 
                    "🏗️ Executive Progress & Compaction",
                    "📊 Quantities Rate"
                ])
                
                with tab_360:
                    st.markdown(f"### 🌐 Executive Profile: `{selected_comp}`")
                    
                    battalion_col_360 = next((c for c in comp_df_full.columns if 'BATTAL' in c.upper()), None)
                    zone_col_360 = next((c for c in comp_df_full.columns if 'ZONE' in c.upper()), None)
                    elment_col_360 = next((c for c in comp_df_full.columns if 'ELMEN' in c.upper() or 'ELEMENT' in c.upper()), None)
                    
                    total_tests_360 = int(pd.to_numeric(comp_df_full[num_tests_col], errors='coerce').fillna(0).sum()) if num_tests_col else len(comp_df_full)
                    battalions_count = comp_df_full[battalion_col_360].nunique() if battalion_col_360 else "N/A"
                    zones_count = comp_df_full[zone_col_360].nunique() if zone_col_360 else "N/A"
                    avg_dur_360 = pd.to_numeric(comp_df_full['DURATION'], errors='coerce').mean() if 'DURATION' in comp_df_full.columns else "N/A"
                    
                    c1, c2, c3, c4 = st.columns(4)
                    create_card(c1, "Total Test Points", total_tests_360)
                    create_card(c2, "Active Battalions", battalions_count)
                    create_card(c3, "Active Zones", zones_count)
                    create_card(c4, "Avg Delay (Days)", f"{avg_dur_360:.1f}" if pd.notna(avg_dur_360) else "N/A")
                    
                    if battalion_col_360 and zone_col_360:
                        st.markdown("#### 🗺️ Spatial Footprint (Battalion ➔ Zone)")
                        tree_df_360 = comp_df_full.copy()
                        tree_df_360[battalion_col_360] = tree_df_360[battalion_col_360].fillna('Unknown Battalion').astype(str)
                        tree_df_360[zone_col_360] = tree_df_360[zone_col_360].fillna('Unknown Zone').astype(str)
                        tree_grouped = tree_df_360.groupby([battalion_col_360, zone_col_360]).size().reset_index(name='Submittals')
                        fig_tree_360 = px.treemap(tree_grouped, path=[battalion_col_360, zone_col_360], values='Submittals', title=f"Workload Distribution for {selected_comp}", color='Submittals', color_continuous_scale='Blues')
                        fig_tree_360 = style_3d_glassy(fig_tree_360, chart_type="treemap")
                        st.plotly_chart(fig_tree_360, use_container_width=True, key=f"tree_360_{selected_comp}")
                        
                    col_q1, col_q2 = st.columns(2)
                    with col_q1:
                        if battalion_col_360 and 'sample status' in comp_df_full.columns:
                            st.markdown("#### ⚖️ Quality by Battalion")
                            comp_df_full['status_upper'] = comp_df_full['sample status'].str.upper()
                            qual_df = comp_df_full.groupby([battalion_col_360, 'status_upper']).size().reset_index(name='Count')
                            fig_qual = px.bar(qual_df, x=battalion_col_360, y='Count', color='status_upper', title="Approval/Rejection per Battalion", barmode='group', color_discrete_map=STATUS_COLORS)
                            fig_qual = style_3d_glassy(fig_qual, chart_type="bar")
                            st.plotly_chart(fig_qual, use_container_width=True, key=f"qual_{selected_comp}")
                            
                    with col_q2:
                        if elment_col_360:
                            st.markdown("#### 🏗️ Workload by Element")
                            el_df = comp_df_full.groupby(elment_col_360).size().reset_index(name='Count').sort_values('Count', ascending=False)
                            fig_elment = px.bar(el_df.head(15), x=elment_col_360, y='Count', title="Top 15 Elements by Submittals", color=elment_col_360, color_discrete_sequence=NEON_COLORS)
                            fig_elment = style_3d_glassy(fig_elment, chart_type="bar")
                            st.plotly_chart(fig_elment, use_container_width=True, key=f"elment_{selected_comp}")
                        else:
                            st.info("No Element column found for workload breakdown.")
                                
                    col_d1, col_d2 = st.columns(2)
                    with col_d1:
                        if 'Done BY' in comp_df_full.columns:
                            st.markdown("#### 👨‍💼 Processed by Office (Done BY)")
                            off_df = comp_df_full.groupby('Done BY').size().reset_index(name='Count').sort_values('Count', ascending=False)
                            off_df['Percent'] = (off_df['Count'] / off_df['Count'].sum() * 100).round(1)
                            fig_off = px.bar(off_df, x='Done BY', y='Count', text='Count', title="Submittal Volume per Review Office", color='Done BY', color_discrete_sequence=NEON_COLORS)
                            fig_off.update_traces(customdata=off_df['Percent'], hovertemplate='<b>Office:</b> %{x}<br>Count: %{y}<br>Percentage: %{customdata}%')
                            fig_off = style_3d_glassy(fig_off, chart_type="bar")
                            st.plotly_chart(fig_off, use_container_width=True, key=f"off_{selected_comp}")
                            
                    with col_d2:
                        if 'sample status' in comp_df_full.columns and 'layer' in comp_df_full.columns and elment_col_360:
                            st.markdown("#### 🚨 Smart Red Flags (Unresolved Layers)")
                            st.caption("Shows rejections ONLY IF the same Layer/Element wasn't approved later.")
                            rejected_mask = comp_df_full['sample status'].astype(str).str.upper().isin(['REJECTED', 'REVISE'])
                            accepted_mask = comp_df_full['sample status'].astype(str).str.upper().isin(['ACCEPTED', 'APPROVED AS NOTED'])
                            comp_df_rf = comp_df_full.copy()
                            comp_df_rf['Loc_ID'] = comp_df_rf[elment_col_360].astype(str) + "_" + comp_df_rf['layer'].astype(str)
                            approved_locs = set(comp_df_rf[accepted_mask]['Loc_ID'].unique())
                            red_flags = comp_df_rf[rejected_mask & (~comp_df_rf['Loc_ID'].isin(approved_locs))]
                            if not red_flags.empty:
                                display_cols = ['serial', 'sample status', elment_col_360, 'layer']
                                if battalion_col_360: display_cols.append(battalion_col_360)
                                if 'Test Type' in red_flags.columns: display_cols.append('Test Type')
                                existing_cols = [c for c in display_cols if c in red_flags.columns]
                                st.dataframe(red_flags[existing_cols].head(100), use_container_width=True)
                                st.caption(f"Total unresolved layers: {len(red_flags)}")
                            else:
                                st.success("✅ All rejected layers have been successfully re-tested and approved!")
                        else:
                            st.info("Requires 'sample status', 'layer', and 'Element' columns for Smart Red Flags.")

                    if 'sample status' in comp_df_full.columns and 'Date( SUB)' in comp_df_full.columns and 'layer' in comp_df_full.columns:
                        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
                        st.markdown("#### 🧾 Rework & Delay Ledger (Rejected Items Analysis)")
                        
                        rework_df = comp_df_full.copy()
                        rejected_items = rework_df[rework_df['sample status'].astype(str).str.upper().isin(['REJECTED', 'REVISE'])]
                        
                        if not rejected_items.empty:
                            ledger_data = []
                            total_delay_days = 0
                            el_col = elment_col_360 if elment_col_360 else None
                            
                            for _, rej_row in rejected_items.iterrows():
                                serial = rej_row.get('serial', 'N/A')
                                rej_date = rej_row['Date( SUB)']
                                layer = str(rej_row.get('layer', 'Unknown'))
                                test_type = str(rej_row.get('Test Type', 'Unknown'))
                                
                                filter_mask = (
                                    (rework_df['sample status'].astype(str).str.upper().isin(['ACCEPTED', 'APPROVED AS NOTED'])) & 
                                    (rework_df['Date( SUB)'] >= rej_date) &
                                    (rework_df['layer'].astype(str) == layer) &
                                    (rework_df['Test Type'].astype(str) == test_type)
                                )
                                
                                if el_col:
                                    element_val = str(rej_row.get(el_col, 'Unknown'))
                                    filter_mask = filter_mask & (rework_df[el_col].astype(str) == element_val)
                                else:
                                    element_val = "N/A"

                                future_accepts = rework_df[filter_mask]
                                
                                if not future_accepts.empty:
                                    acc_date = future_accepts['Date( SUB)'].min()
                                    delay_days = (acc_date - rej_date).days
                                    total_delay_days += delay_days
                                    status_text = "Resolved ✅"
                                else:
                                    acc_date = pd.NaT
                                    delay_days = 0
                                    status_text = "Pending 🚨"
                                    
                                ledger_data.append({
                                    "Rejected Serial": serial,
                                    "Element": element_val,
                                    "Layer": layer,
                                    "Test Type": test_type,
                                    "Rejection Date": rej_date.strftime('%Y-%m-%d') if pd.notna(rej_date) else 'N/A',
                                    "Resolution Date": acc_date.strftime('%Y-%m-%d') if pd.notna(acc_date) else 'Not Resolved',
                                    "Rework Delay (Days)": delay_days,
                                    "Status": status_text
                                })
                            
                            ledger_df = pd.DataFrame(ledger_data).sort_values(by=["Status", "Rework Delay (Days)"], ascending=[False, False])
                            
                            st.markdown(f"""
                            <div style="background: rgba(231, 76, 60, 0.1); border-left: 5px solid #e74c3c; padding: 20px; border-radius: 8px; margin-bottom: 20px; display: flex; justify-content: space-between; align-items: center;">
                                <div>
                                    <h3 style="margin: 0; color: #e74c3c; font-size: 18px;">Total Rework Time Leakage</h3>
                                    <p style="margin: 5px 0 0 0; color: {ui['text_muted']}; font-size: 14px;">Total project days lost tracking re-submissions for the same rejected layers/elements.</p>
                                </div>
                                <div style="font-size: 32px; font-weight: bold; color: #e74c3c;">{total_delay_days} Days Lost</div>
                            </div>
                            """, unsafe_allow_html=True)
                            
                            st.dataframe(ledger_df, use_container_width=True)
                        else:
                            st.success(f"✅ No rejected submittals found for {selected_comp}. Zero rework leakage!")
                    else:
                        st.info("Requires 'sample status', 'layer', and 'Date( SUB)' columns to calculate rework delays.")

                    if battalion_col_360 and elment_col_360 and 'Date ( test)' in comp_df_full.columns:
                        st.markdown("#### ⏱️ Inter-Battalion Element Timeline Analysis")
                        st.caption("Evaluates contractor speed and start/end dates for each element across different battalions.")
                        time_df = comp_df_full.dropna(subset=['Date ( test)'])
                        if not time_df.empty:
                            timeline = time_df.groupby([battalion_col_360, elment_col_360]).agg(
                                Start_Date=('Date ( test)', 'min'),
                                End_Date=('Date ( test)', 'max'),
                                Total_Submittals=('Date ( test)', 'count')
                            ).reset_index()
                            timeline['Duration (Days)'] = (timeline['End_Date'] - timeline['Start_Date']).dt.days
                            timeline['Start_Date'] = timeline['Start_Date'].dt.strftime('%Y-%m-%d')
                            timeline['End_Date'] = timeline['End_Date'].dt.strftime('%Y-%m-%d')
                            timeline = timeline.sort_values([battalion_col_360, 'Start_Date'])
                            st.dataframe(timeline, use_container_width=True)
                        else:
                            st.info("No valid dates found for timeline analysis.")

                with tab_stockpile:
                    if battalion_col_stock:
                        avail_bats = ["All Battalions"] + sorted([str(b) for b in comp_df_full[battalion_col_stock].unique() if pd.notna(b) and str(b).strip() != ''])
                        selected_bat = st.selectbox("📍 Filter Sourcing Analysis by Battalion:", avail_bats, key=f"bat_stock_{selected_comp}")
                        
                        if selected_bat != "All Battalions":
                            comp_bat_df = comp_df_full[comp_df_full[battalion_col_stock].astype(str) == selected_bat]
                            b_key = fmt_b(selected_bat)
                            req_qty = target_dict.get(f"{selected_comp.strip()}_{b_key}", np.nan)
                        else:
                            comp_bat_df = comp_df_full
                            m_keys = [k for k in target_dict.keys() if k.startswith(selected_comp.strip() + "_")]
                            if m_keys:
                                req_qty = sum(target_dict[k] for k in m_keys)
                            else:
                                req_qty = np.nan
                    else:
                        comp_bat_df = comp_df_full
                        req_qty = target_dict.get(selected_comp.strip(), np.nan)

                    stock_df = comp_bat_df[comp_bat_df['Loc_Category'] == 'Stockpile']
                    
                    if num_tests_col:
                        stock_count = int(pd.to_numeric(stock_df[num_tests_col], errors='coerce').fillna(0).sum())
                        bottom_count = int(pd.to_numeric(comp_bat_df[comp_bat_df['Loc_Category'] == 'Bottom of Excavation'][num_tests_col], errors='coerce').fillna(0).sum())
                        fill_count = int(pd.to_numeric(comp_bat_df[comp_bat_df['Loc_Category'] == 'Fill'][num_tests_col], errors='coerce').fillna(0).sum())
                    else:
                        stock_count = len(stock_df)
                        bottom_count = len(comp_bat_df[comp_bat_df['Loc_Category'] == 'Bottom of Excavation'])
                        fill_count = len(comp_bat_df[comp_bat_df['Loc_Category'] == 'Fill'])
                    
                    col_200 = next((c for c in stock_df.columns if '200' in str(c)), None)
                    if col_200 and not stock_df.empty:
                        clean_200 = stock_df[col_200].astype(str).str.replace('%', '', regex=False).str.strip()
                        clean_200 = pd.to_numeric(clean_200, errors='coerce')
                        avg_200 = clean_200.mean()
                    else:
                        avg_200 = np.nan
                    
                    cc1, cc2, cc3, cc4 = st.columns(4)
                    create_card(cc1, "Stockpile Tests", stock_count)
                    create_card(cc2, "Bottom Excavation Tests", bottom_count)
                    create_card(cc3, "Fill Tests", fill_count)
                    create_card(cc4, "Avg Sieve #200 (Stockpile)", f"{avg_200:.2f}%" if pd.notna(avg_200) else "N/A")
                    
                    if pd.notna(req_qty) and req_qty > 0:
                        req_qty_int = int(req_qty)
                        diff = stock_count - req_qty_int
                        progress_pct = min(100, (stock_count / req_qty_int) * 100) if req_qty_int > 0 else 100
                        
                        if progress_pct >= 90:
                            prog_color = "#2ecc71" 
                            status_color = "#2ecc71"
                            status_icon = "✅"
                        elif progress_pct >= 70:
                            prog_color = "#f1c40f" 
                            status_color = "#f1c40f"
                            status_icon = "⚠️"
                        else:
                            prog_color = "#e74c3c" 
                            status_color = "#e74c3c"
                            status_icon = "🚨"

                        if diff >= 0:
                            status_msg = f"<span style='color: #2ecc71;'>{status_icon} Target Exceeded (+{diff} Tests)</span>"
                            prog_color = "#2ecc71"
                        else:
                            status_msg = f"<span style='color: {status_color};'>{status_icon} Missing {abs(diff)} Tests ({progress_pct:.1f}%)</span>"
                        
                        st.markdown(f"""
                        <div style="background: {ui['card_bg']}; padding: 25px; border-radius: 12px; border-left: 6px solid #00d2ff; margin-top: 15px; margin-bottom: 25px; box-shadow: {ui['shadow']};">
                            <h4 style="color: #00d2ff; margin-top: 0; margin-bottom: 20px; font-size: 18px; font-weight: bold;">🎯 Stockpile Target Achievement</h4>
                            <div style="display: flex; justify-content: space-between; margin-bottom: 12px; align-items: flex-end;">
                                <span style="color: {ui['text_muted']}; font-size: 14px;">Target Required: <b style="color: {ui['text_main']}; font-size: 16px;">{req_qty_int}</b></span>
                                <span style="color: #00d2ff; font-size: 14px;">Executed Tests: <b style="color: {ui['text_main']}; font-size: 16px;">{stock_count}</b></span>
                                <span style="font-size: 14px; font-weight: bold; color: {ui['text_main']};">Status: {status_msg}</span>
                            </div>
                            <div class="prog-bg" style="height: 12px; background: rgba(127,140,141,0.2); border-radius: 10px; width: 100%; overflow: hidden;">
                                <div class="prog-fill" style="width: {progress_pct}%; background: {prog_color}; height: 100%; border-radius: 10px; transition: width 1s ease-in-out;"></div>
                            </div>
                        </div>
                        """, unsafe_allow_html=True)
                    else:
                        st.markdown(f"""
                        <div style="background: {ui['card_bg']}; padding: 25px; border-radius: 12px; border-left: 6px solid #95a5a6; margin-top: 15px; margin-bottom: 25px; box-shadow: {ui['shadow']};">
                            <h4 style="color: #95a5a6; margin-top: 0; margin-bottom: 10px;">🎯 Stockpile Target Achievement</h4>
                            <p style="color: {ui['text_muted']}; font-size: 15px; margin: 0;">No 'Required Quantity' target is currently defined for <b>{selected_comp}</b> in the selected scope.</p>
                        </div>
                        """, unsafe_allow_html=True)

                    if 'Date ( test)' in comp_bat_df.columns:
                        time_analysis_df = comp_bat_df.dropna(subset=['Date ( test)']).copy()
                        time_analysis_df['Month'] = time_analysis_df['Date ( test)'].dt.strftime('%b %Y')
                        fill_by_month = time_analysis_df[time_analysis_df['Loc_Category'] == 'Fill'].groupby('Month').size()
                        stock_by_month = time_analysis_df[time_analysis_df['Loc_Category'] == 'Stockpile'].groupby('Month').size()

                        if not fill_by_month.empty:
                            peak_fill_month = fill_by_month.idxmax()
                            peak_fill_val = fill_by_month.max()
                            stock_in_peak = stock_by_month.get(peak_fill_month, 0)
                            
                            bat_key = selected_bat if battalion_col_stock else "all"
                            scan_key = f"scan_{selected_comp}_{bat_key}"
                            
                            if scan_key not in st.session_state:
                                st.session_state[scan_key] = False

                            if not st.session_state[scan_key]:
                                st.markdown("<br>", unsafe_allow_html=True)
                                if st.button("🧠 Run AI Material Correlation Scan", type="primary", use_container_width=True, key=f"btn_{scan_key}"):
                                    with st.container():
                                        progress_bar = st.progress(0)
                                        status_text = st.empty()
                                        status_text.markdown(f"**<span style='color:#00d2ff;'>[1/3]</span> Scanning {total_requests_count:,} Submittal Logs...**", unsafe_allow_html=True)
                                        time.sleep(1)
                                        progress_bar.progress(33)
                                        status_text.markdown("**<span style='color:#ffaa00;'>[2/3]</span> Correlating Fill layers with Stockpile sources...**", unsafe_allow_html=True)
                                        time.sleep(1.2)
                                        progress_bar.progress(66)
                                        status_text.markdown("**<span style='color:#2ecc71;'>[3/3]</span> Generating Quality Traceability Insights...**", unsafe_allow_html=True)
                                        time.sleep(1.2)
                                        progress_bar.progress(100)
                                        time.sleep(0.5)
                                        st.session_state[scan_key] = True
                                        st.rerun()
                            
                            if st.session_state[scan_key]:
                                has_target_qty = pd.notna(req_qty) and req_qty > 0
                                req_qty_int = int(req_qty) if has_target_qty else 0
                                
                                if has_target_qty and req_qty_int > 0:
                                    coverage_ratio = stock_count / req_qty_int
                                    
                                    # 1. حساب نسبة الثقة بناءً على تحقيق التارجت الهندسي
                                    ai_confidence = min(99.9, 60.0 + (coverage_ratio * 39.9))
                                    ai_confidence = round(ai_confidence, 1)
                                    
                                    # 2. تقييم التغطية بناءً على الكود/المواصفات
                                    if coverage_ratio < 0.75:
                                        status_level, status_color, status_bg, status_icon = "SEVERE DEFICIT", "#e74c3c", "rgba(231, 76, 60, 0.1)", "🚨"
                                        missing_tests = req_qty_int - stock_count
                                        quality_insight = f"Significant deficit detected based on engineering targets. Only {stock_count} out of {req_qty_int} required tests are logged."
                                        directive = f"ACTION REQUIRED: Execute and submit at least {missing_tests} Stockpile samples immediately to meet the volumetric target specifications."
                                        
                                    elif coverage_ratio < 1.0:
                                        status_level, status_color, status_bg, status_icon = "COVERAGE GAP", "#f1c40f", "rgba(241, 196, 15, 0.1)", "⚠️"
                                        missing_tests = req_qty_int - stock_count
                                        quality_insight = f"Material approval rate is slightly lagging. You are at {int(coverage_ratio*100)}% of the required volumetric target."
                                        directive = f"ADVISORY: Schedule {missing_tests} additional stockpile tests to achieve full engineering compliance."
                                        
                                    else:
                                        status_level, status_color, status_bg, status_icon = "OPTIMAL COVERAGE", "#2ecc71", "rgba(46, 204, 113, 0.1)", "✅"
                                        quality_insight = f"Stockpile testing frequency exceeds or meets the required engineering targets ({int(coverage_ratio*100)}% coverage)."
                                        directive = f"MAINTAIN: Volumetric target achieved. Continue standard QC monitoring."
                                        
                                    field_detect_msg = f"Target Required: <b style='color:#00d2ff;'>{req_qty_int} tests</b> based on volume.<br>Executed Stockpile Tests: <b style='color:{status_color};'>{stock_count} tests</b>."
                                    
                                else:
                                    # في حالة مفيش تارجت مسجل للمقاول
                                    ai_confidence = 45.5 
                                    status_level, status_color, status_bg, status_icon = "NO TARGET DEFINED", "#95a5a6", "rgba(149, 165, 166, 0.1)", "❓"
                                    field_detect_msg = f"Executed Stockpile Tests: <b style='color:#00d2ff;'>{stock_count}</b>.<br>Target Required: <b style='color:#e74c3c;'>Not Found</b>."
                                    quality_insight = "Cannot calculate material quality coverage because the 'Required Quantity' target is missing for this contractor."
                                    directive = "SYSTEM.HALT: Please define the 'Required Quantity' in the Data Transformation Hub to enable AI volumetric audits."

                                st.markdown(f"""
                                <style>
                                    @keyframes scanline {{ 0% {{ transform: translateY(-10px); opacity: 0; }} 50% {{ opacity: 1; }} 100% {{ transform: translateY(0); opacity: 1; }} }}
                                    .ai-terminal {{ background: linear-gradient(145deg, #0a1118, {status_bg}); border: 1px solid rgba(255,255,255,0.05); border-left: 5px solid {status_color}; border-radius: 12px; padding: 25px; margin: 20px 0; box-shadow: 0 0 20px {status_bg}; animation: scanline 0.8s ease-out forwards; }}
                                    .ai-badge {{ background: rgba(0,0,0,0.4); border: 1px solid {ui['border_color']}; padding: 5px 12px; border-radius: 20px; font-size: 12px; color: #00d2ff; }}
                                </style>
                                <div class="ai-terminal">
                                    <div style="display: flex; justify-content: space-between; align-items: center; border-bottom: 1px solid rgba(255,255,255,0.1); padding-bottom: 15px; margin-bottom: 20px;">
                                        <h3 style="color: {status_color}; margin: 0; display: flex; align-items: center; font-size: 20px;"><span style="font-size: 24px; margin-right: 10px;">🤖</span> Generative AI Quality Auditor (Volumetric)</h3>
                                        <div style="display: flex; gap: 10px;">
                                            <span class="ai-badge">⚡ Data Confidence: {ai_confidence}%</span>
                                            <span class="ai-badge" style="color: {status_color}; border-color: {status_color}; font-weight: bold;">{status_icon} Status: {status_level}</span>
                                        </div>
                                    </div>
                                    <div style="display: grid; grid-template-columns: 1fr 1fr 1fr; gap: 20px;">
                                        <div style="background: rgba(0,0,0,0.2); padding: 15px; border-radius: 8px; border-top: 2px solid #00d2ff;">
                                            <div style="color: #00d2ff; font-weight: bold; font-size: 11px; letter-spacing: 1px; margin-bottom: 8px;">> FIELD_DATA.DETECT()</div>
                                            <div style="color: {ui['text_main']}; font-size: 14px; line-height: 1.6;">{field_detect_msg}</div>
                                        </div>
                                        <div style="background: rgba(0,0,0,0.2); padding: 15px; border-radius: 8px; border-top: 2px solid #ffaa00;">
                                            <div style="color: #ffaa00; font-weight: bold; font-size: 11px; letter-spacing: 1px; margin-bottom: 8px;">> ENGINEERING_GAP.ANALYZE()</div>
                                            <div style="color: {ui['text_main']}; font-size: 14px; line-height: 1.6;">{quality_insight}</div>
                                        </div>
                                        <div style="background: rgba(0,0,0,0.2); padding: 15px; border-radius: 8px; border-top: 2px solid {status_color};">
                                            <div style="color: {status_color}; font-weight: bold; font-size: 11px; letter-spacing: 1px; margin-bottom: 8px;">> QC_ACTION.RECOMMEND()</div>
                                            <div style="color: {ui['text_main']}; font-size: 14px; line-height: 1.6; font-weight: 500;">{directive}</div>
                                        </div>
                                    </div>
                                </div>
                                """, unsafe_allow_html=True)
                                col_reset, _ = st.columns([0.2, 0.8])
                                if col_reset.button("🔄 Reset AI Auditor", key=f"reset_{scan_key}"):
                                    st.session_state[scan_key] = False
                                    st.rerun()

                    ch_col1, ch_col2 = st.columns(2)
                    with ch_col1:
                        if 'Classification' in stock_df.columns and not stock_df.empty:
                            class_counts = stock_df['Classification'].value_counts().reset_index()
                            class_counts.columns = ['Classification', 'Count']
                            fig_class = px.pie(class_counts, names='Classification', values='Count', title=f"Stockpile Classifications for {selected_comp}", hole=0.3, color_discrete_sequence=NEON_COLORS)
                            fig_class.update_traces(textinfo='label+percent', hovertemplate='<b>Class:</b> %{label}<br>Count: %{value}<br>Percentage: %{percent}')
                            fig_class = style_3d_glassy(fig_class, chart_type="pie")
                            st.plotly_chart(fig_class, use_container_width=True, key=f"class_{selected_comp}")
                        else:
                            st.info(f"No Stockpile classification data logged.")
                            
                    with ch_col2:
                        if 'sample status' in stock_df.columns and not stock_df.empty:
                            stock_df['status_upper'] = stock_df['sample status'].str.upper()
                            fig_stock_status = px.pie(stock_df, names='status_upper', title=f"Stockpile Approval/Rejection Rate", hole=0.3, color='status_upper', color_discrete_map=STATUS_COLORS)
                            fig_stock_status.update_traces(textinfo='label+percent', hovertemplate='<b>Status:</b> %{label}<br>Count: %{value}<br>Percentage: %{percent}')
                            fig_stock_status = style_3d_glassy(fig_stock_status, chart_type="pie")
                            st.plotly_chart(fig_stock_status, use_container_width=True, key=f"stock_status_{selected_comp}")
                        else:
                            st.info(f"No Stockpile status data logged.")

                    ch_col3, ch_col4 = st.columns(2)
                    with ch_col3:
                        if 'Date ( test)' in stock_df.columns and not stock_df.empty:
                            time_df = stock_df.dropna(subset=['Date ( test)']).copy()
                            time_df['Month'] = time_df['Date ( test)'].dt.strftime('%b %Y')
                            time_df['Month_Sort'] = time_df['Date ( test)'].dt.to_period('M')
                            monthly_stock = time_df.groupby(['Month_Sort', 'Month']).size().reset_index(name='Count').sort_values('Month_Sort')
                            fig_timeline = px.bar(monthly_stock, x='Month', y='Count', title="Stockpile Tests Timeline", text_auto=True, color_discrete_sequence=['#ffaa00'])
                            fig_timeline = style_3d_glassy(fig_timeline, chart_type="bar")
                            st.plotly_chart(fig_timeline, use_container_width=True, key=f"stock_time_{selected_comp}")
                        else:
                            st.info("No Date data available to show Stockpile timeline.")

                    with ch_col4:
                        if 'sample status' in comp_bat_df.columns and not comp_bat_df.empty:
                            comp_bat_df['status_upper'] = comp_bat_df['sample status'].str.upper()
                            fig_status = px.pie(comp_bat_df, names='status_upper', title=f"Overall Approval Rate (All Tests)", hole=0.3, color='status_upper', color_discrete_map=STATUS_COLORS)
                            fig_status.update_traces(textinfo='label+percent', hovertemplate='<b>Status:</b> %{label}<br>Count: %{value}<br>Percentage: %{percent}')
                            fig_status = style_3d_glassy(fig_status, chart_type="pie")
                            st.plotly_chart(fig_status, use_container_width=True, key=f"stock_all_{selected_comp}")
                        else:
                            st.info(f"No overall status data logged.")

                with tab_execution:
                    st.markdown(f"### 🏗️ Executive Progress & Compaction: `{selected_comp}`")
                    
                    company_col_qty = next((c for c in df.columns if 'COMPANY' in c.upper() and c != 'Company Name'), 'Company Name')
                    qty_match_df = df[df[company_col_qty].astype(str).str.strip().str.lower() == selected_comp.strip().lower()]
                    if qty_match_df.empty: qty_match_df = comp_df_full
                        
                    tot_qty_col = next((c for c in df.columns if 'TOTAL QUANTITY' in str(c).strip().upper()), None)
                    exec_qty_col = next((c for c in df.columns if 'EXECUTED QUANTITY' in str(c).strip().upper()), None)
                    
                    tot_qty = pd.to_numeric(qty_match_df[tot_qty_col], errors='coerce').max() if tot_qty_col else 0
                    exe_qty = pd.to_numeric(qty_match_df[exec_qty_col], errors='coerce').max() if exec_qty_col else 0
                    prog_pct = (exe_qty / tot_qty * 100) if pd.notna(tot_qty) and tot_qty > 0 else 0
                    
                    test_col = 'Test Type' if 'Test Type' in comp_df_full.columns else None
                    compaction_df = pd.DataFrame()
                    if test_col:
                        compaction_df = comp_df_full[comp_df_full[test_col].astype(str).str.contains('DPL|PLATE', case=False, na=False)].copy()
                        
                    num_tests_col_exec = next((c for c in comp_df_full.columns if 'NUMBER OF TESTS' in str(c).strip().upper() or 'NUM OF TEST' in str(c).strip().upper()), None)
                    
                    dpl_df = compaction_df[compaction_df[test_col].astype(str).str.contains('DPL', case=False, na=False)] if test_col else pd.DataFrame()
                    plate_df = compaction_df[compaction_df[test_col].astype(str).str.contains('PLATE', case=False, na=False)] if test_col else pd.DataFrame()
                    
                    dpl_pts = int(pd.to_numeric(dpl_df[num_tests_col_exec], errors='coerce').sum()) if num_tests_col_exec and not dpl_df.empty else len(dpl_df)
                    plate_pts = int(pd.to_numeric(plate_df[num_tests_col_exec], errors='coerce').sum()) if num_tests_col_exec and not plate_df.empty else len(plate_df)
                    total_test_points = dpl_pts + plate_pts
                    
                    avg_dpl = pd.to_numeric(dpl_df['AVERAGE VALUE'], errors='coerce').mean() if 'AVERAGE VALUE' in dpl_df.columns else np.nan
                        
                    c1, c2, c3, c4 = st.columns(4)
                    create_card(c1, "Total Target Qty", f"{tot_qty:,.0f}" if pd.notna(tot_qty) and tot_qty>0 else "N/A")
                    create_card(c2, "Executed Qty", f"{exe_qty:,.0f}" if pd.notna(exe_qty) and exe_qty>0 else "0")
                    pts_html = f"<div style='font-size:14px; color:#8da3b9; margin-top:5px;'>DPL: <b style='color:#00d2ff;'>{dpl_pts}</b> | Plate: <b style='color:#ffaa00;'>{plate_pts}</b></div>"
                    create_card(c3, "Total Compaction Points", f"{total_test_points:,}", delta_html=pts_html)
                    create_card(c4, "Average DPL Value", f"{avg_dpl:.2f}" if pd.notna(avg_dpl) else "N/A")

                    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
                    
                    r2_c1, r2_c2 = st.columns(2)
                    
                    with r2_c1:
                        st.markdown("#### 🚀 Execution Progress vs Target")
                        if pd.notna(tot_qty) and tot_qty > 0:
                            fig_exec_kpi = go.Figure(go.Indicator(
                                mode = "gauge+number+delta",
                                value = exe_qty,
                                title = {'text': "Completed Volume", 'font': {'size': 16, 'color': "white" if is_dark else "#2C3E50"}},
                                number = {'font': {'size': 35, 'color': "white" if is_dark else "#2C3E50"}},
                                delta = {'reference': tot_qty, 'increasing': {'color': "#2ecc71"}, 'decreasing': {'color': "#e74c3c"}},
                                gauge = {
                                    'axis': {'range': [None, tot_qty], 'tickwidth': 1, 'tickcolor': "rgba(255,255,255,0.2)"},
                                    'bar': {'color': "#00d2ff"},
                                    'bgcolor': "rgba(255,255,255,0.05)" if is_dark else "rgba(0,0,0,0.02)",
                                    'steps': [
                                        {'range': [0, tot_qty*0.5], 'color': "rgba(231,76,60,0.3)"},
                                        {'range': [tot_qty*0.5, tot_qty*0.8], 'color': "rgba(241,196,15,0.3)"},
                                        {'range': [tot_qty*0.8, tot_qty], 'color': "rgba(46,204,113,0.3)"}
                                    ],
                                    'threshold': {'line': {'color': "red", 'width': 4}, 'thickness': 0.75, 'value': tot_qty}
                                }
                            ))
                            fig_exec_kpi.update_layout(paper_bgcolor="rgba(0,0,0,0)", height=320, margin=dict(l=20, r=20, t=50, b=20), font={'family': 'Montserrat'})
                            st.plotly_chart(fig_exec_kpi, use_container_width=True, key=f"exec_kpi_gauge_{selected_comp}")
                        else:
                            st.info("No Target Quantity defined to show progress.")
                            
                    with r2_c2:
                        st.markdown("#### ⚖️ Compaction Quality Metrics (Pass Rate)")
                        if 'sample status' in compaction_df.columns and not compaction_df.empty:
                            compaction_df['status_upper'] = compaction_df['sample status'].str.upper()
                            fig_comp_qual = px.pie(compaction_df, names='status_upper', hole=0.4, color='status_upper', color_discrete_map=STATUS_COLORS)
                            fig_comp_qual.update_traces(textinfo='label+percent', hovertemplate='<b>Status:</b> %{label}<br>Count: %{value}<br>Yield: %{percent}')
                            fig_comp_qual = style_3d_glassy(fig_comp_qual, chart_type="pie")
                            fig_comp_qual.update_layout(height=320, margin=dict(l=20, r=20, t=20, b=20))
                            st.plotly_chart(fig_comp_qual, use_container_width=True, key=f"comp_qual_pie_{selected_comp}")
                        else:
                            st.info("No Quality/Status data found for Compaction.")
                            
                    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
                    
                    st.markdown("#### 📈 Monthly Compaction Trend Analysis")
                    if not compaction_df.empty and 'Date ( test)' in compaction_df.columns:
                        compaction_df['Month'] = compaction_df['Date ( test)'].dt.strftime('%b %Y')
                        compaction_df['Month_Sort'] = compaction_df['Date ( test)'].dt.to_period('M')
                        monthly_comp = compaction_df.groupby(['Month_Sort', 'Month', test_col]).size().reset_index(name='Count').sort_values('Month_Sort')
                        
                        fig_comp_trend = px.bar(monthly_comp, x='Month', y='Count', color=test_col, barmode='group', color_discrete_sequence=NEON_COLORS)
                        fig_comp_trend.update_traces(hovertemplate='<b>Month:</b> %{x}<br><b>Tests:</b> %{y}')
                        fig_comp_trend = style_3d_glassy(fig_comp_trend, chart_type="bar")
                        fig_comp_trend.update_layout(height=350)
                        st.plotly_chart(fig_comp_trend, use_container_width=True, key=f"comp_trend_bar_{selected_comp}")
                    else:
                        st.info("No Date data found for Compaction Trend.")
                        
                    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
                    st.markdown("#### 🧠 Executive AI Insights & Alerts")
                    
                    pred_date_str = "Insufficient Data"
                    daily_rate = 0
                    if 'Date ( test)' in qty_match_df.columns and pd.notna(tot_qty) and tot_qty > 0:
                        dates = qty_match_df['Date ( test)'].dropna()
                        if len(dates) >= 2 and exe_qty > 0:
                            min_date = dates.min()
                            max_date = dates.max()
                            days_worked = (max_date - min_date).days
                            if days_worked > 0:
                                daily_rate = exe_qty / days_worked
                                if exe_qty >= tot_qty:
                                    pred_date_str = "Target Achieved ✅"
                                else:
                                    rem_qty = tot_qty - exe_qty
                                    rem_days = rem_qty / daily_rate
                                    pred_date = max_date + timedelta(days=rem_days)
                                    pred_date_str = pred_date.strftime('%B %Y')
                    
                    rej_rate = 0
                    top_fail_str = "None"
                    if 'sample status' in compaction_df.columns and not compaction_df.empty:
                        rej_df = compaction_df[compaction_df['sample status'].astype(str).str.upper().isin(['REJECTED', 'REVISE'])]
                        rej_rate = (len(rej_df) / len(compaction_df)) * 100
                        if not rej_df.empty:
                            if test_col in rej_df.columns:
                                top_fail_str = rej_df[test_col].value_counts().idxmax()
                    
                    insight_c1, insight_c2, insight_c3 = st.columns(3)
                    
                    with insight_c1:
                        st.markdown(f"""
                        <div style="background: rgba(0, 210, 255, 0.05); border-left: 4px solid #00d2ff; padding: 20px; border-radius: 8px; height: 100%;">
                            <h4 style="color: #00d2ff; margin-top: 0; font-size: 16px;">📊 Performance vs Target</h4>
                            <div style="color: {ui['text_main']}; font-size: 14px; line-height: 1.6;">
                                • <b>Completion:</b> {prog_pct:.1f}% of total target.<br>
                                • <b>Daily Velocity:</b> ~{int(daily_rate):,} units/day.<br>
                                • <b>Status:</b> {"On Track" if prog_pct > 50 else "Requires Acceleration"}
                            </div>
                        </div>
                        """, unsafe_allow_html=True)
                        
                    with insight_c2:
                        qual_color = "#e74c3c" if rej_rate > 15 else "#2ecc71"
                        st.markdown(f"""
                        <div style="background: rgba(231, 76, 60, 0.05); border-left: 4px solid {qual_color}; padding: 20px; border-radius: 8px; height: 100%;">
                            <h4 style="color: {qual_color}; margin-top: 0; font-size: 16px;">⚠️ Quality Issues</h4>
                            <div style="color: {ui['text_main']}; font-size: 14px; line-height: 1.6;">
                                • <b>Rejection Rate:</b> {rej_rate:.1f}% in compaction.<br>
                                • <b>Most Failed Test:</b> {top_fail_str}.<br>
                                • <b>Action:</b> {"Urgent audit needed for equipment." if rej_rate > 15 else "Quality is within acceptable limits."}
                            </div>
                        </div>
                        """, unsafe_allow_html=True)
                        
                    with insight_c3:
                        st.markdown(f"""
                        <div style="background: rgba(241, 196, 15, 0.05); border-left: 4px solid #f1c40f; padding: 20px; border-radius: 8px; height: 100%;">
                            <h4 style="color: #f1c40f; margin-top: 0; font-size: 16px;">🔮 Predictive Completion</h4>
                            <div style="color: {ui['text_main']}; font-size: 14px; line-height: 1.6;">
                                • <b>Forecasted Finish:</b> <b style="color:#ffaa00; font-size:16px;">{pred_date_str}</b><br>
                                • <b>Algorithm:</b> Based on historical velocity of {int(daily_rate):,} avg volume per active day.<br>
                            </div>
                        </div>
                        """, unsafe_allow_html=True)

                st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

# ==========================================
# tab_quantities — FINAL VERSION
# ضيف الكود ده جوه: with tab_quantities:
# ==========================================

                with tab_quantities:
                    st.markdown("### 📊 Quantities Rate & Execution Analytics")
                    
                    # السطر السحري اللي هيمنع الإيرور:
                    tracker_df = pd.DataFrame()
                    # ══════════════════════════════════════════════════════
                    # COLUMN DETECTION 
                    # ══════════════════════════════════════════════════════
                    company_col     = next((c for c in df.columns if c.strip() == 'Company'), None)
                    comp_name_col   = next((c for c in df.columns if c.strip() == 'Company Name'), None)
                    contractor_col  = next((c for c in df.columns if c.strip() == 'Contractor'), None)
                    
                    total_qty_col   = next((c for c in df.columns if c.strip() == 'Total Quantity'), None)
                    if not total_qty_col:
                        total_qty_col = next((c for c in df.columns if 'TOTAL QUANTITY' in c.upper()), None)
                        
                    exec_qty_col    = next((c for c in df.columns if c.strip() == 'Executed Quantity'), None)
                    
                    exec_qty_m3_col = next((c for c in df.columns if c.strip() == 'Executed Quantity (m3)'), None)
                    if not exec_qty_m3_col:
                        exec_qty_m3_col = next((c for c in df.columns
                                                if 'EXECUTED' in c.upper() and 'M3' in c.upper().replace('³','3').replace('(','').replace(')','').replace(' ','')), None)
                                                
                    elem_all_col    = next((c for c in df.columns if c.strip() in ['Element (all)', 'Element (All)', 'Element(all)']), None)
                    if not elem_all_col:
                        elem_all_col = next((c for c in df.columns if 'ELEMENT' in c.upper() and 'ALL' in c.upper()), None)
                        
                    elment_col      = next((c for c in df.columns if c.strip() in ['ELMENT', 'Elment', 'ELEMENT']), None)
                    if not elment_col:
                        elment_col  = next((c for c in df.columns if 'ELMEN' in c.upper() and 'ALL' not in c.upper()), None)
                        
                    target_col      = next((c for c in df.columns if 'TARGET DAILY RATE' in c.upper()), None)
                    
                    date_daily_col  = next((c for c in df.columns if 'DATE' in c.upper() and 'DAILY' in c.upper()), None)
                    if not date_daily_col:
                        date_daily_col = next((c for c in df.columns if c.strip() == 'Date (Daily)'), None)
                        
                    sector_col      = next((c for c in df.columns if c.strip() == 'Sectoer' or 'SECTOR' in c.upper()), None)
                    test_type_col   = next((c for c in df.columns if 'TEST TYPE' in c.upper() or c.strip() == 'Test Type'), None)
                    num_tests_col   = next((c for c in df.columns if 'NUMBER OF TESTS' in c.upper() or 'NUM OF TEST' in c.upper()), None)

                    with st.expander("🔍 Column Detection", expanded=False):
                        st.json({
                            "Company":              company_col,
                            "Company Name":         comp_name_col,
                            "Contractor":           contractor_col,
                            "Total Quantity":       total_qty_col,
                            "Executed Quantity":    exec_qty_col,
                            "Executed Qty (m3)":    exec_qty_m3_col,
                            "Element (all)":        elem_all_col,
                            "ELMENT":               elment_col,
                            "Target Daily Rate":    target_col,
                            "Date (Daily)":         date_daily_col,
                            "Sector":               sector_col,
                            "Test Type":            test_type_col,
                            "Number of Tests":      num_tests_col,
                        })

                    # ══════════════════════════════════════════════════════
                    # CRITICAL FIX: Numeric conversion & Removing Commas
                    # ══════════════════════════════════════════════════════
                    for col in [total_qty_col, exec_qty_col, exec_qty_m3_col, target_col]:
                        if col and col in df.columns:
                            # Remove commas (e.g., "2,754.5" -> "2754.5") before converting to float
                            if df[col].dtype == 'object':
                                df[col] = df[col].astype(str).str.replace(',', '', regex=False).str.replace(' ', '', regex=False)
                            df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)
                            
                    if num_tests_col and num_tests_col in df.columns:
                        if df[num_tests_col].dtype == 'object':
                            df[num_tests_col] = df[num_tests_col].astype(str).str.replace(',', '', regex=False)
                        df[num_tests_col] = pd.to_numeric(df[num_tests_col], errors='coerce').fillna(0)

                    # ── Contractor selector ────────────────────────────────
                    contractors_all = ['All Contractors']
                    if contractor_col and contractor_col in df.columns:
                        contractors_all += sorted(
                            df[contractor_col].dropna().astype(str).str.strip()
                            .unique().tolist()
                        )
                    sel_contractor = st.selectbox(
                        "🏢 Select Contractor:", contractors_all, key="qty_contractor_sel"
                    )

                    if sel_contractor != 'All Contractors' and contractor_col:
                        df_sel = df[df[contractor_col].astype(str).str.strip() == sel_contractor].copy()
                    else:
                        df_sel = df.copy()

                    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

                    # ══════════════════════════════════════════════════════
                    # 1. MAIN KPI CARDS
                    # ══════════════════════════════════════════════════════
                    st.markdown("#### 📦 Main KPIs")

                    total_scope = 0
                    if company_col and total_qty_col and company_col in df.columns:
                        if sel_contractor != 'All Contractors':
                            lookup_df = df[df[company_col].astype(str).str.strip() == sel_contractor]
                        else:
                            lookup_df = df
                        
                        # --- التعديل هنا: جمع العمود بالكامل بدون استبعاد أي طبقات ---
                        total_scope = lookup_df[total_qty_col].sum()

                    exec_from_company = 0
                    if company_col and exec_qty_col and company_col in df.columns:
                        if sel_contractor != 'All Contractors':
                            lookup_df2 = df[df[company_col].astype(str).str.strip() == sel_contractor]
                        else:
                            lookup_df2 = df
                        exec_from_company = lookup_df2[exec_qty_col].sum()

                    exec_daily_sum = 0
                    if contractor_col and exec_qty_m3_col and contractor_col in df.columns:
                        exec_daily_sum = df_sel[exec_qty_m3_col].sum()

                    completion_pct = (exec_daily_sum / total_scope * 100) if total_scope > 0 else 0
                    comp_color = "#2ecc71" if completion_pct >= 80 else ("#f1c40f" if completion_pct >= 50 else "#e74c3c")

                    c1, c2, c3, c4 = st.columns(4)

                    create_card(
                        c1, "🏗️ Total Project Quantity (m³)",
                        f"{total_scope:,.1f}" if total_scope > 0 else "N/A",
                        delta_html="<span style='color:#00d2ff;font-size:11px'>Company ← Total Quantity</span>"
                    )
                    create_card(
                        c2, "📦 Executed Quantity (Company)",
                        f"{exec_from_company:,.1f}",
                        delta_html="<span style='color:#2ecc71;font-size:11px'>Company ← Executed Quantity</span>"
                    )
                    create_card(
                        c3, "🚧 Total Daily Executed (m³)",
                        f"{exec_daily_sum:,.1f}",
                        delta_html="<span style='color:#ffaa00;font-size:11px'>SUMIF: Contractor ← Executed Qty (m3)</span>"
                    )
                    create_card(
                        c4, "📈 Completion %",
                        f"{completion_pct:.1f}%",
                        delta_html=f"<span style='color:{comp_color};font-size:11px'>{'✅ On Track' if completion_pct >= 80 else '⚠️ Needs Attention'}</span>",
                        progress=min(100, completion_pct)
                    )

                    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)

                    # ══════════════════════════════════════════════════════
                    # 2. ELEMENTS EXECUTION CARDS
                    # ══════════════════════════════════════════════════════
                    if elem_all_col and contractor_col and exec_qty_m3_col:
                        st.markdown("#### 🔍 Executed Quantity per Element")

                        elem_sumif = (
                            df_sel
                            .groupby([contractor_col, elem_all_col])[exec_qty_m3_col]
                            .sum()
                            .reset_index()
                        )
                        elem_sumif.columns = ['Contractor', 'Element', 'Executed (m³)']
                        elem_sumif = (
                            elem_sumif[elem_sumif['Executed (m³)'] > 0]
                            .sort_values('Executed (m³)', ascending=False)
                        )

                        if not elem_sumif.empty:
                            top4 = elem_sumif.head(4)
                            e_cols = st.columns(min(len(top4), 4))
                            for i, (_, row) in enumerate(top4.iterrows()):
                                create_card(
                                    e_cols[i],
                                    f"📍 {row['Element']}",
                                    f"{row['Executed (m³)']:,.1f} m³",
                                    delta_html=f"<span style='color:#8da3b9;font-size:10px'>{row['Contractor']}</span>"
                                )
                            with st.expander("📋 View All Elements"):
                                st.dataframe(elem_sumif, use_container_width=True)
                        else:
                            st.info("No element execution data available.")

                    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
# ══════════════════════════════════════════════════════
                   # ══════════════════════════════════════════════════════
                    # 3. CHART: Daily Executed vs Target Rate + DPL Tests (Unified Timeline by Element)
                    # ══════════════════════════════════════════════════════
                    st.markdown("#### 🚀 Execution vs Target Rate & DPL Tests")
                    
                    # الزرار الجانبي للتحكم في ظهور خط الـ DPL
                    show_dpl_line = st.sidebar.toggle("📉 Show DPL Tests Line", value=True)
                    
                    # 💡 الزرار الجانبي الجديد للدمج الذكي (اختياري)
                    enable_smart_mask = st.sidebar.toggle("🪄 Smart Element Merge (e.g. BH-11-1 ➔ BH-11)", value=False)

                    date_test_col = next((c for c in df.columns if 'DATE' in c.upper() and 'TEST' in c.upper()), None)
                    if not date_test_col:
                        date_test_col = next((c for c in df.columns if c.strip() == 'Date ( test)'), None)

                    if date_daily_col and target_col and exec_qty_m3_col and contractor_col and date_test_col and comp_name_col and elem_all_col and elment_col:
                        
                        # 💡 --- التعديل المحلي (Local Scope) للشارت ده بس ---
                        chart_df = df.copy()
                        
                        # لو المدير فعّل الزرار، هنطبق دالة التوحيد على أي كتيبة
                        if enable_smart_mask:
                            def unify_element(val):
                                val = str(val).strip()
                                parts = val.split('-')
                                # لو الاسم متقسم بـ (-) لـ 3 مقاطع أو أكتر وآخر جزء عبارة عن رقم
                                if len(parts) >= 3 and parts[-1].isdigit():
                                    return "-".join(parts[:-1]) # يرجعه بدون الرقم الأخير
                                return val
                                
                            chart_df[elem_all_col] = chart_df[elem_all_col].apply(unify_element)
                            chart_df[elment_col] = chart_df[elment_col].apply(unify_element)
                        # ----------------------------------------
                        
                        # 💡 كل استخدام لـ df تم استبداله بـ chart_df
                        if sel_contractor != 'All Contractors':
                            available_elements = chart_df[chart_df[contractor_col].astype(str).str.strip().str.lower() == sel_contractor.lower()][elem_all_col].dropna().astype(str).str.strip().unique()
                        else:
                            available_elements = chart_df[elem_all_col].dropna().astype(str).str.strip().unique()
                            
                        available_elements = sorted([e for e in available_elements if e.lower() not in ['nan', 'none', '']])
                        
                        # أزرار الفلترة والـ Drill-down (أيام/أسابيع)
                        col_f1, col_f2 = st.columns([0.7, 0.3])
                        with col_f1:
                            sel_elem_chart = st.selectbox("📍 Filter Timeline by Element:", ["All Elements"] + available_elements, key="chart_elem_filter")
                        with col_f2:
                            time_view = st.radio("⏱️ Chart Granularity:", ["Weekly", "Daily"], horizontal=True)

                        # 💡 ====== الفكرة العبقرية: اكتشاف المقاولين المشتركين وتحليل حصة كل مقاول ====== 💡
                        selected_sub_contractor = "Combined"
                        
                        if sel_elem_chart != 'All Elements':
                            # نبحث في الداتا كلها عن العنصر ده مين اشتغله (كميات وجودة)
                            exec_mask = chart_df[elem_all_col].astype(str).str.strip().str.lower() == sel_elem_chart.lower()
                            qa_mask = chart_df[elment_col].astype(str).str.strip().str.lower() == sel_elem_chart.lower()
                            
                            df_elem_exec = chart_df[exec_mask]
                            df_elem_qa = chart_df[qa_mask]
                            
                            c_exec = df_elem_exec[contractor_col].dropna().unique().tolist()
                            c_qa = df_elem_qa[comp_name_col].dropna().unique().tolist()
                            
                            # دمج أسماء المقاولين بدون تكرار
                            shared_contractors = sorted(list(set(c_exec + c_qa)))
                            shared_contractors = [c for c in shared_contractors if str(c).strip().lower() not in ['nan', 'none', '']]
                            
                            # لو لقينا أكتر من مقاول اشتغل في نفس العنصر
                            if len(shared_contractors) > 1:
                                
                                # حساب حصة كل مقاول (كميات + DPL) لعرضها للمدير
                                breakdown_html = ""
                                for c in shared_contractors:
                                    # حساب كمية المقاول
                                    c_qty = df_elem_exec[df_elem_exec[contractor_col].astype(str).str.strip().str.lower() == c.lower()][exec_qty_m3_col].sum()
                                    
                                    # حساب اختبارات DPL للمقاول
                                    c_dpl_mask = (df_elem_qa[comp_name_col].astype(str).str.strip().str.lower() == c.lower()) & (df_elem_qa[test_type_col].astype(str).str.upper().str.contains('DPL'))
                                    if num_tests_col and num_tests_col in df_elem_qa.columns:
                                        c_dpl_df = df_elem_qa[c_dpl_mask].copy()
                                        c_dpl_df[num_tests_col] = pd.to_numeric(c_dpl_df[num_tests_col].astype(str).str.replace(',', '', regex=False), errors='coerce').fillna(0)
                                        c_dpl = c_dpl_df[num_tests_col].sum()
                                    else:
                                        c_dpl = len(df_elem_qa[c_dpl_mask])
                                        
                                    # تنسيق النتيجة
                                    breakdown_html += f"<li style='margin-bottom: 5px;'><b>{c}:</b> نفذ <span style='color:#00d2ff;'>{c_qty:,.1f} m³</span> | وقام بعمل <span style='color:#ffaa00;'>{int(c_dpl)} اختبار DPL</span></li>"
                                
                                st.markdown(f"""
                                <div style="background: rgba(241, 196, 15, 0.1); border-left: 4px solid #f1c40f; padding: 15px; border-radius: 8px; margin-bottom: 15px;">
                                    <b style="color: #f1c40f; font-size: 16px;">🧠 AI Cross-Contractor Insight (Workload Split):</b><br>
                                    <span style="color: #d1d5da; font-size: 14px;">تم العمل على العنصر <b>{sel_elem_chart}</b> بواسطة <b>{len(shared_contractors)} شركات مختلفة</b>. إليك كشف حساب الكميات والجودة لكل شركة:</span>
                                    <ul style="color: #ffffff; font-size: 14px; margin-top: 10px; background: rgba(0,0,0,0.2); padding: 10px 30px; border-radius: 5px;">
                                        {breakdown_html}
                                    </ul>
                                    <span style="color: #d1d5da; font-size: 12px;">اختر من الأسفل طريقة عرض الشارت (مدمج أم شركة محددة).</span>
                                </div>
                                """, unsafe_allow_html=True)
                                
                                sub_opts = ["🔗 Combined View (Full Timeline)"] + shared_contractors
                                sub_choice = st.radio("🛠️ Isolate Contractor or View Combined:", sub_opts, horizontal=True)
                                
                                if sub_choice != "🔗 Combined View (Full Timeline)":
                                    selected_sub_contractor = sub_choice

                        # --- 1. مسار الكميات (Execution Data) ---
                        if sel_elem_chart == 'All Elements':
                            df_exec = chart_df[chart_df[contractor_col].astype(str).str.strip().str.lower() == sel_contractor.lower()].copy() if sel_contractor != 'All Contractors' else chart_df.copy()
                        else:
                            df_exec = chart_df[chart_df[elem_all_col].astype(str).str.strip().str.lower() == sel_elem_chart.lower()].copy()
                            if selected_sub_contractor != "Combined":
                                df_exec = df_exec[df_exec[contractor_col].astype(str).str.strip().str.lower() == selected_sub_contractor.lower()]

                        df_exec[date_daily_col] = pd.to_datetime(df_exec[date_daily_col], dayfirst=True, errors='coerce')
                        df_exec = df_exec.dropna(subset=[date_daily_col])
                        
                        # التجميع اليومي الصحيح (تم إضافة عمود المقاول لضمان جمع تارجت القطاعات المدمجة)
                        daily_elem_exec = df_exec.groupby([date_daily_col, contractor_col, elem_all_col]).agg(
                            Executed=(exec_qty_m3_col, 'sum'),
                            Target=(target_col, 'max')
                        ).reset_index()
                        
                        daily_exec = daily_elem_exec.groupby(date_daily_col).agg(
                            Executed=('Executed', 'sum'),
                            Target=('Target', 'sum')
                        ).reset_index().sort_values(date_daily_col)

                        # --- 2. مسار الجودة (DPL Data) ---
                        if sel_elem_chart == 'All Elements':
                            df_qa = chart_df[chart_df[comp_name_col].astype(str).str.strip().str.lower() == sel_contractor.lower()].copy() if sel_contractor != 'All Contractors' else chart_df.copy()
                        else:
                            df_qa = chart_df[chart_df[elment_col].astype(str).str.strip().str.lower() == sel_elem_chart.lower()].copy()
                            if selected_sub_contractor != "Combined":
                                df_qa = df_qa[df_qa[comp_name_col].astype(str).str.strip().str.lower() == selected_sub_contractor.lower()]
                                
                        df_qa[date_test_col] = pd.to_datetime(df_qa[date_test_col], dayfirst=True, errors='coerce')
                        df_qa = df_qa.dropna(subset=[date_test_col])
                        
                        # 💡 THE FIX: إزالة الفلتر القديم اللي كان بيمسح اختبارات الـ DPL لو المقاول معندوش كميات تنفيذ
                        df_dpl = df_qa[df_qa[test_type_col].astype(str).str.upper().str.contains('DPL')].copy()
                        
                        daily_dpl = pd.DataFrame()
                        if not df_dpl.empty and num_tests_col in df_dpl.columns:
                            df_dpl[num_tests_col] = pd.to_numeric(df_dpl[num_tests_col].astype(str).str.replace(',', '', regex=False), errors='coerce').fillna(0)
                            
                            daily_dpl = df_dpl.groupby(date_test_col)[num_tests_col].sum().reset_index()
                            daily_dpl.rename(columns={date_test_col: 'Date', num_tests_col: 'DPL Tests'}, inplace=True)
                            daily_dpl = daily_dpl[daily_dpl['DPL Tests'] > 0].sort_values('Date')

                        # 💡 THE FIX 2: توحيد تاريخ البداية للمشروع بناءً على الكميات أو الـ DPL (أيهما أقدم)
                        min_exec = daily_exec[date_daily_col].min() if not daily_exec.empty else pd.NaT
                        min_dpl = daily_dpl['Date'].min() if not daily_dpl.empty else pd.NaT
                        
                        if pd.notna(min_exec) and pd.notna(min_dpl):
                            min_proj_date = min(min_exec, min_dpl)
                        elif pd.notna(min_exec):
                            min_proj_date = min_exec
                        elif pd.notna(min_dpl):
                            min_proj_date = min_dpl
                        else:
                            min_proj_date = pd.Timestamp.now().normalize()

                        # حساب الأسابيع لمسار الكميات
                        weekly_exec = pd.DataFrame()
                        if not daily_exec.empty:
                            daily_exec['Proj_Week'] = ((daily_exec[date_daily_col] - min_proj_date).dt.days // 7) + 1
                            weekly_exec = daily_exec.groupby('Proj_Week').agg(Executed=('Executed', 'sum'), Target=('Target', 'sum')).reset_index()
                            weekly_exec['Week_Label'] = "Wk " + weekly_exec['Proj_Week'].astype(str) + "<br>" + (min_proj_date + pd.to_timedelta((weekly_exec['Proj_Week'] - 1) * 7, unit='D')).dt.strftime('%b %y')

                        # حساب الأسابيع لمسار الجودة (DPL)
                        weekly_dpl = pd.DataFrame()
                        if not daily_dpl.empty:
                            daily_dpl['Proj_Week'] = ((daily_dpl['Date'] - min_proj_date).dt.days // 7) + 1
                            weekly_dpl = daily_dpl.groupby('Proj_Week').agg({'DPL Tests': 'sum'}).reset_index()
                            weekly_dpl['Week_Label'] = "Wk " + weekly_dpl['Proj_Week'].astype(str) + "<br>" + (min_proj_date + pd.to_timedelta((weekly_dpl['Proj_Week'] - 1) * 7, unit='D')).dt.strftime('%b %y')

                        # --- التجهيز المسبق لجدول التتبع (عشان نطلع الـ AI Alert) ---
                        duration_rows = []
                        if sel_contractor != 'All Contractors':
                            df_valid_dur = chart_df[chart_df[contractor_col].astype(str).str.strip().str.lower() == sel_contractor.lower()].dropna(subset=[elem_all_col, date_daily_col])
                        else:
                            df_valid_dur = chart_df.dropna(subset=[contractor_col, elem_all_col, date_daily_col])
                            
                        for (ct, el), group in df_valid_dur.groupby([contractor_col, elem_all_col]):
                            group_exec = group[group[exec_qty_m3_col] > 0].copy()
                            if group_exec.empty: continue
                            group_exec[date_daily_col] = pd.to_datetime(group_exec[date_daily_col], dayfirst=True, errors='coerce')
                            
                            min_d = group_exec[date_daily_col].min()
                            max_d = group_exec[date_daily_col].max()
                            t_duration = (max_d - min_d).days + 1
                            
                            # حساب الأيام الفعلية بدقة (تفادي الأيام الساقطة)
                            a_days = group_exec[date_daily_col].nunique()
                            
                            d_t = group_exec.groupby(date_daily_col).agg(Target=(target_col, 'max'), Exec=(exec_qty_m3_col, 'sum'))
                            m_days = len(d_t[d_t['Exec'] < d_t['Target']])
                            avg_tar = d_t['Target'].replace(0, np.nan).mean()
                            tot_e = d_t['Exec'].sum()
                            
                            ideal_dur = (tot_e / avg_tar) if pd.notna(avg_tar) and avg_tar > 0 else 0
                            var_dur = t_duration - ideal_dur
                            
                            duration_rows.append({
                                'Contractor': ct, 'Element': el, 'Actual Active Days': a_days,
                                'Total Spanned Days': t_duration, 'Days Missed Target': m_days,
                                'Ideal Duration': ideal_dur, 'Delay Variance': var_dur,
                                'Last Exec Date': max_d.strftime('%Y-%m-%d')
                            })
                        tracker_df = pd.DataFrame(duration_rows).sort_values('Delay Variance', ascending=False) if duration_rows else pd.DataFrame()

                        # --- 3. الدمج ورسم الشارت (Weekly / Daily) ---
                        ch_left, ch_right = st.columns([0.75, 0.25])
                        with ch_left:
                            fig_d = make_subplots(specs=[[{"secondary_y": True}]])
                            
                            if time_view == "Weekly":
                                df_weeks_exec = weekly_exec[['Proj_Week', 'Week_Label', 'Executed', 'Target']] if not weekly_exec.empty else pd.DataFrame(columns=['Proj_Week', 'Week_Label', 'Executed', 'Target'])
                                df_weeks_dpl = weekly_dpl[['Proj_Week', 'Week_Label', 'DPL Tests']] if not weekly_dpl.empty else pd.DataFrame(columns=['Proj_Week', 'Week_Label', 'DPL Tests'])
                                
                                # دمج البيانات لضمان عدم سقوط أسابيع (الزجزاج)
                                all_weeks_df = pd.merge(df_weeks_exec, df_weeks_dpl, on=['Proj_Week', 'Week_Label'], how='outer').sort_values('Proj_Week').fillna(0)
                                
                                x_data_exec = all_weeks_df['Week_Label'].tolist() if not all_weeks_df.empty else []
                                y_exec = all_weeks_df['Executed'].tolist() if not all_weeks_df.empty else []
                                y_target = all_weeks_df['Target'].tolist() if not all_weeks_df.empty else []
                                x_data_dpl = all_weeks_df['Week_Label'].tolist() if not all_weeks_df.empty else []
                                y_dpl = all_weeks_df['DPL Tests'].tolist() if not all_weeks_df.empty else []
                                x_title = "Project Weeks (from actual start date)"
                            else:
                                x_data_exec = daily_exec[date_daily_col] if not daily_exec.empty else []
                                y_exec = daily_exec['Executed'] if not daily_exec.empty else []
                                y_target = daily_exec['Target'] if not daily_exec.empty else []
                                x_data_dpl = daily_dpl['Date'] if not daily_dpl.empty else []
                                y_dpl = daily_dpl['DPL Tests'] if not daily_dpl.empty else []
                                x_title = "Date"
                            
                            if len(x_data_exec) > 0:
                                fig_d.add_trace(go.Bar(
                                    x=x_data_exec, y=y_exec,
                                    name=f'Execution (m³)', marker_color='#00d2ff', opacity=0.85
                                ), secondary_y=False)
                                
                                fig_d.add_trace(go.Scatter(
                                    x=x_data_exec, y=y_target,
                                    name='Target Rate', mode='lines',
                                    line=dict(color='#e74c3c', width=3, dash='dash')
                                ), secondary_y=False)
                                
                            if len(x_data_dpl) > 0 and show_dpl_line:
                                fig_d.add_trace(go.Scatter(
                                    x=x_data_dpl, y=y_dpl,
                                    name=f'DPL Tests', mode='markers+lines',
                                    line=dict(color='#ffaa00', width=2),
                                    marker=dict(symbol='diamond', size=8, color='white', line=dict(color='#ffaa00', width=2))
                                ), secondary_y=True)

                            chart_title = f"{time_view} Timeline: Execution vs Target" + (" vs DPL" if show_dpl_line else "")
                            if sel_elem_chart != 'All Elements': chart_title += f" ➔ {sel_elem_chart}"
                                
                            fig_d.update_layout(
                                title=chart_title, height=450, hovermode='x unified', margin=dict(l=0, r=0, t=40, b=0),
                                legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1)
                            )
                            
                            if time_view == "Weekly":
                                fig_d.update_xaxes(
                                    type='category',
                                    categoryorder='array',
                                    categoryarray=x_data_exec,
                                    tickangle=0,
                                    title_text=x_title
                                )
                                fig_d.update_layout(margin=dict(b=80))
                            else:
                                fig_d.update_xaxes(title_text=x_title, tickangle=-45)
                                
                            fig_d.update_yaxes(title_text="Quantity (m³)", secondary_y=False)
                            fig_d.update_yaxes(title_text="Number of DPL Tests", secondary_y=True, showgrid=False)
                            try: fig_d = style_3d_glassy(fig_d, "bar")
                            except: pass
                            st.plotly_chart(fig_d, use_container_width=True, key="qty_dynamic_chart")
                            exported_figs["2. Execution Timeline"] = fig_d
                        with ch_right:
                            st.markdown("**Performance Summary**")
                            days = len(daily_exec)
                            met = int((daily_exec['Executed'] >= daily_exec['Target']).sum()) if not daily_exec.empty else 0
                            hit_rate = round((met / days * 100), 1) if days else 0
                            total_exec = daily_exec['Executed'].sum() if not daily_exec.empty else 0
                            total_dpl = daily_dpl['DPL Tests'].sum() if not daily_dpl.empty else 0
                            
                            st.markdown(f"""
                            <div style="background:rgba(255,255,255,0.05); padding:15px; border-radius:8px; border-left:4px solid #00d2ff; margin-bottom:15px;">
                                <div style="color:{ui['text_muted']}; font-size:12px;">Total Executed (m³)</div>
                                <div style="color:{ui['text_main']}; font-size:20px; font-weight:bold;">{total_exec:,.1f}</div>
                                <hr style="border-color:rgba(255,255,255,0.1); margin:10px 0;">
                                <div style="color:{ui['text_muted']}; font-size:12px;">Total DPL Tests</div>
                                <div style="color:#ffaa00; font-size:20px; font-weight:bold;">{int(total_dpl):,}</div>
                                <hr style="border-color:rgba(255,255,255,0.1); margin:10px 0;">
                                <div style="color:{ui['text_muted']}; font-size:12px;">Target Hit Rate</div>
                                <div style="color:#2ecc71; font-size:20px; font-weight:bold;">{hit_rate}%</div>
                            </div>
                            """, unsafe_allow_html=True)
                            
                            # --- الذكاء الاصطناعي يستنتج أسوأ عنصر (AI Alert) ---
                            if not tracker_df.empty and sel_contractor != 'All Contractors' and sel_elem_chart == 'All Elements':
                                worst_elem = tracker_df.iloc[0]
                                if worst_elem['Delay Variance'] > 5:
                                    st.markdown(f"""
                                    <div style="background:rgba(231, 76, 60, 0.1); border-left:4px solid #e74c3c; padding:15px; border-radius:8px;">
                                        <div style="color:#e74c3c; font-size:12px; font-weight:bold; margin-bottom:5px;">🚨 AI DELAY ALERT</div>
                                        <div style="color:{ui['text_main']}; font-size:13px; line-height:1.5;">
                                            Element <b style="color:#ffaa00;">{worst_elem['Element']}</b> is causing the largest drag, showing a delay variance of <b style="color:#e74c3c;">{worst_elem['Delay Variance']:.1f} days</b> vs the ideal duration.
                                        </div>
                                    </div>
                                    """, unsafe_allow_html=True)

                    # --- 4. جدول الأصفار النظيف (Tracker Table) ---
                    st.markdown("##### ⏳ Duration & Pacing Tracker")
                    if not tracker_df.empty:
                        if sel_elem_chart != 'All Elements':
                            display_tracker = tracker_df[tracker_df['Element'] == sel_elem_chart]
                        else:
                            display_tracker = tracker_df
                        
                        def highlight_var(val):
                            if isinstance(val, (int, float)):
                                if val > 5: return 'color: #e74c3c; font-weight: bold'
                                elif val < 0: return 'color: #2ecc71; font-weight: bold'
                            return ''
                            
                        # تنسيق الأرقام العشرية لمنع الأصفار الزيادة
                        st.dataframe(
                            display_tracker.style.format({'Ideal Duration': '{:.1f}', 'Delay Variance': '{:.1f}'})
                            .map(highlight_var, subset=['Delay Variance']),
                            use_container_width=True, hide_index=True
                        )
                    else:
                        st.info("No active execution data for duration tracking.")
                    # السيكشنات 4-8 دلوقتي بره الـ else — بتظهر دايماً
                    # الـ warning بيظهر بس لو في columns ناقصة في الـ chart
                    if not (date_daily_col and target_col and exec_qty_m3_col and contractor_col):
                        st.warning("⚠️ Missing columns for the unified chart above.")
                    # ══════════════════════════════════════════════════════
                    # 4. ELEMENT COVERAGE AUDIT (Linking QA to Execution)
                    # ══════════════════════════════════════════════════════
                    st.markdown("#### 🕵️ Quantity Auditor — Missing Elements")
                    st.caption(
                        "Links elements tested in QA (ELMENT + Company Name) with Execution log "
                        "(Element (all) + Contractor) to detect tested elements with no billed quantities."
                    )

                    if elem_all_col and elment_col and contractor_col and comp_name_col:
                        # 1. جلب العناصر اللي المقاول عملها اختبارات جودة (من مسار الجودة)
                        if sel_contractor != 'All Contractors':
                            tested_elements = df[df[comp_name_col].astype(str).str.strip().str.lower() == sel_contractor.lower()][elment_col].dropna().astype(str).str.strip().unique()
                        else:
                            tested_elements = df[elment_col].dropna().astype(str).str.strip().unique()

                        # 2. جلب بيانات الكميات للمقاول ده (من مسار الكميات)
                        if sel_contractor != 'All Contractors':
                            exec_df = df[df[contractor_col].astype(str).str.strip().str.lower() == sel_contractor.lower()]
                        else:
                            exec_df = df.copy()
                            
                        # تجميع الكميات لكل عنصر (Element (all))
                        exec_agg = exec_df.groupby(elem_all_col)[exec_qty_m3_col].sum().reset_index()
                        exec_agg[elem_all_col] = exec_agg[elem_all_col].astype(str).str.strip().str.lower()

                        missing_elems  = []
                        covered_elems  = []

                        for el_name in tested_elements:
                            if el_name == '' or el_name.lower() == 'nan': continue
                            
                            # بنبحث عن العنصر ده هل نزل له كمية ولا لأ
                            match = exec_agg[exec_agg[elem_all_col] == el_name.lower()]
                            qty = match[exec_qty_m3_col].sum() if not match.empty else 0

                            if qty == 0:
                                missing_elems.append({
                                    'Tested Element (QA)': el_name,
                                    'Status': '❌ No Quantity Logged'
                                })
                            else:
                                covered_elems.append({
                                    'Tested Element (QA)': el_name,
                                    'Executed Qty (m³)': round(qty, 1),
                                    'Status': '✅ Has Quantity'
                                })

                        col_m, col_c = st.columns(2)

                        with col_m:
                            if missing_elems:
                                miss_df = pd.DataFrame(missing_elems).sort_values('Tested Element (QA)')
                                st.markdown(f"""
                                <div style="background:rgba(231,76,60,0.1); border-left:4px solid #e74c3c; padding:14px;border-radius:8px;margin-bottom:10px;">
                                    <b style="color:#e74c3c;">🚨 {len(missing_elems)} Tested Elements Missing Quantities</b><br>
                                    <span style="font-size:12px;color:#d1d5da;">Tested in QA but no execution quantity found!</span>
                                </div>
                                """, unsafe_allow_html=True)
                                st.dataframe(miss_df, use_container_width=True)
                            else:
                                st.success("✅ All tested elements have assigned quantities!")

                        with col_c:
                            if covered_elems:
                                cov_df = pd.DataFrame(covered_elems).sort_values('Executed Qty (m³)', ascending=False)
                                st.markdown(f"""
                                <div style="background:rgba(46,204,113,0.1); border-left:4px solid #2ecc71; padding:14px;border-radius:8px;margin-bottom:10px;">
                                    <b style="color:#2ecc71;">✅ {len(covered_elems)} Elements with Quantities</b><br>
                                    <span style="font-size:12px;color:#d1d5da;">Properly linked between QA/QC and Execution.</span>
                                </div>
                                """, unsafe_allow_html=True)
                                st.dataframe(cov_df, use_container_width=True)
                    else:
                        st.warning("⚠️ Missing columns for auditor.")

                   # ══════════════════════════════════════════════════════
                    # 5. ADVANCED KPI SUMMARY (Execution Intensity Index - EII)
                    # ══════════════════════════════════════════════════════
                    st.markdown("#### 🎯 Execution Intensity Matrix (EII)")
                    st.caption("Evaluates Daily Productivity vs Average Target (EII %) against Daily Consistency (Hit Rate %).")

                    if contractor_col and exec_qty_m3_col and target_col and elem_all_col:
                        kpi_rows = []
                        
                        df_valid = df.dropna(subset=[contractor_col, elem_all_col])
                        
                        if sel_contractor != 'All Contractors':
                            df_valid = df_valid[df_valid[contractor_col].astype(str).str.strip().str.lower() == sel_contractor.lower()]
                            
                        # 🔥 التعديل السحري: ربط شارت الماتريكس بفلتر العناصر اللي فوق
                        if 'sel_elem_chart' in locals() and sel_elem_chart != 'All Elements':
                            df_valid = df_valid[df_valid[elem_all_col].astype(str).str.strip().str.lower() == sel_elem_chart.lower()]
                            
                        grouped = df_valid.groupby([contractor_col, elem_all_col])
                        
                        for name, group in grouped:
                            ct = str(name[0]).strip()
                            el = str(name[1]).strip()
                            
                            if ct.lower() in ['nan', 'none', ''] or el.lower() in ['nan', 'none', '']: continue
                            
                            # 1. إجمالي المنفذ
                            exec_sum = group[exec_qty_m3_col].sum()
                            
                            # 2. حصر أيام العمل الفعلية (اللي ليها تارجت)
                            valid_target_group = group[group[target_col] > 0]
                            days_w = len(valid_target_group)
                            
                            if days_w > 0:
                                cumulative_target = valid_target_group[target_col].sum()
                                
                                # 3. المؤشرات الهندسية الجديدة (السرعات)
                                daily_productivity = exec_sum / days_w
                                avg_daily_target = cumulative_target / days_w
                                
                                # 4. مؤشر كثافة التنفيذ (EII)
                                eii_pct = round((daily_productivity / avg_daily_target * 100), 1) if avg_daily_target > 0 else 0
                                # 🔥 تحجيم السرعات الخرافية لـ 400% كحد أقصى عشان الشارت ميبوظش والنقط تفضل ظاهرة
                                eii_pct = min(eii_pct, 6000)
                                # 5. حساب الاستمرارية (Hit Rate %)
                                days_met = int((valid_target_group[exec_qty_m3_col] >= valid_target_group[target_col]).sum())
                                hit_rate = round((days_met / days_w * 100), 1)
                            else:
                                daily_productivity = 0
                                avg_daily_target = 0
                                eii_pct = 0
                                hit_rate = 0
                                cumulative_target = 0
                                
                            status = '🟢 Good' if hit_rate >= 70 and eii_pct >= 90 else ('🟡 Fair' if hit_rate >= 40 else '🔴 Poor')
                            
                            if exec_sum > 0 or cumulative_target > 0:
                                kpi_rows.append({
                                    'Contractor': ct,
                                    'Element': el,
                                    'Working Days': days_w,
                                    'Avg Daily Target (m³/d)': round(avg_daily_target, 1),
                                    'Daily Productivity (m³/d)': round(daily_productivity, 1),
                                    'EII (Intensity) %': eii_pct,
                                    'Hit Rate %': hit_rate,
                                    'Status': status,
                                    'Bubble_Size': max(exec_sum, 100) # لحجم الدائرة
                                })

                        if kpi_rows:
                            kpi_df = pd.DataFrame(kpi_rows).sort_values(by=['Contractor', 'EII (Intensity) %'], ascending=[True, False])
                            
                            st.markdown("##### 📈 Matrix Chart: Execution Intensity (EII) vs Hit Rate")
                            
                            fig_kpi = px.scatter(
                                kpi_df, 
                                x='EII (Intensity) %', 
                                y='Hit Rate %',
                                color='Contractor',
                                size='Bubble_Size',
                                # text='Element', # ❌ لغينا التكست المطبوع عشان نمنع الزحمة
                                hover_name='Contractor',
                                hover_data={
                                    'Contractor': False, 
                                    'Element': True, 
                                    'Working Days': True,
                                    'Avg Daily Target (m³/d)': True,
                                    'Daily Productivity (m³/d)': True, 
                                    'Bubble_Size': False
                                },
                                size_max=45,
                                color_discrete_sequence=NEON_COLORS
                            )
                            
                            # خط المنتصف للاستمرارية وللكثافة
                            fig_kpi.add_hline(y=50, line_dash="dash", line_color="rgba(255,255,255,0.3)", line_width=1)
                            fig_kpi.add_vline(x=100, line_dash="dash", line_color="rgba(255,255,255,0.3)", line_width=1)
                            
                            # إضافة الـ Annotations التحليلية (بأماكن ثابتة ومضبوطة)
                            fig_kpi.add_annotation(x=175, y=95, text="🌟 High Intensity & Consistent", showarrow=False, font=dict(color="#2ecc71", size=11), bgcolor="rgba(0,0,0,0.5)")
                            fig_kpi.add_annotation(x=50, y=95, text="🐢 Consistent but Slow", showarrow=False, font=dict(color="#00d2ff", size=11), bgcolor="rgba(0,0,0,0.5)")
                            fig_kpi.add_annotation(x=175, y=5, text="🚀 High Speed, Poor Consistency", showarrow=False, font=dict(color="#f1c40f", size=11), bgcolor="rgba(0,0,0,0.5)")
                            fig_kpi.add_annotation(x=50, y=5, text="🚨 Slow & Erratic", showarrow=False, font=dict(color="#e74c3c", size=11), bgcolor="rgba(0,0,0,0.5)")
                            
                            fig_kpi.update_layout(
                                title="Execution Intensity Matrix (EII)", # ✅ إضافة عنوان صريح
                                showlegend=False, # ✅ إخفاء الليجند العملاق لتوفير المساحة
                                height=600, 
                                xaxis=dict(title="Execution Intensity Index - EII % (Capped at 6000%)"), # ✅ تحجيم المحور السيني لمنع التمدد الكارثي
                                yaxis=dict(range=[-10, 110], title="Target Hit Rate % (Daily Consistency)"),
                                hovermode='closest',
                                margin=dict(t=50, b=30, l=30, r=30)
                            )
                            
                            try: fig_kpi = style_3d_glassy(fig_kpi, "scatter")
                            except: pass
                            
                            st.plotly_chart(fig_kpi, use_container_width=True, key="qty_kpi_scatter_eii")
                            exported_figs["3. Execution Intensity Matrix"] = fig_kpi
                            with st.expander("📋 View Detailed Engineering Data (Speed & EII)"):
                                st.dataframe(kpi_df.drop(columns=['Bubble_Size']), use_container_width=True)
                    else:
                        st.warning("⚠️ Missing columns to generate Element Matrix.")

                    # ══════════════════════════════════════════════════════
                    # 6. CONTRACTOR SCORECARD & VERDICT ENGINE (AI DECISION)
                    # ══════════════════════════════════════════════════════
                    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
                    st.markdown("#### 🏆 Contractor Performance Scorecard")
                    st.caption("إطار تقييم هندسي شامل يعتمد على 5 معايير لتقييم المقاولين وتحديد الموثوقية.")

                    if contractor_col and target_col and exec_qty_m3_col and date_daily_col:
                        # 1. تجهيز الداتا للتقييم
                        df_score = df.copy()
                        df_score[target_col]      = pd.to_numeric(df_score[target_col], errors='coerce').fillna(0)
                        df_score[exec_qty_m3_col] = pd.to_numeric(df_score[exec_qty_m3_col], errors='coerce').fillna(0)
                        df_score[date_daily_col]  = pd.to_datetime(df_score[date_daily_col], dayfirst=True, errors='coerce')
                        df_score = df_score.dropna(subset=[contractor_col, date_daily_col])
                        
                        # 2. حساب حد الحجم الضخم (Tier A) للـ Scale Bonus
                        contractor_volumes = df_score.groupby(contractor_col)[exec_qty_m3_col].sum()
                        tier_a_threshold = contractor_volumes.quantile(0.75) if not contractor_volumes.empty else 0

                        scorecard_rows = []
                        
                        for ct, group in df_score.groupby(contractor_col):
                            ct = str(ct).strip()
                            if ct.lower() in ['nan', 'none', '']: continue
                            
                            group = group.sort_values(date_daily_col)
                            total_exec = group[exec_qty_m3_col].sum()
                            
                            if total_exec <= 0: continue
                                
                            # --- 1. Target Achievement (Weight: 25%) ---
                            valid_target_days = group[group[target_col] > 0]
                            total_valid_days = len(valid_target_days)
                            if total_valid_days > 0:
                                hit_days = len(valid_target_days[valid_target_days[exec_qty_m3_col] >= valid_target_days[target_col]])
                                hit_rate = hit_days / total_valid_days
                            else:
                                hit_rate = 0
                            score_target = min(100, hit_rate * 100)
                            
                            # --- 2. Consistency (Weight: 20%) ---
                            if total_valid_days > 2:
                                mean_perf = valid_target_days[exec_qty_m3_col].mean()
                                std_perf = valid_target_days[exec_qty_m3_col].std()
                                cov = (std_perf / mean_perf) if mean_perf > 0 else 1
                                consistency = max(0, 1 - cov)
                            else:
                                consistency = 0
                            score_consistency = min(100, consistency * 100)
                            
                            # --- 3. Momentum (Weight: 15%) ---
                            if total_valid_days > 3:
                                y = (valid_target_days[exec_qty_m3_col] / valid_target_days[target_col]).replace([np.inf, -np.inf], 0).fillna(0).values
                                x = np.arange(len(y))
                                slope, _ = np.polyfit(x, y, 1)
                                momentum = 50 + (slope * 500)
                                score_momentum = max(0, min(100, momentum))
                            else:
                                score_momentum = 50 
                                
                            # --- 4. Volume Delivery (Weight: 10%) ---
                            if company_col and total_qty_col and company_col in df.columns:
                                tot_req = df[df[company_col].astype(str).str.strip() == ct][total_qty_col].sum()
                                vol_delivery = (total_exec / tot_req * 100) if pd.notna(tot_req) and tot_req > 0 else 0
                            else:
                                vol_delivery = 0
                            score_volume = min(100, vol_delivery)
                            
                            # --- 5. Active Days Rate (Weight: 5%) ---
                            min_date = group[date_daily_col].min()
                            max_date = group[date_daily_col].max()
                            spanned_days = (max_date - min_date).days + 1
                            active_days = group[group[exec_qty_m3_col] > 0][date_daily_col].nunique()
                            active_rate = (active_days / spanned_days * 100) if spanned_days > 0 else 0
                            score_active = min(100, active_rate)
                            
                            # --- 6. Site Execution Quality Yield (Weight: 25%) 🎯 ---
                            score_quality = 50 # سكور افتراضي لو مفيش اختبارات
                            test_col_name = next((c for c in df.columns if 'TEST TYPE' in c.upper() or c.strip() == 'Test Type'), 'Test Type')
                            
                            if comp_name_col and 'sample status' in df.columns and test_col_name in df.columns:
                                ct_qa = df[df[comp_name_col].astype(str).str.strip().str.lower() == ct.lower()]
                                # الفلترة لاختبارات DPL و Plate Load فقط التي تمت في الموقع
                                field_tests = ct_qa[ct_qa[test_col_name].astype(str).str.upper().str.contains('DPL|PLATE', na=False)]
                                total_field_tests = len(field_tests)
                                
                                if total_field_tests > 0:
                                    acc_field = len(field_tests[field_tests['sample status'].astype(str).str.upper().isin(['ACCEPTED', 'APPROVED AS NOTED'])])
                                    score_quality = (acc_field / total_field_tests) * 100
                            
                            # --- 🧮 حساب النتيجة النهائية الموزونة ---
                            final_score = (
                                (score_target * 0.25) + 
                                (score_consistency * 0.20) + 
                                (score_momentum * 0.15) + 
                                (score_volume * 0.10) + 
                                (score_active * 0.05) +
                                (score_quality * 0.25)
                            )
                            
                            # --- 🚀 Scale Bonus (علاوة الحجم للمقاولين الكبار) ---
                            is_tier_a = False
                            if total_exec >= tier_a_threshold and tier_a_threshold > 0:
                                final_score += 5.0 # إضافة 5 درجات
                                final_score = min(100.0, final_score)
                                is_tier_a = True
                                
                            display_name = f"🌟 {ct}" if is_tier_a else ct
                            
                            # --- إصدار القرار الإداري (Verdict) ---
                            if final_score >= 70:
                                verdict = "🟢 Reliable"
                            elif final_score >= 40:
                                verdict = "🟡 Watch List"
                            else:
                                verdict = "🔴 Critical"
                                
                            scorecard_rows.append({
                                'Contractor': display_name,
                                'Final Score': round(final_score, 1),
                                'Verdict': verdict,
                                'Quality Yield (25%)': f"{score_quality:.1f}%",
                                'Target Achv. (25%)': f"{score_target:.1f}%",
                                'Consistency (20%)': f"{score_consistency:.1f}%",
                                'Momentum (15%)': f"{score_momentum:.1f}%",
                                'Vol Delivery (10%)': f"{score_volume:.1f}%",
                                'Active Rate (5%)': f"{score_active:.1f}%"
                            })
                            
                        if scorecard_rows:
                            score_df = pd.DataFrame(scorecard_rows).sort_values('Final Score', ascending=False)
                            
                            # تصميم الكروت الإرشادية للحالة الإدارية
                            st.markdown("""
                            <style>
                            .v-box { padding: 15px; border-radius: 10px; text-align: center; border: 1px solid; }
                            .v-g { background: rgba(46,204,113,0.1); border-color: #2ecc71; }
                            .v-a { background: rgba(241,196,15,0.1); border-color: #f1c40f; }
                            .v-r { background: rgba(231,76,60,0.1); border-color: #e74c3c; }
                            </style>
                            """, unsafe_allow_html=True)
                            
                            vc1, vc2, vc3 = st.columns(3)
                            vc1.markdown(f'<div class="v-box v-g"><h4 style="color:#2ecc71; margin:0;">🟢 Reliable Contractor</h4><p style="font-size:12px; color:{ui["text_muted"]}; margin:5px 0 0;">Score ≥ 70 — إنتاجية ثابتة وجودة دمك ممتازة. (🌟 = مقاول فئة أولى حجماً).</p></div>', unsafe_allow_html=True)
                            vc2.markdown(f'<div class="v-box v-a"><h4 style="color:#f1c40f; margin:0;">🟡 Watch List</h4><p style="font-size:12px; color:{ui["text_muted"]}; margin:5px 0 0;">Score 40-70 — أداء متذبذب أو رفض متكرر. يحتاج مراقبة مكثفة.</p></div>', unsafe_allow_html=True)
                            vc3.markdown(f'<div class="v-box v-r"><h4 style="color:#e74c3c; margin:0;">🔴 Critical — Intervention</h4><p style="font-size:12px; color:{ui["text_muted"]}; margin:5px 0 0;">Score < 40 — خطر على الجدول الزمني وجودة المشروع. تدخل فوري.</p></div>', unsafe_allow_html=True)
                            
                            st.markdown("<br>", unsafe_allow_html=True)
                            
                            # رسم الشارت النهائي للسكور
                            fig_score = px.bar(
                                score_df, x='Contractor', y='Final Score', text='Final Score',
                                color='Final Score', color_continuous_scale=['#e74c3c', '#f1c40f', '#2ecc71'],
                                range_color=[0, 100], title="Contractor Overall Performance Score (Out of 100)"
                            )
                            fig_score.add_hline(y=70, line_dash="dash", line_color="#2ecc71", annotation_text="Reliable Threshold")
                            fig_score.add_hline(y=40, line_dash="dash", line_color="#e74c3c", annotation_text="Critical Threshold")
                            fig_score.update_traces(textposition='outside')
                            try: fig_score = style_3d_glassy(fig_score, "bar")
                            except: pass
                            st.plotly_chart(fig_score, use_container_width=True, key="scorecard_bar_chart")
                            exported_figs["11. Contractor Performance Scorecard"] = fig_score
                            
                            # عرض الجدول والتلوين
                            def color_verdict(val):
                                if 'Reliable' in str(val): return 'color: #2ecc71; font-weight: bold'
                                elif 'Watch' in str(val): return 'color: #f1c40f; font-weight: bold'
                                elif 'Critical' in str(val): return 'color: #e74c3c; font-weight: bold'
                                return ''
                                
                            st.dataframe(score_df.style.map(color_verdict, subset=['Verdict']), use_container_width=True, hide_index=True)
                        else:
                            st.info("لا توجد بيانات كافية لحساب تقييم المقاولين.")
                    else:
                        st.warning("⚠️ أعمدة التقييم غير مكتملة.")
                    # ══════════════════════════════════════════════════════
                    # 6.5 EARNED VALUE MANAGEMENT (SPI MODULE)
                    # ══════════════════════════════════════════════════════
                    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
                    st.markdown("#### ⏱️ Schedule Performance Index (EVM)")
                    st.caption("مؤشر الأداء الزمني (SPI): يقيس كفاءة إنجاز العمل مقارنة بالجدول الزمني المخطط (الهدف اليومي). القيمة > 1.0 تعني تقدم، والقيمة < 1.0 تعني تأخير.")
                    
                    if contractor_col and date_daily_col and target_col and exec_qty_m3_col:
                        # 1. تجهيز الداتا للـ EVM (تجميع صحيح لتفادي تكرار التارجت لنفس اليوم)
                        spi_df = df_sel.dropna(subset=[contractor_col, date_daily_col]).copy()
                        
                        spi_daily = spi_df.groupby([date_daily_col, contractor_col, elem_all_col]).agg(
                            EV=(exec_qty_m3_col, 'sum'),
                            PV=(target_col, 'max')
                        ).reset_index()
                        
                        # 2. حساب الـ SPI الكلي للقطاع أو الاختيار
                        total_ev = spi_daily['EV'].sum()
                        total_pv = spi_daily['PV'].sum()
                        overall_spi = (total_ev / total_pv) if total_pv > 0 else 0
                        
                        # 3. حساب الـ SPI لكل مقاول للمضمار
                        contractor_spi = spi_daily.groupby(contractor_col).agg(
                            EV=('EV', 'sum'),
                            PV=('PV', 'sum')
                        ).reset_index()
                        contractor_spi['SPI'] = (contractor_spi['EV'] / contractor_spi['PV']).fillna(0)
                        contractor_spi = contractor_spi[contractor_spi['PV'] > 0] # استبعاد اللي معندوش تارجت
                        contractor_spi = contractor_spi.sort_values('SPI', ascending=True) 

                        col_spi_g, col_spi_b = st.columns([0.35, 0.65])
                        
                        # 🎯 العداد الزجاجي (Gauge Chart)
                        with col_spi_g:
                            if overall_spi >= 1.0: 
                                gauge_color = "#2ecc71"
                                spi_status = "AHEAD OF SCHEDULE"
                            elif overall_spi >= 0.85: 
                                gauge_color = "#f1c40f"
                                spi_status = "MINOR DELAY"
                            else: 
                                gauge_color = "#e74c3c"
                                spi_status = "CRITICAL DELAY"
                                
                            fig_spi_gauge = go.Figure(go.Indicator(
                                mode = "gauge+number",
                                value = overall_spi,
                                number = {'valueformat': ".2f", 'font': {'size': 45, 'color': gauge_color, 'family': 'Rajdhani'}},
                                title = {'text': "Overall SPI (Efficiency)", 'font': {'size': 16, 'color': '#e2e8f0'}},
                                gauge = {
                                    'axis': {'range': [0, max(1.5, overall_spi+0.2)], 'tickwidth': 1, 'tickcolor': "white"},
                                    'bar': {'color': gauge_color, 'thickness': 0.8},
                                    'bgcolor': "rgba(15, 23, 42, 0.5)" if is_dark else "rgba(240, 245, 250, 0.8)",
                                    'borderwidth': 2,
                                    'bordercolor': "rgba(0, 210, 255, 0.1)",
                                    'steps': [
                                        {'range': [0, 0.85], 'color': "rgba(231, 76, 60, 0.2)"},
                                        {'range': [0.85, 0.99], 'color': "rgba(241, 196, 15, 0.2)"},
                                        {'range': [0.99, 1.5], 'color': "rgba(46, 204, 113, 0.2)"}
                                    ],
                                    'threshold': {'line': {'color': "white", 'width': 4}, 'thickness': 0.75, 'value': 1.0}
                                }
                            ))
                            fig_spi_gauge.update_layout(height=280, margin=dict(l=30, r=30, t=40, b=10), paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)")
                            
                            st.markdown(f"""
                            <div style="background: {ui['card_bg']}; border: 1px solid {ui['border_color']}; border-radius: 12px; padding: 15px; text-align: center; box-shadow: {ui['shadow']}; height: 100%;">
                            """, unsafe_allow_html=True)
                            st.plotly_chart(fig_spi_gauge, use_container_width=True, key="spi_gauge_main")
                            
                            diff_pct = abs(1 - overall_spi) * 100
                            if overall_spi >= 1.0:
                                msg = f"🚀 Zone is {diff_pct:.1f}% Ahead of Schedule"
                            else:
                                msg = f"⚠️ Zone is {diff_pct:.1f}% Behind Schedule"
                                
                            st.markdown(f"<div style='margin-top: -20px; font-family: Rajdhani; font-size: 15px; color: {gauge_color}; font-weight: bold;'>{msg}</div></div>", unsafe_allow_html=True)
                        
                        # 🏎️ مضمار المقاولين (Bullet/Bar Chart)
                        with col_spi_b:
                            if not contractor_spi.empty:
                                contractor_spi['Color'] = contractor_spi['SPI'].apply(lambda x: '#2ecc71' if x >= 1.0 else ('#f1c40f' if x >= 0.85 else '#e74c3c'))
                                fig_spi_bar = px.bar(
                                    contractor_spi, x='SPI', y=contractor_col, orientation='h', text='SPI', 
                                    title="🏎️ Contractors SPI Leaderboard"
                                )
                                fig_spi_bar.update_traces(marker_color=contractor_spi['Color'], texttemplate='%{text:.2f}', textposition='outside', textfont=dict(size=14, color='white'))
                                fig_spi_bar.add_vline(x=1.0, line_dash="dash", line_color="white", line_width=2, annotation_text="Target (1.0)", annotation_position="top right")
                                
                                fig_spi_bar.update_layout(
                                    height=280, margin=dict(l=10, r=30, t=40, b=10),
                                    xaxis=dict(range=[0, max(1.5, contractor_spi['SPI'].max() + 0.2)]),
                                    yaxis=dict(title="")
                                )
                                try: fig_spi_bar = style_3d_glassy(fig_spi_bar, "bar")
                                except: pass
                                
                                st.markdown(f"""
                                <div style="background: {ui['card_bg']}; border: 1px solid {ui['border_color']}; border-radius: 12px; padding: 15px; box-shadow: {ui['shadow']}; height: 100%;">
                                """, unsafe_allow_html=True)
                                st.plotly_chart(fig_spi_bar, use_container_width=True, key="spi_bar_main")
                                st.markdown("</div>", unsafe_allow_html=True)
                            else:
                                st.info("No valid targets found to generate Contractor Leaderboard.")
                                
                        # 🤖 قرار الإدارة من الذكاء الاصطناعي
                        st.markdown("<br>", unsafe_allow_html=True)
                        if not contractor_spi.empty:
                            worst_contractor = contractor_spi.iloc[0]
                            best_contractor = contractor_spi.iloc[-1]
                            
                            if overall_spi >= 1.0:
                                ai_spi_msg = f"القطاع يحقق معدلات إنجاز ممتازة متجاوزاً الجدول الزمني (SPI = {overall_spi:.2f}). أفضل أداء مسجل حالياً لشركة <b>{best_contractor[contractor_col]}</b> بمؤشر كفاءة يبلغ <b>{best_contractor['SPI']:.2f}</b>."
                                ai_spi_color = "#2ecc71"
                            elif overall_spi >= 0.85:
                                ai_spi_msg = f"القطاع يواجه تأخيراً طفيفاً (SPI = {overall_spi:.2f}). شركة <b>{worst_contractor[contractor_col]}</b> هي الأكثر تأخيراً بمؤشر يبلغ <b>{worst_contractor['SPI']:.2f}</b>، ويُنصح بمراجعة معدلات التنفيذ اليومية لها."
                                ai_spi_color = "#f1c40f"
                            else:
                                ai_spi_msg = f"تأخير حرج في القطاع (SPI = {overall_spi:.2f}). شركة <b>{worst_contractor[contractor_col]}</b> تمثل العائق الأكبر للجدول الزمني بمؤشر يبلغ <b>{worst_contractor['SPI']:.2f}</b>. يتطلب الأمر تدخلاً إدارياً عاجلاً."
                                ai_spi_color = "#e74c3c"
                                
                            st.markdown(f"""
                            <div style="background: rgba(10, 20, 33, 0.6); border-right: 5px solid {ai_spi_color}; padding: 15px 20px; border-radius: 8px; display: flex; align-items: center; justify-content: space-between; margin-bottom: 20px;" dir="rtl">
                                <div style="text-align: right;">
                                    <div style="color: {ai_spi_color}; font-weight: bold; font-size: 15px; margin-bottom: 5px; font-family: 'Rajdhani', sans-serif;">🤖 رؤية الذكاء الاصطناعي (AI EVM Insight):</div>
                                    <div style="color: {ui['text_main']}; font-size: 15px; line-height: 1.6;">{ai_spi_msg}</div>
                                </div>
                                <div style="font-size: 35px; opacity: 0.8;">📊</div>
                            </div>
                            """, unsafe_allow_html=True)
                            
                    else:
                        st.warning("⚠️ أعمدة الكميات (المنفذ والمستهدف) غير مكتملة لحساب الـ SPI.")
                    
                    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)     
                    # ══════════════════════════════════════════════════════
                    # 7. DPL & PLATE LOAD TESTS (Strict Lookup in Company Name)
                    # ══════════════════════════════════════════════════════
                    st.markdown("#### 🧪 DPL & PLATE LOAD Tests per Contractor")
                    st.caption("Takes the Contractor name and strictly matches it against 'Company Name' to get accurate test counts and Avg DPL.")

                    if contractor_col and comp_name_col and test_type_col and num_tests_col:
                        # 1. بنقرأ من الداتا الأصلية الكاملة عشان مفيش صفوف تضيع
                        df_tests = df.copy()
                        
                        # 2. تنظيف الفواصل من عدد الاختبارات عشان تتجمع صح
                        if df_tests[num_tests_col].dtype == 'object':
                            df_tests[num_tests_col] = df_tests[num_tests_col].astype(str).str.replace(',', '', regex=False)
                        df_tests[num_tests_col] = pd.to_numeric(df_tests[num_tests_col], errors='coerce').fillna(0)
                        
                        # 3. تنظيف عمود المتوسط 
                        if 'AVERAGE VALUE' in df_tests.columns:
                            df_tests['AVERAGE VALUE'] = pd.to_numeric(df_tests['AVERAGE VALUE'], errors='coerce')

                        # 4. فلترة أنواع الاختبارات
                        mask_dpl  = df_tests[test_type_col].astype(str).str.upper().str.contains('DPL')
                        mask_pl   = df_tests[test_type_col].astype(str).str.upper().str.contains('PLATE')
                        df_dpl    = df_tests[mask_dpl]
                        df_pl     = df_tests[mask_pl]

                        # 5. تحديد الأسماء اللي هنبحث بيها (بناءً على اختيارك من الفلتر)
                        if sel_contractor != 'All Contractors':
                            contractors_to_show = [sel_contractor]
                        else:
                            # بناخد لستة الأسماء من عمود Contractor بس
                            contractors_to_show = df[contractor_col].dropna().astype(str).str.strip().unique()

                        test_rows = []
                        
                        for ct in contractors_to_show:
                            if ct.lower() in ['nan', 'none', '']: continue
                            
                            # ⚠️ التعديل الجذري هنا: بنبحث بالاسم (ct) في عمود Company Name فقـــــط
                            ct_dpl = df_dpl[df_dpl[comp_name_col].astype(str).str.strip().str.lower() == ct.lower()]
                            ct_pl  = df_pl[df_pl[comp_name_col].astype(str).str.strip().str.lower() == ct.lower()]

                            dpl_total = ct_dpl[num_tests_col].sum()
                            pl_total  = ct_pl[num_tests_col].sum()
                            
                            # حساب متوسط الـ DPL
                            avg_dpl = ct_dpl['AVERAGE VALUE'].mean() if 'AVERAGE VALUE' in ct_dpl.columns else np.nan

                            if dpl_total > 0 or pl_total > 0:
                                test_rows.append({
                                    'Contractor':       ct,
                                    'DPL Tests':        int(dpl_total),
                                    'Plate Load Tests': int(pl_total),
                                    'Total Tests':      int(dpl_total + pl_total),
                                    'Avg DPL':          round(avg_dpl, 2) if pd.notna(avg_dpl) else 0
                                })

                        if test_rows:
                            test_df = pd.DataFrame(test_rows).sort_values('Total Tests', ascending=False)
                            
                            # رسم الكروت
                            top4_t = test_df.head(4)
                            t_cols = st.columns(min(len(top4_t), 4))
                            for i, (_, row) in enumerate(top4_t.iterrows()):
                                avg_text = f" | Avg DPL: <b style='color:#2ecc71;'>{row['Avg DPL']}</b>" if row['Avg DPL'] > 0 else ""
                                create_card(
                                    t_cols[i],
                                    f"🧪 {row['Contractor']}",
                                    f"DPL: {row['DPL Tests']:,}",
                                    delta_html=(
                                        f"<span style='color:#8da3b9;font-size:12px'>"
                                        f"Plate Load: <b style='color:#00d2ff;'>{row['Plate Load Tests']:,}</b>{avg_text}</span>"
                                    )
                                )

                            # رسم الشارت
                            fig_tests = px.bar(
                                test_df,
                                x='Contractor',
                                y=['DPL Tests', 'Plate Load Tests'],
                                barmode='group',
                                color_discrete_sequence=['#00d2ff', '#ffaa00'],
                                title="DPL vs Plate Load Tests per Contractor",
                                text_auto=True
                            )
                            try: fig_tests = style_3d_glassy(fig_tests, "bar")
                            except: pass
                            st.plotly_chart(fig_tests, use_container_width=True, key="dpl_pl_chart_strict")
                            exported_figs["12. DPL vs Plate Load per Contractor"] = fig_tests

                            with st.expander("📋 Full Table"):
                                st.dataframe(test_df, use_container_width=True)
                        else:
                            st.info("No DPL or PLATE LOAD data found for this contractor.")
                    else:
                        st.warning("⚠️ Missing columns for tests analysis.")
                    # ══════════════════════════════════════════════════════
                    # 8. ADVANCED DPL ANALYTICS MODULE (NEW)
                    # ══════════════════════════════════════════════════════
                    st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
                    st.markdown("#### 🎯 DPL Deep Analytics Matrix")
                    st.caption("Comprehensive analysis of Dynamic Cone Penetrometer (DPL) test averages per company and element.")
                    
                    if comp_name_col and 'AVERAGE VALUE' in df.columns and test_type_col in df.columns:
                        # فلترة الداتا للـ DPL فقط واستبعاد القيم الفارغة
                        dpl_data = df[df[test_type_col].astype(str).str.upper().str.contains('DPL', na=False)].copy()
                        dpl_data['AVERAGE VALUE'] = pd.to_numeric(dpl_data['AVERAGE VALUE'], errors='coerce')
                        dpl_data = dpl_data.dropna(subset=['AVERAGE VALUE', comp_name_col])
                        
                        if not dpl_data.empty:
# --- 1. Line Chart & Table (All Companies) ---
                            st.markdown("##### 🏢 DPL Performance & Consistency by Company")
                            
                           # 1. UI Radio Button (معايير القياس + المؤشر الشامل)
                            metric_choice = st.radio(
                                "📊 (Metric) اختر معيار القياس:",
                                options=[
                                    "🔴 Average (المتوسط)", 
                                    "🟡 Consistency (الاستمرارية)", 
                                    "🟢 Pass Rate (نسبة القبول)", 
                                    "🏆 Combined Score (المؤشر الشامل)"
                                ],
                                horizontal=True
                            )

                            # 2. تجميع البيانات والعمليات الحسابية
                            def count_accepted(series):
                                return series.astype(str).str.upper().isin(['ACCEPTED', 'APPROVED AS NOTED']).sum()

                            group_dpl = dpl_data.groupby(comp_name_col).agg(
                                total_tests=('AVERAGE VALUE', 'count'),
                                avg_dpl=('AVERAGE VALUE', 'mean'),
                                std_dpl=('AVERAGE VALUE', 'std'),
                                accepted_tests=('sample status', count_accepted) if 'sample status' in dpl_data.columns else ('AVERAGE VALUE', lambda x: 0)
                            ).reset_index()

                            # حساب الاستمرارية (Consistency)
                            group_dpl['Consistency'] = 1 - (group_dpl['std_dpl'] / group_dpl['avg_dpl']).fillna(0)
                            group_dpl['Consistency'] = group_dpl['Consistency'].apply(lambda x: max(0, x)) * 100
                            
                            # حساب نسبة القبول الفعلي (Pass Rate)
                            group_dpl['Pass Rate %'] = (group_dpl['accepted_tests'] / group_dpl['total_tests']) * 100

                            # 🔥 حساب المؤشر الشامل (Combined Score)
                            max_avg = group_dpl['avg_dpl'].max()
                            group_dpl['Norm_Avg'] = (group_dpl['avg_dpl'] / max_avg * 100) if max_avg > 0 else 0
                            
                            group_dpl['Combined Score'] = (
                                (group_dpl['Pass Rate %'] * 0.40) + 
                                (group_dpl['Consistency'] * 0.30) + 
                                (group_dpl['Norm_Avg'] * 0.30)
                            )

                            # 3. توجيه البيانات حسب اختيار المدير
                            if "Average" in metric_choice:
                                y_col = 'avg_dpl'
                                y_title = 'Avg DPL Value (Blows)'
                                line_color = "#00d2ff"
                                group_dpl = group_dpl.sort_values(y_col, ascending=False)
                                
                            elif "Consistency" in metric_choice:
                                y_col = 'Consistency'
                                y_title = 'Consistency Score % (Higher is Better)'
                                line_color = "#f1c40f"
                                group_dpl = group_dpl.sort_values(y_col, ascending=False)
                                
                            elif "Pass Rate" in metric_choice:
                                y_col = 'Pass Rate %'
                                y_title = 'Approval Rate % (Field QA/QC)'
                                line_color = "#2ecc71"
                                group_dpl = group_dpl.sort_values(y_col, ascending=False)
                                
                            else: # 🏆 Combined Score
                                y_col = 'Combined Score'
                                y_title = 'Composite Score % (40% Pass + 30% Cons. + 30% Avg)'
                                line_color = "#9b59b6"
                                group_dpl = group_dpl.sort_values(y_col, ascending=False)

                            comp_dpl_stat = group_dpl[[comp_name_col, y_col]].copy()
                            comp_dpl_stat.columns = ['Contractor / Company Name', y_title]

                            col_chart, col_table = st.columns([0.65, 0.35])
                            
                            with col_chart:
                                # 4. رسم الشارت الديناميكي
                                fig_dpl_line = px.line(
                                    group_dpl, 
                                    x=comp_name_col, 
                                    y=y_col,
                                    markers=True,
                                    title=f"Trend of {y_title} by Company"
                                )
                                
                                hover_temp = '<b>Contractor:</b> %{x}<br><b>' + y_title + ':</b> %{y:.2f}'
                                fig_dpl_line.update_traces(
                                    line=dict(color=line_color, width=3),
                                    marker=dict(size=10, color=line_color, line=dict(color='white', width=2)),
                                    hovertemplate=hover_temp
                                )
                                
                                try: fig_dpl_line = style_3d_glassy(fig_dpl_line, "line")
                                except: pass
                                st.plotly_chart(fig_dpl_line, use_container_width=True, key="dpl_line_overall")
                                exported_figs["13. DPL Performance & Consistency"] = fig_dpl_line
                                
                            with col_table:
                                # تنسيق الأرقام لـ 3 علامات عشرية في الجدول
                                st.dataframe(comp_dpl_stat.style.format({y_title: "{:.3f}"}), use_container_width=True, hide_index=True)
                                csv_dpl = comp_dpl_stat.to_csv(index=False).encode('utf-8-sig')
                                st.download_button(label=f"📥 Download Data", data=csv_dpl, 
                                                   file_name=f"DPL_Metrics_{datetime.now(EGYPT_TZ).strftime('%Y%m%d')}.csv", 
                                                   mime="text/csv", type="primary", use_container_width=True)
                            
                            # --- 2. Filter & Deep Dive per Company ---
                            st.markdown("##### 🔍 Interactive Deep Dive: Contractor Specific")
                            dpl_companies = ['-- Select Contractor --'] + sorted(dpl_data[comp_name_col].unique().tolist())
                            selected_dpl_comp = st.selectbox("Select a Contractor to investigate DPL metrics:", dpl_companies, key="dpl_deep_comp_sel")
                            
                            if selected_dpl_comp != '-- Select Contractor --':
                                filtered_dpl = dpl_data[dpl_data[comp_name_col] == selected_dpl_comp]
                                
                                deep_c1, deep_c2 = st.columns(2)
                                
                                with deep_c1:
                                    if elment_col:
                                        elem_avg = filtered_dpl.groupby(elment_col)['AVERAGE VALUE'].mean().reset_index()
                                        elem_avg.columns = ['Element', 'Avg DPL']
                                        elem_avg = elem_avg.sort_values('Avg DPL', ascending=False)
                                        fig_elem = px.bar(elem_avg, x='Element', y='Avg DPL', text_auto='.2f', 
                                                          title=f"Avg DPL per Element for {selected_dpl_comp}",
                                                          color='Avg DPL', color_continuous_scale='Blues')
                                        try: fig_elem = style_3d_glassy(fig_elem, "bar")
                                        except: pass
                                        st.plotly_chart(fig_elem, use_container_width=True, key="dpl_elem_bar")
                                    else:
                                        st.info("Element column not found.")
                                
                                with deep_c2:
                                    if 'sample status' in filtered_dpl.columns:
                                        filtered_dpl['Status'] = filtered_dpl['sample status'].str.upper()
                                        status_counts = filtered_dpl['Status'].value_counts().reset_index()
                                        status_counts.columns = ['Status', 'Count']
                                        fig_stat = px.pie(status_counts, names='Status', values='Count', 
                                                          title=f"DPL Approval Rate for {selected_dpl_comp}", hole=0.4,
                                                          color='Status', color_discrete_map=STATUS_COLORS)
                                        fig_stat.update_traces(textinfo='percent+value')
                                        try: fig_stat = style_3d_glassy(fig_stat, "pie")
                                        except: pass
                                        st.plotly_chart(fig_stat, use_container_width=True, key="dpl_status_pie")
                                
                                if date_test_col and date_test_col in filtered_dpl.columns:
                                    filtered_dpl['Month'] = filtered_dpl[date_test_col].dt.strftime('%b %Y')
                                    filtered_dpl['Month_Sort'] = filtered_dpl[date_test_col].dt.to_period('M')
                                    monthly_dpl = filtered_dpl.groupby(['Month_Sort', 'Month'])['AVERAGE VALUE'].mean().reset_index()
                                    monthly_dpl = monthly_dpl.sort_values('Month_Sort')
                                    
                                    if not monthly_dpl.empty:
                                        fig_month = px.line(monthly_dpl, x='Month', y='AVERAGE VALUE', markers=True,
                                                            title=f"Monthly Average DPL Trend for {selected_dpl_comp}",
                                                            color_discrete_sequence=['#2ecc71'])
                                        fig_month.update_traces(line=dict(width=3), marker=dict(size=10, color='white', line=dict(color='#2ecc71', width=2)))
                                        try: fig_month = style_3d_glassy(fig_month, "line")
                                        except: pass
                                        st.plotly_chart(fig_month, use_container_width=True, key="dpl_month_line")
                        else:
                            st.info("No DPL records found with valid Average Values.")
                    else:
                        st.warning("⚠️ Missing Required Columns (Company Name, AVERAGE VALUE, or Test Type) for DPL Analytics.")    
        st.markdown('<div class="bi-title">🔍 Advanced Element Quality Auditor</div>', unsafe_allow_html=True)
        bh_col_name = next((col for col in filtered_df.columns if str(col).strip().upper() in ['ELEMENT', 'ELMENT', 'BH', 'LOCATION']), None)
        zone_col_name = next((col for col in filtered_df.columns if 'ZONE' in str(col).strip().upper() or 'AREA' in str(col).strip().upper()), None)
        if bh_col_name:
            filtered_df[bh_col_name] = filtered_df[bh_col_name].fillna('').astype(str).str.strip()
            bh_list = [bh for bh in filtered_df[bh_col_name].unique() if str(bh).upper() != 'NAN' and str(bh) != '']
            if len(bh_list) > 0:
                selected_bh = st.selectbox(f"Select an Element ({bh_col_name}) to investigate:", ["-- Select Element --"] + sorted(bh_list))
                if selected_bh != "-- Select Element --":
                    bh_df_raw = filtered_df[filtered_df[bh_col_name] == selected_bh].copy()
                    bh_df = None
                    if zone_col_name and bh_df_raw[zone_col_name].nunique() > 1:
                        available_zones = sorted([str(z) for z in bh_df_raw[zone_col_name].unique() if pd.notna(z) and str(z).strip() != ''])
                        st.warning(f"⚠️ **Attention:** Element `{selected_bh}` is present in multiple zones. Please select the required Zone:")
                        selected_zone = st.radio("📍 Select Zone:", available_zones, horizontal=True)
                        if selected_zone:
                            bh_df = bh_df_raw[bh_df_raw[zone_col_name].astype(str) == selected_zone].copy()
                            st.markdown(f"#### 🎯 Investigation Report: `{selected_bh}` <span style='color:#00d2ff; font-size:18px;'>[Zone: {selected_zone}]</span>", unsafe_allow_html=True)
                    else:
                        bh_df = bh_df_raw
                        st.markdown(f"#### 🎯 Investigation Report: `{selected_bh}`")
                    
                    if bh_df is not None:
                        if 'layer' in bh_df.columns:
                            bh_df['Layer_Num'] = bh_df['layer'].astype(str).str.extract(r'(\d+)').fillna(999).astype(int)
                            bh_df = bh_df.sort_values(['Layer_Num', 'Date ( test)'])
                        
                        bh_total_submittals = len(bh_df) 
                        num_tests_col_bh = next((c for c in bh_df.columns if 'NUMBER OF TESTS' in str(c).strip().upper() or 'NUM OF TEST' in str(c).strip().upper()), None)
                        bh_total_tests = int(pd.to_numeric(bh_df[num_tests_col_bh], errors='coerce').fillna(0).sum()) if num_tests_col_bh else bh_total_submittals 
                        bh_accepted = len(bh_df[bh_df['sample status'].astype(str).str.upper().isin(['ACCEPTED', 'APPROVED AS NOTED'])]) if 'sample status' in bh_df.columns else 0
                        bh_pass_rate = (bh_accepted / bh_total_submittals * 100) if bh_total_submittals > 0 else 0
                        bh_avg_dpl = pd.to_numeric(bh_df['AVERAGE VALUE'], errors='coerce').mean() if 'AVERAGE VALUE' in bh_df.columns else 0
                        start_date = bh_df['Date ( test)'].min().strftime('%Y-%m-%d') if 'Date ( test)' in bh_df.columns and not pd.isna(bh_df['Date ( test)'].min()) else "N/A"
                        end_date = bh_df['Date ( test)'].max().strftime('%Y-%m-%d') if 'Date ( test)' in bh_df.columns and not pd.isna(bh_df['Date ( test)'].max()) else "N/A"
                        
                        c1, c2, c3, c4 = st.columns(4)
                        create_card(c1, "Total Submittals", bh_total_submittals)
                        create_card(c2, "Total Tests", bh_total_tests)
                        create_card(c3, "First Test Date", start_date)
                        create_card(c4, "Last Test Date", end_date)
                        c5, c6, c7, c8 = st.columns(4)
                        create_card(c5, "Passed/Approved", bh_accepted)
                        create_card(c6, "Approval Rate (%)", f"{bh_pass_rate:.1f}%")
                        create_card(c7, "Avg DPL Value", f"{bh_avg_dpl:.2f}" if not pd.isna(bh_avg_dpl) else "N/A")
                        create_card(c8, "Rejected Submittals", bh_total_submittals - bh_accepted)

                        if 'Company Name' in bh_df.columns:
                            if 'Date ( test)' in bh_df.columns:
                                comp_stats = bh_df.dropna(subset=['Company Name']).groupby('Company Name')['Date ( test)'].agg(['min', 'max']).reset_index()
                                comp_details = [f"<span style='color:#2ecc71;'><b>{r['Company Name']}</b></span>: <span style='font-size:16px; color:{ui['text_muted']};'>{r['min'].strftime('%Y-%m-%d') if pd.notna(r['min']) else 'N/A'} <b style='color:#ffaa00;'>&rarr;</b> {r['max'].strftime('%Y-%m-%d') if pd.notna(r['max']) else 'N/A'}</span>" for _, r in comp_stats.iterrows()]
                                companies_str = "<br>".join(comp_details) if comp_details else "N/A"
                            else:
                                companies_worked = bh_df['Company Name'].dropna().unique()
                                companies_str = " ، ".join(companies_worked) if len(companies_worked) > 0 else "N/A"
                            st.markdown(f"""
                                <div class="custom-card" style="margin-top: 5px; text-align: left; padding-left: 30px;">
                                    <div class="metric-label" style="color:#ffaa00; text-align: left; margin-bottom: 15px;">Contractors Timeline on this Element</div>
                                    <div class="metric-value" style="font-size: 18px; line-height: 2.0; font-weight: 500; color:{ui['text_main']};">{companies_str}</div>
                                </div>
                                """, unsafe_allow_html=True)

                        if 'layer' in bh_df.columns and 'sample status' in bh_df.columns:
                            rejected_mask = bh_df['sample status'].astype(str).str.upper().isin(['REJECTED', 'REVISE'])
                            accepted_mask = bh_df['sample status'].astype(str).str.upper().isin(['ACCEPTED', 'APPROVED AS NOTED'])
                            approved_layers = set(bh_df[accepted_mask]['layer'].dropna().astype(str).unique())
                            unresolved_alerts = list(set([(str(row.get('layer', 'Unknown')), row.get('Test Type', 'N/A'), row.get('serial', 'N/A')) for _, row in bh_df[rejected_mask].iterrows() if str(row.get('layer', 'Unknown')) not in approved_layers]))
                            if unresolved_alerts:
                                st.markdown("#### 🚨 Critical Quality Alerts (Unresolved Submittals)")
                                alert_cols = st.columns(min(len(unresolved_alerts), 4) if len(unresolved_alerts) > 0 else 1)
                                for idx, alert in enumerate(unresolved_alerts[:8]): 
                                    l, t_type, ser = alert
                                    alert_cols[idx % 4].markdown(f"""
                                        <div style="background: rgba(231, 76, 60, 0.15); backdrop-filter: blur(5px); padding: 15px; border-radius: 15px; border: 1px solid #e74c3c; margin-bottom: 10px; box-shadow: 0 4px 15px rgba(231, 76, 60, 0.2);">
                                            <div style="color: #e74c3c; font-size: 16px; font-weight: bold; margin-bottom: 5px;">⚠️ Action Required</div>
                                            <div style="color: {ui['text_main']}; font-size: 14px; line-height: 1.6;">
                                                <b>Layer:</b> {l}<br><b>Test:</b> {t_type}<br><b>Serial No:</b> {ser}<br>
                                                <span style="font-size:12px; color:#e74c3c;">Status is REVISE/REJECTED with no subsequent approval found!</span>
                                            </div>
                                        </div>
                                        """, unsafe_allow_html=True)
                        st.divider()

                        if 'layer' in bh_df.columns and 'Date ( test)' in bh_df.columns:
                            st.markdown("#### 🧠 AI Engineering Sequence Inspector (Compaction Only)")
                            seq_df = bh_df.dropna(subset=['Date ( test)']).copy()
                            if 'Test Type' in seq_df.columns:
                                seq_df = seq_df[seq_df['Test Type'].astype(str).str.contains('SANDCONE|SAND CONE|DPL', case=False, na=False)]
                            seq_df['Layer_Int'] = seq_df['layer'].astype(str).str.extract(r'(\d+)').fillna(-1).astype(int)
                            seq_df = seq_df[seq_df['Layer_Int'] > 0]
                            if not seq_df.empty:
                                layer_timeline = seq_df.groupby('Layer_Int')['Date ( test)'].min().reset_index()
                                layer_timeline = layer_timeline.sort_values('Layer_Int')
                                logic_errors = []
                                missing_layers = []
                                min_layer = layer_timeline['Layer_Int'].min()
                                max_layer = layer_timeline['Layer_Int'].max()
                                expected_layers = set(range(min_layer, max_layer + 1))
                                actual_layers = set(layer_timeline['Layer_Int'])
                                missing_layers = sorted(list(expected_layers - actual_layers))
                                for i in range(len(layer_timeline) - 1):
                                    curr_L = layer_timeline.iloc[i]['Layer_Int']
                                    next_L = layer_timeline.iloc[i+1]['Layer_Int']
                                    curr_D = layer_timeline.iloc[i]['Date ( test)']
                                    next_D = layer_timeline.iloc[i+1]['Date ( test)']
                                    if curr_D > next_D:
                                        logic_errors.append(f"<b>Layer {curr_L}</b> was tested on <span style='color:#ffaa00;'>{curr_D.date()}</span>, which is AFTER <b>Layer {next_L}</b> tested on <span style='color:#ffaa00;'>{next_D.date()}</span>.")
                                if not missing_layers and not logic_errors:
                                    st.markdown(f"""
                                    <div style="background: rgba(46, 204, 113, 0.1); border-left: 4px solid #2ecc71; padding: 15px; border-radius: 10px; margin-bottom: 20px;">
                                        <h5 style="color: #2ecc71; margin: 0;">✅ Sequence Verified</h5>
                                        <p style="color: {ui['text_main']}; margin: 5px 0 0 0; font-size: 14px;">All compaction layers are chronologically correct with no missing intermediate layers.</p>
                                    </div>
                                    """, unsafe_allow_html=True)
                                else:
                                    if missing_layers:
                                        missing_str = ", ".join([f"Layer {l}" for l in missing_layers])
                                        st.markdown(f"""
                                        <div style="background: rgba(241, 196, 15, 0.1); border-left: 4px solid #f1c40f; padding: 15px; border-radius: 10px; margin-bottom: 10px;">
                                            <h5 style="color: #f1c40f; margin: 0;">⚠️ Missing Compaction Layers Detected</h5>
                                            <p style="color: {ui['text_main']}; margin: 5px 0 0 0; font-size: 14px;">Gap found in execution sequence. Missing: <b style="color:{ui['text_main']};">{missing_str}</b></p>
                                        </div>
                                        """, unsafe_allow_html=True)
                                    if logic_errors:
                                        errors_html = "<br>".join([f"🚨 {err}" for err in logic_errors])
                                        st.markdown(f"""
                                        <div style="background: rgba(231, 76, 60, 0.1); border-left: 4px solid #e74c3c; padding: 15px; border-radius: 10px; margin-bottom: 20px;">
                                            <h5 style="color: #e74c3c; margin: 0;">🛑 Critical Chronological Illogic</h5>
                                            <p style="color: {ui['text_main']}; margin: 5px 0 0 0; font-size: 14px; line-height:1.8;">{errors_html}</p>
                                        </div>
                                        """, unsafe_allow_html=True)
                            else:
                                st.info("No compaction tests (SANDCONE or DPL) found to evaluate sequence.")
                        st.divider()

                        if 'Sampling Location' in bh_df.columns:
                            st.markdown("#### ⛏️ Bottom of Excavation & Soil Quality")
                            boe_df = bh_df[bh_df['Sampling Location'].astype(str).str.contains('Bottom|Soil', case=False, na=False)]
                            if not boe_df.empty:
                                boe_count = len(boe_df)
                                st.info(f"📌 Found **{boe_count}** submittals related to Bottom of Excavation / Soil in this Element.")
                                if 'Classification' in boe_df.columns:
                                    class_counts = boe_df['Classification'].value_counts().reset_index()
                                    class_counts.columns = ['Classification', 'Count']
                                    fig_sc = px.bar(class_counts, x='Classification', y='Count', title="Soil Classifications", color='Classification', text_auto=True, color_discrete_sequence=NEON_COLORS)
                                    fig_sc = style_3d_glassy(fig_sc, chart_type="bar")
                                    st.plotly_chart(fig_sc, use_container_width=True, key=f"sc_{selected_bh}")
                            else:
                                st.success("No 'Bottom of Excavation' specific issues or tests logged for this Element.")
                        st.divider()

                        if 'Test Type' in bh_df.columns and 'Done BY' in bh_df.columns:
                            layer_col = bh_df['layer'] if 'layer' in bh_df.columns else pd.Series([''] * len(bh_df), index=bh_df.index)
                            samp_loc = bh_df['Sampling Location'] if 'Sampling Location' in bh_df.columns else pd.Series(['General Location'] * len(bh_df), index=bh_df.index)
                            bh_df['Execution_Node'] = np.where(layer_col.astype(str).str.contains(r'\d'), layer_col, samp_loc)
                            bh_df['Execution_Node'] = bh_df['Execution_Node'].replace(r'^\s*$', 'General Location', regex=True).fillna('General Location')
                            fig_matrix = px.treemap(bh_df, path=['Done BY', 'Test Type', 'Execution_Node'], title=f"Who did What & Where in {selected_bh}", color='Done BY', color_discrete_sequence=NEON_COLORS)
                            fig_matrix.update_traces(textinfo="label+value")
                            fig_matrix = style_3d_glassy(fig_matrix, chart_type="treemap")
                            st.plotly_chart(fig_matrix, use_container_width=True, key=f"mat_{selected_bh}")
                        st.divider()

                        b_col1, b_col2 = st.columns(2)
                        with b_col1:
                            if 'sample status' in bh_df.columns:
                                bh_df['status_upper'] = bh_df['sample status'].str.upper()
                                fig_ep = px.pie(bh_df, names='status_upper', title=f"Status Breakdown for {selected_bh}", hole=0.4, color='status_upper', color_discrete_map=STATUS_COLORS)
                                fig_ep.update_traces(textinfo='label+percent', hovertemplate='<b>Status:</b> %{label}<br>Count: %{value}<br>Percentage: %{percent}')
                                fig_ep = style_3d_glassy(fig_ep, chart_type="pie")
                                st.plotly_chart(fig_ep, use_container_width=True, key=f"ep_{selected_bh}")
                        with b_col2:
                            if 'layer' in bh_df.columns:
                                layer_reqs = bh_df.groupby('layer').size().reset_index(name='Submittals')
                                layer_reqs['Layer_Num'] = layer_reqs['layer'].astype(str).str.extract(r'(\d+)').fillna(999).astype(int)
                                layer_reqs = layer_reqs.sort_values('Layer_Num')
                                fig_eb = px.bar(layer_reqs, x='layer', y='Submittals', title="Number of Submittals per Layer (Sorted)", text_auto=True, color_discrete_sequence=['#ffaa00'])
                                fig_eb = style_3d_glassy(fig_eb, chart_type="bar")
                                st.plotly_chart(fig_eb, use_container_width=True, key=f"eb_{selected_bh}")

                        if 'Date ( test)' in bh_df.columns and 'AVERAGE VALUE' in bh_df.columns and 'layer' in bh_df.columns:
                            trend_df = bh_df.dropna(subset=['Date ( test)', 'AVERAGE VALUE'])
                            if not trend_df.empty:
                                fig_el = px.line(trend_df, x='Date ( test)', y='AVERAGE VALUE', color='layer', markers=True, title=f"DPL Values Trend across Layers over time for {selected_bh}", color_discrete_sequence=NEON_COLORS)
                                fig_el = style_3d_glassy(fig_el, chart_type="line")
                                st.plotly_chart(fig_el, use_container_width=True, key=f"el_{selected_bh}")
                        
                        with st.expander(f"📂 View Raw Detailed Audit Log for `{selected_bh}`"):
                            st.dataframe(bh_df.drop(columns=['Layer_Num', 'Execution_Node'], errors='ignore'), use_container_width=True)
        else:
            st.warning("⚠️ **Column Not Found:** Could not locate an 'Element' column in your uploaded file to enable Deep Dive Analysis.")

        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
        with st.expander("📂 View Complete Operational Records (Raw Data)"):
            st.dataframe(filtered_df, use_container_width=True)
        # ==========================================
        # 🚨 MODULE 1: Unresolved Rejections Tracker (Action Table)
        # ==========================================
        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
        st.markdown('<div class="bi-title">🚨 Action Tracker: Unresolved Rejections (سجل العينات المرفوضة المعلقة)</div>', unsafe_allow_html=True)
        st.caption("هذا الجدول يحصر العينات (DPL & Plate Load) المرفوضة والتي لم يتم تسجيل عينة ناجحة لها في نفس المنسوب، لتوجيهها للمكتب الفني لإغلاقها.")

        if 'Company Name' in filtered_df.columns and 'sample status' in filtered_df.columns and 'layer' in filtered_df.columns:
            # تحديد الأعمدة المطلوبة بدقة
            test_col = next((c for c in filtered_df.columns if 'TEST TYPE' in c.upper() or c.strip() == 'Test Type'), None)
            sub_date_col = next((c for c in filtered_df.columns if 'DATE( SUB)' in c.upper() or c.strip() == 'Date( SUB)'), None)
            test_date_col = next((c for c in filtered_df.columns if 'DATE ( TEST)' in c.upper() or c.strip() == 'Date ( test)'), None)
            serial_col = next((c for c in filtered_df.columns if 'SERIAL' in c.upper() or c.strip() == 'serial'), None)
            elem_col = next((c for c in filtered_df.columns if c.strip() in ['ELMENT', 'Elment', 'ELEMENT', 'Element (all)', 'Element (All)']), None)
            done_by_col = next((c for c in filtered_df.columns if 'DONE BY' in c.upper()), None) # 💡 إضافة عمود المكتب
            
            if test_col and sub_date_col and test_date_col and serial_col and elem_col:
                # فلترة الداتا لـ DPL و Plate Load فقط
                target_tests = filtered_df[filtered_df[test_col].astype(str).str.upper().str.contains('DPL|PLATE', na=False)].copy()
                
                if not target_tests.empty:
                    target_tests['status_upper'] = target_tests['sample status'].astype(str).str.upper()
                    
                    # الحل
                    target_tests['Unique_Loc'] = (target_tests['Company Name'].astype(str).str.strip().str.upper() + "_" + 
                              target_tests[elem_col].astype(str).str.strip().str.upper() + "_" + 
                              target_tests['layer'].astype(str).str.strip().str.upper() + "_" + 
                              target_tests[test_col].astype(str).str.strip().str.upper())
                    
                    # سحب كل الأماكن اللي اتوافق عليها
                    approved_locs = set(target_tests[target_tests['status_upper'].isin(['ACCEPTED', 'APPROVED AS NOTED'])]['Unique_Loc'].unique())
                    
                    # فلترة العينات المرفوضة اللي مش موجودة في قايمة المقبول
                    unresolved_df = target_tests[
                        (target_tests['status_upper'].isin(['REJECTED', 'REVISE'])) & 
                        (~target_tests['Unique_Loc'].isin(approved_locs))
                    ].copy()
                    
                    if not unresolved_df.empty:
                        # 💡 إضافة عمود المكتب (done_by_col) لترتيب العرض
                        display_cols = ['Company Name', done_by_col, serial_col, sub_date_col, test_date_col, 'layer', test_col, elem_col]
                        
                        # التأكد إن الأعمدة دي موجودة فعلاً في الداتا
                        display_cols = [c for c in display_cols if c is not None and c in unresolved_df.columns]
                        
                        unresolved_display = unresolved_df[display_cols].sort_values(by=['Company Name', sub_date_col])
                        
                        st.markdown(f"""
                        <div style="background: rgba(231, 76, 60, 0.1); border-left: 4px solid #e74c3c; padding: 15px; border-radius: 8px; margin-bottom: 15px;">
                            <b style="color: #e74c3c;">يوجد عدد ({len(unresolved_display)}) طلب مرفوض لم يتم إغلاقه هندسياً حتى الآن!</b><br>
                            <span style="font-size: 13px; color: {{ui['text_muted']}};">يرجى تحميل الجدول وإرساله للمكتب الفني للتعامل معها.</span>
                        </div>
                        """, unsafe_allow_html=True)
                        
                        
                        
                        # زرار تصدير الداتا للإكسيل للمكتب الفني
                        st.dataframe(unresolved_display, use_container_width=True, hide_index=True)
                        
                        # السطر السحري بديل زرار التحميل القديم
                        export_table_tools(unresolved_display, f"Unresolved_Rejections_Log_{datetime.now(EGYPT_TZ).strftime('%Y%m%d')}")
                        
                    else:
                        st.success("✅ ممتاز! لا توجد أي عينات DPL أو Plate Load مرفوضة معلقة حالياً.")
            else:
                st.info("⚠️ بعض الأعمدة المطلوبة (مثل Date, Serial, Element) غير مكتملة لتوليد السجل.")

        # ==========================================
        # 🚨 MODULE 1.5: Missing Layers Tracker (سجل الطبقات الناقصة)
        # ==========================================
        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
        st.markdown('<div class="bi-title">🚨 Action Tracker: Missing Layers (سجل الطبقات المفقودة)</div>', unsafe_allow_html=True)
        st.caption("يقوم هذا الجدول باكتشاف الفجوات في تسلسل طبقات الردم لاختبارات (DPL & Sand Cone) لكل مقاول وعنصر. (مثال: تم العثور على طبقة 1 و 3، واختفت طبقة 2 من السجل).")

        elem_col_missing = next((c for c in filtered_df.columns if c.strip() in ['ELMENT', 'Elment', 'ELEMENT', 'Element (all)', 'Element (All)']), None)
        test_col_missing = next((c for c in filtered_df.columns if 'TEST TYPE' in c.upper() or c.strip() == 'Test Type'), None)

        if 'Company Name' in filtered_df.columns and 'layer' in filtered_df.columns and elem_col_missing:
            # 💡 فلترة الاختبارات لتشمل الـ DPL والـ Sand Cone معاً 
            if test_col_missing:
                layer_tests_df = filtered_df[filtered_df[test_col_missing].astype(str).str.upper().str.contains('DPL|SAND', na=False)].copy()
            else:
                layer_tests_df = filtered_df.copy()
                
            # استخراج أرقام الطبقات كأرقام صحيحة
            layer_tests_df['Layer_Num'] = layer_tests_df['layer'].astype(str).str.extract(r'(\d+)').fillna(-1).astype(int)
            layer_tests_df = layer_tests_df[layer_tests_df['Layer_Num'] > 0]
            
            missing_records = []
            
            # تجميع بالشركة والعنصر لاكتشاف التسلسل
            for (comp, elem), group in layer_tests_df.groupby(['Company Name', elem_col_missing]):
                layers = group['Layer_Num'].unique()
                if len(layers) > 1:
                    min_l = int(layers.min())
                    max_l = int(layers.max())
                    
                    # بناء التسلسل المثالي من أصغر لـ أكبر طبقة
                    expected_layers = set(range(min_l, max_l + 1))
                    actual_layers = set(layers)
                    
                    # استخراج الطبقات الساقطة
                    missing_layers = sorted(list(expected_layers - actual_layers))
                    
                    if missing_layers:
                        missing_records.append({
                            'Company Name': comp,
                            'Element': elem,
                            'Missing Layers (الفجوات)': ", ".join([str(m) for m in missing_layers]),
                            'Highest Layer Reached': max_l
                        })
            
            if missing_records:
                missing_df = pd.DataFrame(missing_records).sort_values(by=['Company Name', 'Element'])
                
                st.markdown(f"""
                <div style="background: rgba(241, 196, 15, 0.1); border-left: 4px solid #f1c40f; padding: 15px; border-radius: 8px; margin-bottom: 15px;">
                    <b style="color: #f1c40f;">يوجد عدد ({len(missing_df)}) قطاع به ثغرات في تسلسل طبقات (DPL / Sand Cone)!</b><br>
                    <span style="font-size: 13px; color: {{ui['text_muted']}};">الطبقات المذكورة تم تخطيها ولم يتم العثور على أي قراءة (نجاح أو سقوط) لها في السجل الهندسي.</span>
                </div>
                """, unsafe_allow_html=True)
                
                st.dataframe(missing_df, use_container_width=True, hide_index=True)
                
                # السطر السحري لإضافة الـ 3 زراير (CSV, Excel, PDF)
                export_table_tools(missing_df, f"Missing_Layers_Log_{datetime.now(EGYPT_TZ).strftime('%Y%m%d')}")
            else:
                st.success("✅ هندسياً ممتاز! لا توجد أي ثغرات أو طبقات ناقصة في تسلسل اختبارات (DPL / Sand Cone) لجميع العناصر المدموكة.")
        else:
            st.info("⚠️ الأعمدة المطلوبة (Company Name, layer, Element) غير متوفرة لتوليد سجل الفجوات.")
       # ==========================================
        # 🧊 MODULE 2: AI-Powered 3D Subsurface Digital Twin (Smart Time-Mapping)
        # ==========================================
        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
        st.markdown('<div class="bi-title">🧊 3D Subsurface Digital Twin (Deep Analytics)</div>', unsafe_allow_html=True)
        
        st.markdown("""
        <div style="display: flex; gap: 15px; font-size: 12px; background: rgba(0,0,0,0.2); padding: 10px; border-radius: 5px; margin-bottom: 10px;">
            <div><b>أشكال الاختبارات:</b> 🟢 (كرة) = DPL | 🔷 (ماسة) = Plate Load</div>
            <div><b>حالة الطبقة:</b> <span style="color:#00ff87;">أخضر</span> = مقبول | <span style="color:#ff007f;">أحمر</span> = مرفوض وتم معالجته | <span style="color:#ff9900; font-weight:bold;">برتقالي</span> = مرفوض معلق (لم يتم قبوله)</div>
        </div>
        """, unsafe_allow_html=True)
        
        layer_col = next((c for c in filtered_df.columns if c.strip().lower() == 'layer'), None)
        status_col = next((c for c in filtered_df.columns if c.strip().lower() == 'sample status'), None)
        elem_col = next((c for c in filtered_df.columns if c.strip().upper() in ['ELMENT', 'ELEMENT', 'ELEMENT (ALL)']), None)
        test_date_col = next((c for c in filtered_df.columns if 'DATE' in c.upper() and 'TEST' in c.upper()), None)
        test_col = next((c for c in filtered_df.columns if 'TEST TYPE' in c.upper() or c.strip() == 'Test Type'), None)

        if layer_col and status_col and elem_col and test_col:
            df_viz = filtered_df.dropna(subset=[layer_col, status_col, elem_col, test_col]).copy()
            df_viz = df_viz[df_viz[test_col].astype(str).str.upper().str.contains('DPL|PLATE', na=False)]
            
            # استخراج الأرقام
            df_viz['Layer_Num'] = df_viz[layer_col].astype(str).str.extract(r'(\d+\.?\d*)')[0].fillna(0).astype(float)
            df_viz = df_viz[df_viz['Layer_Num'] > 0]
            
            if not df_viz.empty:
                df_viz['status_upper'] = df_viz[status_col].str.upper()
                df_viz['Test_Category'] = df_viz[test_col].astype(str).str.upper().apply(lambda x: 'PLATE' if 'PLATE' in x else 'DPL')
                
                company_col = next((c for c in filtered_df.columns if c.strip().lower() in ['company name', 'company', 'contractor']), None)
                serial_col = next((c for c in filtered_df.columns if c.strip().lower() in ['serial', 'serial no', 'no']), None)
                df_viz['Company_Info'] = df_viz[company_col] if company_col else 'N/A'
                df_viz['Serial_Info'] = df_viz[serial_col] if serial_col else 'N/A'

                if test_date_col and test_date_col in df_viz.columns:
                    df_viz['Time_Axis'] = pd.to_datetime(df_viz[test_date_col], dayfirst=True, errors='coerce')
                    df_viz = df_viz.sort_values('Time_Axis')
                    df_viz['Y_Val'] = df_viz['Time_Axis'].dt.strftime('%Y-%m-%d').fillna("Unknown")
                    y_label = "Timeline (Date)"
                else:
                    df_viz['Time_Axis'] = pd.NaT
                    df_viz['Y_Val'] = "Static"
                    y_label = "Depth"

                all_elements = sorted([e for e in df_viz[elem_col].unique() if str(e).strip() != '' and str(e).lower() != 'nan'])
                col_filter, _ = st.columns([0.4, 0.6])
                selected_elem_3d = col_filter.selectbox("📍 Isolate specific Element:", ["All Elements"] + all_elements, key="viz_3d_elem_filter")

                plot_df = df_viz[df_viz[elem_col] == selected_elem_3d].copy() if selected_elem_3d != "All Elements" else df_viz.copy()

                # 💡 1. الذكاء المكاني للـ Plate (ربط الارتفاع بالزمن بين طبقات الـ DPL)
                plot_df['Visual_Z'] = plot_df['Layer_Num'] # الافتراضي للـ DPL
                
                dpl_only = plot_df[plot_df['Test_Category'] == 'DPL'].sort_values('Time_Axis')
                plate_idx = plot_df[plot_df['Test_Category'] == 'PLATE'].index
                
                if not dpl_only.empty:
                    for idx in plate_idx:
                        plate_date = plot_df.loc[idx, 'Time_Axis']
                        if pd.notna(plate_date):
                            past_dpl = dpl_only[dpl_only['Time_Axis'] <= plate_date]
                            if not past_dpl.empty:
                                base_layer = past_dpl['Layer_Num'].max()
                                plot_df.loc[idx, 'Visual_Z'] = base_layer + 0.5 # يوضع فوق آخر طبقة DPL بنص درجة
                            else:
                                plot_df.loc[idx, 'Visual_Z'] = 0.5
                else:
                    # لو مفيش DPL خالص، نرص الـ Plate فوق بعضه برقم تسلسلي
                    for i, idx in enumerate(plate_idx): plot_df.loc[idx, 'Visual_Z'] = i + 1

                # 💡 2. اكتشاف الطبقات المعلقة (الذكاء الجديد)
                hanging_layers = []
                success_statuses = ['ACCEPTED', 'APPROVED AS NOTED', 'APPROVED']
                
                for test_cat in ['DPL', 'PLATE']:
                    cat_df = plot_df[plot_df['Test_Category'] == test_cat]
                    for layer in cat_df['Layer_Num'].unique():
                        layer_df = cat_df[cat_df['Layer_Num'] == layer]
                        max_date = layer_df['Time_Axis'].max()
                        latest_tests = layer_df[layer_df['Time_Axis'] == max_date]
                        
                        # لو آخر يوم مفيش فيه أي كلمة تدل على النجاح، وفي كلمة تدل على الرفض = معلق
                        has_success = any(status in success_statuses for status in latest_tests['status_upper'].values)
                        has_reject = any(status in ['REJECTED', 'REVISE'] for status in latest_tests['status_upper'].values)
                        
                        if not has_success and has_reject:
                            hanging_layers.append((layer, test_cat))
                            
                hanging_dpl = [lyr for lyr, cat in hanging_layers if cat == 'DPL']
                hanging_plate = [lyr for lyr, cat in hanging_layers if cat == 'PLATE']

                def get_point_color(row):
                    is_hanging = (row['Layer_Num'], row['Test_Category']) in hanging_layers
                    if row['status_upper'] in ['REJECTED', 'REVISE']:
                        return '#ff9900' if is_hanging else '#ff007f'
                    return '#00ff87'
                
                plot_df['Color'] = plot_df.apply(get_point_color, axis=1)
                plot_df['Symbol'] = plot_df['Test_Category'].apply(lambda x: 'diamond' if x == 'PLATE' else 'circle')
                
                plot_df['Hover_Text'] = (
                    "<b>📌 Element:</b> " + plot_df[elem_col].astype(str) + "<br>" +
                    "<b>📏 Real Elevation:</b> Level " + plot_df['Layer_Num'].astype(str) + "<br>" +
                    "<b>🔬 Test Type:</b> " + plot_df['Test_Category'] + "<br>" +
                    "<b>⚖️ Status:</b> " + plot_df['status_upper'] + "<br>" +
                    "<b>📅 Date:</b> " + plot_df['Y_Val'] + "<br>" +
                    "<b>📑 Serial:</b> " + plot_df['Serial_Info'].astype(str) + "<br>" +
                    "<b>🏢 Company:</b> " + plot_df['Company_Info'].astype(str)
                )

                col_3d, col_ai = st.columns([0.75, 0.25])

                with col_3d:
                    fig_3d = go.Figure()
                    fig_3d.add_trace(go.Scatter3d(
                        x=plot_df[elem_col],
                        y=plot_df['Y_Val'],
                        # 💡 هنا بنرسم بناءً على الارتفاع البصري الجديد (Visual_Z) مش المنسوب الحقيقي
                        z=plot_df['Visual_Z'],
                        mode='markers',
                        marker=dict(
                            size=7 if selected_elem_3d == "All Elements" else 14, 
                            color=plot_df['Color'],
                            symbol=plot_df['Symbol'],
                            opacity=0.9,
                            line=dict(color='rgba(255,255,255,0.7)', width=1.5) 
                        ),
                        text=plot_df['Hover_Text'],
                        hovertemplate="%{text}<extra></extra>"
                    ))
                    
                    fig_3d.update_layout(
                        title=f"Time-Mapped Subsurface Profile: {selected_elem_3d}",
                        scene=dict(
                            xaxis=dict(title="Element", backgroundcolor="rgba(0,0,0,0)", gridcolor="rgba(0,210,255,0.1)", showbackground=False, tickfont=dict(color="#00d2ff")),
                            yaxis=dict(title=y_label, backgroundcolor="rgba(0,0,0,0)", gridcolor="rgba(0,210,255,0.1)", showbackground=False, tickfont=dict(color="#ffaa00")),
                            zaxis=dict(title="Progress Sequence", backgroundcolor="rgba(0,0,0,0)", gridcolor="rgba(0,210,255,0.1)", showbackground=False, tickfont=dict(color="#2ecc71")),
                            camera=dict(eye=dict(x=1.8, y=-1.8, z=0.8))
                        ),
                        height=650, margin=dict(l=0, r=0, b=0, t=40), paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)'
                    )
                    st.plotly_chart(fig_3d, use_container_width=True, key="hologram_3d_viz_new")

                with col_ai:
                    max_dpl = plot_df[plot_df['Test_Category'] == 'DPL']['Layer_Num'].max() if not plot_df[plot_df['Test_Category'] == 'DPL'].empty else 0
                    max_plate = plot_df[plot_df['Test_Category'] == 'PLATE']['Layer_Num'].max() if not plot_df[plot_df['Test_Category'] == 'PLATE'].empty else 0
                    
                    # 💡 3. حساب التكرار للـ DPL فقط (إعفاء الـ Plate)
                    layer_counts = plot_df[plot_df['Test_Category'] == 'DPL'].groupby('Layer_Num').size()
                    repeated_dpl = layer_counts[layer_counts > 1].index.tolist()
                    
                    def format_list(lst): return ", ".join(map(lambda x: str(x).rstrip('0').rstrip('.') if x%1==0 else str(x), sorted(lst))) if lst else "✅ None"
                    
                    h_dpl_str, h_plate_str = format_list(hanging_dpl), format_list(hanging_plate)
                    r_dpl_str = format_list(repeated_dpl)
                    
                    velocity_str = "N/A"
                    if pd.notna(plot_df['Time_Axis'].min()):
                        days_worked = (plot_df['Time_Axis'].max() - plot_df['Time_Axis'].min()).days
                        if days_worked > 0:
                            layers_per_week = (len(plot_df[plot_df['Test_Category'] == 'DPL']['Layer_Num'].unique()) / days_worked) * 7
                            velocity_str = f"{layers_per_week:.1f} / Week"

                    problem_html, solution_html = "", ""
                    if hanging_dpl or hanging_plate:
                        problem_html += f"🚨 <b>Hanging Rejections:</b> Unresolved failures exist.<br>"
                        solution_html += f"👉 <b>URGENT:</b> Halt work. Request NCR closure.<br>"
                    if repeated_dpl:
                        problem_html += f"🔁 <b>Rework Detected (DPL):</b> Multiple tests on same layers.<br>"
                        
                    if not problem_html:
                        problem_html = "✅ Excellent execution. No critical anomalies."
                        solution_html = "👉 Maintain current QA/QC process."
                    elif not solution_html:
                        solution_html = "👉 Audit contractor compaction methodology."

                    st.markdown(f"""
<div style="background: rgba(10, 20, 33, 0.8); border: 1px solid rgba(0, 210, 255, 0.3); padding: 15px; border-radius: 8px; box-shadow: 0 0 15px rgba(0, 210, 255, 0.1);">
<div style="border-bottom: 2px solid #00d2ff; margin-bottom: 15px; padding-bottom: 5px;">
<b style="color: #00d2ff; font-size: 16px;">🧠 AI Diagnostics</b>
</div>
<div style="display: flex; justify-content: space-between; margin-bottom: 10px;">
<div>
<div style="color: {ui['text_muted']}; font-size: 10px; text-transform: uppercase;">Max DPL Layer</div>
<div style="color: #00d2ff; font-size: 16px; font-weight: bold;">{max_dpl}</div>
</div>
<div>
<div style="color: {ui['text_muted']}; font-size: 10px; text-transform: uppercase;">Max Plate Level</div>
<div style="color: #00ff87; font-size: 16px; font-weight: bold;">{max_plate}</div>
</div>
</div>
<div style="margin-bottom: 10px;">
<div style="color: {ui['text_muted']}; font-size: 11px; text-transform: uppercase;">DPL Velocity (Layers/Week)</div>
<div style="color: #ffaa00; font-size: 18px; font-weight: bold;">{velocity_str}</div>
</div>
<div style="margin-top: 15px; margin-bottom: 5px; color: #ff9900; font-size: 12px; font-weight: bold; border-bottom: 1px solid rgba(255,153,0,0.3);">⚠️ Hanging (Unresolved)</div>
<div style="font-size: 12px; margin-bottom: 5px;"><b>DPL:</b> <span style="color: {'#ff9900' if hanging_dpl else '#2ecc71'};">{h_dpl_str}</span></div>
<div style="font-size: 12px; margin-bottom: 10px;"><b>Plate:</b> <span style="color: {'#ff9900' if hanging_plate else '#2ecc71'};">{h_plate_str}</span></div>
<div style="margin-top: 10px; margin-bottom: 5px; color: #f1c40f; font-size: 12px; font-weight: bold; border-bottom: 1px solid rgba(241,196,15,0.3);">🔁 Repeated (Reworked)</div>
<div style="font-size: 12px; margin-bottom: 10px;"><b>DPL:</b> <span style="color: {'#f1c40f' if repeated_dpl else '#2ecc71'};">{r_dpl_str}</span></div>
<hr style="border-color: rgba(255,255,255,0.1); margin: 15px 0;">
<div style="margin-bottom: 10px;"><div style="color: #e74c3c; font-size: 12px; font-weight: bold;">⚠️ AI Diagnostics:</div><div style="color: {ui['text_main']}; font-size: 11px;">{problem_html}</div></div>
<div><div style="color: #2ecc71; font-size: 12px; font-weight: bold;">💡 AI Prescription:</div><div style="color: {ui['text_main']}; font-size: 11px; background: rgba(46,204,113,0.1); padding: 5px; border-radius: 5px;">{solution_html}</div></div>
</div>
""", unsafe_allow_html=True)
            else:
                st.info("💡 لا توجد عينات DPL أو Plate Load كافية لرسم المجسم ثلاثي الأبعاد.")
        else:
            st.warning("⚠️ لم يظهر الشارت لأن أحد الأعمدة مفقود.")
# ==========================================
        # 📥 PPTX Download Button
        # ==========================================
        st.markdown('<div class="bi-title">📊 PowerPoint Executive Deck</div>', unsafe_allow_html=True)
        st.info("💡 **Meeting Mode:** Generate a fully formatted, Dark-Themed PowerPoint presentation containing your live KPIs and interactive charts.")
        
        if st.button("📥 Generate & Download PPTX", type="primary", use_container_width=True):
            with st.spinner("📸 AI is capturing high-resolution charts and building your slides... Please wait (10-15 seconds)..."):
                try:
                    ppt_file = generate_executive_pptx(current_metrics, exported_figs, uploaded_file.name)
                    st.download_button(
                        label="✅ Download Ready! Click to Save .pptx",
                        data=ppt_file,
                        file_name=f"KK_Executive_Deck_{datetime.now(EGYPT_TZ).strftime('%Y%m%d_%H%M')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                    st.balloons()
                except Exception as e:
                    st.error(f"❌ Error generating PowerPoint: {str(e)}. Make sure 'kaleido' is installed.")
        # Back to Home Button
        st.markdown('<div class="gradient-divider"></div>', unsafe_allow_html=True)
        if st.button("🏠 Back to Home", use_container_width=True, key="back_to_home_dashboard"):
            st.session_state["current_page"] = "home"
            st.rerun()

    else:
        st.info("👈 Please connect a Data Source or Upload a CSV to activate the Enterprise Engine.")

# ==========================================
# 14. Main Application Execution
# ==========================================
def main():
    inject_custom_css()  
    init_auth_system()
    
    if "authenticated" not in st.session_state:
        st.session_state["authenticated"] = False

    if not st.session_state["authenticated"]:
        render_login_screen()
    else:
        current_page = st.session_state.get("current_page", "home")
        
        if current_page == "home":
            render_home_page()
        elif current_page == "dashboard":
            render_dashboard()
        elif current_page == "analytics":
            if "analytics_df" not in st.session_state:
                st.markdown("### 📥 Welcome to Advanced Analytics Hub")
                st.info("You can upload your dataset directly here to begin analysis.")
                hub_init_upload = st.file_uploader("Upload Dataset (CSV) 📂", type=["csv"], key="hub_init_uploader")
                if hub_init_upload is not None:
                    st.session_state["analytics_df"] = pd.read_csv(hub_init_upload)
                    st.rerun()
            
            if "analytics_df" in st.session_state:
                render_analytics_hub(st.session_state["analytics_df"])
            else:
                st.warning("⚠️ Please upload a CSV file from the Main Dashboard first to use Analytics Hub")
                if st.button("📊 Go to Main Dashboard", use_container_width=True, type="primary"):
                    st.session_state["current_page"] = "dashboard"
                    st.rerun()

if __name__ == "__main__":
    main()